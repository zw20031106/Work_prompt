#!/usr/bin/env python3
"""审计参考 PPTX 模板的页面系统、文字内容与字号分布。

定位与作用
----------
这个脚本服务 `ppt-polished-deck-collab` 的模板取证流程。
当用户给一个现成 `pptx` 并要求“照着这个模板做”时，agent 需要先把模板当成
页面系统对象来观察，而不是直接开始重画页面。

大致流程
----------
1. 读取 `pptx`，统计 slide / master / layout 基本信息；
2. 提取 slide、layout、master 三层中的文本对象；
3. 统计显式字号分布，并按字号保留样本文本；
4. 输出结构化 JSON，必要时再写一份 Markdown 摘要；
5. 让后续 `brief.md` / `deck_narrative.md` 能基于事实记录模板取证结果。
"""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter, defaultdict
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


DEFAULT_SAMPLE_LIMIT = 3
DEFAULT_TEXT_PREVIEW_LIMIT = 90


def parse_args() -> argparse.Namespace:
    """解析命令行参数。"""
    parser = argparse.ArgumentParser(description="审计参考 PPTX 模板的 layout、文字与字号分布")
    parser.add_argument("--pptx", required=True, type=Path, help="输入参考 PPTX")
    parser.add_argument("--json-out", type=Path, help="可选：写出完整审计 JSON")
    parser.add_argument("--md-out", type=Path, help="可选：写出 Markdown 摘要")
    parser.add_argument(
        "--sample-limit",
        type=int,
        default=DEFAULT_SAMPLE_LIMIT,
        help="每个字号保留多少条样本文本",
    )
    parser.add_argument(
        "--text-preview-limit",
        type=int,
        default=DEFAULT_TEXT_PREVIEW_LIMIT,
        help="单条样本文本最多保留多少字符",
    )
    return parser.parse_args()


def normalize_text(text: str) -> str:
    """压缩空白并返回更适合统计的单行文本。"""
    collapsed = re.sub(r"\s+", " ", text or "").strip()
    return collapsed


def shorten_text(text: str, limit: int) -> str:
    """截断长文本，便于在 JSON / Markdown 中阅读。"""
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def length_to_pt(length) -> float | None:
    """把 python-pptx 的长度对象转成 pt。"""
    if length is None:
        return None
    return round(float(length.pt), 2)


def enum_label(value) -> str:
    """把枚举或普通对象转成稳定字符串。"""
    name = getattr(value, "name", None)
    if name:
        return str(name)
    return str(value)


def shape_type_label(shape) -> str:
    """返回 shape type 的可读标签。"""
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "UNKNOWN"
    return enum_label(shape_type)


def placeholder_info(shape) -> dict | None:
    """返回 placeholder 的元信息。"""
    if not getattr(shape, "is_placeholder", False):
        return None
    info = {
        "idx": shape.placeholder_format.idx,
        "type": enum_label(shape.placeholder_format.type),
    }
    return info


def collect_text_items_from_text_frame(
    *,
    text_frame,
    layer: str,
    container_name: str,
    shape_name: str,
    shape_type: str,
    placeholder: dict | None,
    sample_limit: int,
    preview_limit: int,
) -> list[dict]:
    """从 text frame 提取逐段文本与显式字号。"""
    paragraphs: list[dict] = []
    explicit_sizes: Counter[float] = Counter()

    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        raw_text = paragraph.text or ""
        normalized = normalize_text(raw_text)
        if not normalized:
            continue

        paragraph_sizes_set: set[float] = set()
        paragraph_level_size = length_to_pt(paragraph.font.size)
        if paragraph_level_size is not None:
            paragraph_sizes_set.add(paragraph_level_size)

        for run in paragraph.runs:
            run_size = length_to_pt(run.font.size)
            if run_size is not None:
                paragraph_sizes_set.add(run_size)

        paragraph_sizes = sorted(paragraph_sizes_set)
        explicit_sizes.update(paragraph_sizes)
        paragraphs.append(
            {
                "paragraph_index": paragraph_index,
                "text": shorten_text(normalized, preview_limit),
                "font_sizes_pt": paragraph_sizes,
            }
        )

    if not paragraphs:
        return []

    combined_text = " | ".join(item["text"] for item in paragraphs[:sample_limit])
    return [
        {
            "layer": layer,
            "container_name": container_name,
            "shape_name": shape_name,
            "shape_type": shape_type,
            "placeholder": placeholder,
            "text": shorten_text(combined_text, preview_limit),
            "paragraphs": paragraphs,
            "font_sizes_pt": sorted(explicit_sizes),
            "explicit_font_size_count": sum(explicit_sizes.values()),
            "inherits_font_size": not explicit_sizes,
        }
    ]


def collect_text_items_from_table(
    *,
    table,
    layer: str,
    container_name: str,
    shape_name: str,
    shape_type: str,
    placeholder: dict | None,
    preview_limit: int,
) -> list[dict]:
    """从表格中提取 cell 文本与显式字号。"""
    items: list[dict] = []
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            text = normalize_text(cell.text)
            if not text:
                continue

            paragraph_sizes: list[float] = []
            for paragraph in cell.text_frame.paragraphs:
                paragraph_level_size = length_to_pt(paragraph.font.size)
                if paragraph_level_size is not None:
                    paragraph_sizes.append(paragraph_level_size)
                for run in paragraph.runs:
                    size = length_to_pt(run.font.size)
                    if size is not None:
                        paragraph_sizes.append(size)

            item = {
                "layer": layer,
                "container_name": container_name,
                "shape_name": shape_name,
                "shape_type": shape_type,
                "placeholder": placeholder,
                "cell": {"row": row_index, "col": col_index},
                "text": shorten_text(text, preview_limit),
                "font_sizes_pt": sorted(set(paragraph_sizes)),
                "explicit_font_size_count": len(paragraph_sizes),
                "inherits_font_size": len(paragraph_sizes) == 0,
            }
            items.append(item)
    return items


def collect_shape_text_items(
    *,
    shape,
    layer: str,
    container_name: str,
    sample_limit: int,
    preview_limit: int,
) -> list[dict]:
    """递归提取单个 shape 中的文本对象。"""
    shape_name = getattr(shape, "name", "")
    shape_type = shape_type_label(shape)
    placeholder = placeholder_info(shape)

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        items: list[dict] = []
        for child in shape.shapes:
            items.extend(
                collect_shape_text_items(
                    shape=child,
                    layer=layer,
                    container_name=container_name,
                    sample_limit=sample_limit,
                    preview_limit=preview_limit,
                )
            )
        return items

    items: list[dict] = []
    if getattr(shape, "has_text_frame", False):
        items.extend(
            collect_text_items_from_text_frame(
                text_frame=shape.text_frame,
                layer=layer,
                container_name=container_name,
                shape_name=shape_name,
                shape_type=shape_type,
                placeholder=placeholder,
                sample_limit=sample_limit,
                preview_limit=preview_limit,
            )
        )

    if getattr(shape, "has_table", False):
        items.extend(
            collect_text_items_from_table(
                table=shape.table,
                layer=layer,
                container_name=container_name,
                shape_name=shape_name,
                shape_type=shape_type,
                placeholder=placeholder,
                preview_limit=preview_limit,
            )
        )
    return items


def summarize_text_items(
    *,
    text_items: Iterable[dict],
    sample_limit: int,
) -> dict:
    """统计文本对象的字号分布与样本文本。"""
    font_counter: Counter[float] = Counter()
    font_samples: dict[str, list[str]] = defaultdict(list)
    inherited_count = 0

    for item in text_items:
        if item["inherits_font_size"]:
            inherited_count += 1
        for size in item["font_sizes_pt"]:
            font_counter[size] += 1
            sample_bucket = font_samples[str(size)]
            if len(sample_bucket) < sample_limit:
                sample_bucket.append(item["text"])

    return {
        "text_item_count": len(list(text_items)) if not isinstance(text_items, list) else len(text_items),
        "font_size_distribution_pt": {str(size): count for size, count in sorted(font_counter.items())},
        "font_size_samples": dict(font_samples),
        "inherits_font_size_count": inherited_count,
    }


def summarize_layout_usage(prs: Presentation) -> dict[str, int]:
    """统计每个 layout 被多少张 slide 使用。"""
    usage: Counter[str] = Counter()
    for slide in prs.slides:
        usage[slide.slide_layout.name] += 1
    return {name: count for name, count in sorted(usage.items())}


def analyze_presentation(
    prs: Presentation,
    *,
    sample_limit: int,
    preview_limit: int,
) -> dict:
    """对整个 Presentation 做模板审计。"""
    layout_usage = summarize_layout_usage(prs)

    masters_payload: list[dict] = []
    layouts_payload: list[dict] = []
    slides_payload: list[dict] = []
    default_slide_layouts = [
        {
            "layout_index": index,
            "name": layout.name,
            "partname": str(layout.part.partname),
        }
        for index, layout in enumerate(prs.slide_layouts, start=1)
    ]

    all_layer_items: dict[str, list[dict]] = {"master": [], "layout": [], "slide": []}

    for master_index, master in enumerate(prs.slide_masters, start=1):
        master_name = getattr(master, "name", None) or f"Master {master_index}"
        master_container_name = f"master:{master_name}"
        master_items: list[dict] = []
        for shape in master.shapes:
            master_items.extend(
                collect_shape_text_items(
                    shape=shape,
                    layer="master",
                    container_name=master_container_name,
                    sample_limit=sample_limit,
                    preview_limit=preview_limit,
                )
            )
        all_layer_items["master"].extend(master_items)
        master_summary = summarize_text_items(text_items=master_items, sample_limit=sample_limit)
        masters_payload.append(
            {
                "master_index": master_index,
                "name": master_name,
                "layout_count": len(master.slide_layouts),
                **master_summary,
            }
        )

        for layout_index, layout in enumerate(master.slide_layouts, start=1):
            layout_name = layout.name
            layout_container_name = f"layout:{layout_name}"
            layout_items: list[dict] = []
            for shape in layout.shapes:
                layout_items.extend(
                    collect_shape_text_items(
                        shape=shape,
                        layer="layout",
                        container_name=layout_container_name,
                        sample_limit=sample_limit,
                        preview_limit=preview_limit,
                    )
                )
            all_layer_items["layout"].extend(layout_items)
            layout_summary = summarize_text_items(text_items=layout_items, sample_limit=sample_limit)
            layouts_payload.append(
                {
                    "master_index": master_index,
                    "layout_index": layout_index,
                    "name": layout_name,
                    "used_by_slide_count": layout_usage.get(layout_name, 0),
                    **layout_summary,
                }
            )

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_items: list[dict] = []
        for shape in slide.shapes:
            slide_items.extend(
                collect_shape_text_items(
                    shape=shape,
                    layer="slide",
                    container_name=f"slide:{slide_number}",
                    sample_limit=sample_limit,
                    preview_limit=preview_limit,
                )
            )
        all_layer_items["slide"].extend(slide_items)
        slide_summary = summarize_text_items(text_items=slide_items, sample_limit=sample_limit)
        slides_payload.append(
            {
                "slide_number": slide_number,
                "layout_name": slide.slide_layout.name,
                **slide_summary,
                "text_items": slide_items,
            }
        )

    layer_summaries = {
        layer: summarize_text_items(text_items=items, sample_limit=sample_limit)
        for layer, items in all_layer_items.items()
    }

    return {
        "slide_size_in": {
            "width": round(prs.slide_width / 914400, 3),
            "height": round(prs.slide_height / 914400, 3),
        },
        "summary": {
            "slide_count": len(prs.slides),
            "master_count": len(prs.slide_masters),
            "default_slide_layout_count": len(prs.slide_layouts),
            "all_master_layout_count": sum(len(master.slide_layouts) for master in prs.slide_masters),
            "layout_usage": layout_usage,
            "layer_summaries": layer_summaries,
        },
        "default_slide_layouts": default_slide_layouts,
        "masters": masters_payload,
        "layouts": layouts_payload,
        "slides": slides_payload,
    }


def format_markdown_summary(result: dict, pptx_path: Path) -> str:
    """把审计结果转成精简 Markdown 摘要。"""
    summary = result["summary"]
    slide_layer = summary["layer_summaries"]["slide"]
    layout_layer = summary["layer_summaries"]["layout"]
    master_layer = summary["layer_summaries"]["master"]

    lines = [
        f"# PPTX 模板审计摘要",
        "",
        f"- 文件：`{pptx_path}`",
        f"- 页面尺寸：`{result['slide_size_in']['width']}` x `{result['slide_size_in']['height']}` in",
        f"- slide 数：`{summary['slide_count']}`",
        f"- master 数：`{summary['master_count']}`",
        f"- 默认可用 layout 数：`{summary['default_slide_layout_count']}`",
        f"- 全部 master-layout part 数：`{summary['all_master_layout_count']}`",
        "",
        "## Default Slide Layouts",
        "",
        "| Index | Layout |",
        "| --- | --- |",
    ]
    for layout in result["default_slide_layouts"]:
        lines.append(f"| {layout['layout_index']} | `{layout['name']}` |")

    lines.extend(
        [
            "",
        "## Layout Usage",
        "",
        "| Layout | Used By Slides |",
        "| --- | ---: |",
        ]
    )
    for layout_name, count in summary["layout_usage"].items():
        lines.append(f"| `{layout_name}` | {count} |")

    lines.extend(
        [
            "",
            "## Font Size Distribution",
            "",
            "| Layer | Distribution | Inherited Text Items |",
            "| --- | --- | ---: |",
            f"| `slide` | `{slide_layer['font_size_distribution_pt']}` | {slide_layer['inherits_font_size_count']} |",
            f"| `layout` | `{layout_layer['font_size_distribution_pt']}` | {layout_layer['inherits_font_size_count']} |",
            f"| `master` | `{master_layer['font_size_distribution_pt']}` | {master_layer['inherits_font_size_count']} |",
            "",
            "## Slide Font Samples",
            "",
        ]
    )

    for size, samples in slide_layer["font_size_samples"].items():
        lines.append(f"### `{size}pt`")
        lines.append("")
        for sample in samples:
            lines.append(f"- {sample}")
        lines.append("")

    return "\n".join(lines).rstrip() + "\n"


def main() -> int:
    """执行模板审计。"""
    args = parse_args()
    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        raise SystemExit(f"未找到 PPTX: {pptx_path}")

    prs = Presentation(pptx_path)
    result = {
        "pptx": str(pptx_path),
        **analyze_presentation(
            prs,
            sample_limit=args.sample_limit,
            preview_limit=args.text_preview_limit,
        ),
    }

    summary = result["summary"]
    print(f"[INFO] pptx={pptx_path}")
    print(
        "[INFO] slides={slides} masters={masters} default_layouts={default_layouts} all_master_layouts={all_master_layouts}".format(
            slides=summary["slide_count"],
            masters=summary["master_count"],
            default_layouts=summary["default_slide_layout_count"],
            all_master_layouts=summary["all_master_layout_count"],
        )
    )
    print("[INFO] layout_usage=" + ", ".join(f"{name}:{count}" for name, count in summary["layout_usage"].items()))
    for layer_name, layer_summary in summary["layer_summaries"].items():
        print(
            f"[INFO] {layer_name}_font_sizes={layer_summary['font_size_distribution_pt']} "
            f"inherited={layer_summary['inherits_font_size_count']}"
        )

    if args.json_out:
        args.json_out.parent.mkdir(parents=True, exist_ok=True)
        args.json_out.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[INFO] 写入 JSON: {args.json_out}")

    if args.md_out:
        args.md_out.parent.mkdir(parents=True, exist_ok=True)
        args.md_out.write_text(format_markdown_summary(result, pptx_path), encoding="utf-8")
        print(f"[INFO] 写入 Markdown: {args.md_out}")

    print("[OK] 模板审计完成")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
