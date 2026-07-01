#!/usr/bin/env python3
"""`ppt-polished-deck-collab` 的通用 PPT 资产 helper。

定位与作用
----------
本文件为高质量 deck 的通用构图场景提供最小但稳定的 Python helper，
覆盖标题区、panel、文本、图片卡片、原生 Office chart 和真绑定 connector。
它的目标不是替代业务脚本，而是避免每个 deck 反复复制同一批低层实现。

大致流程
----------
1. 创建统一宽屏 Presentation；
2. 写入统一标题头、字幕和页脚说明；
3. 用 panel、文本块、图片卡片、原生 chart 和节点 shape 组织页面；
4. 在需要时用 glued connector 把线真正粘到业务节点；
5. 把这些 helper 组合进具体 deck 的 build 脚本。
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SITE_INDEX = {
    "top": 0,
    "left": 1,
    "bottom": 2,
    "right": 3,
}

PANEL_LABEL_HEIGHT = 0.40
PANEL_LABEL_GAP = 0.04
PANEL_CONTENT_INSET_X = 0.12
PANEL_CONTENT_INSET_Y = 0.12
DEFAULT_LATIN_FONT_NAME = "Arial"
DEFAULT_EAST_ASIA_FONT_NAME = "黑体"
DEFAULT_SERIF_LATIN_FONT_NAME = "Times New Roman"
DEFAULT_SERIF_EAST_ASIA_FONT_NAME = "宋体"
DEFAULT_FONT_NAME = DEFAULT_LATIN_FONT_NAME
DEFAULT_BODY_LINE_SPACING_MULTIPLE = 1.5
DEFAULT_TITLE_LINE_SPACING_MULTIPLE = 1.0
DEFAULT_TITLE_PARAGRAPH_SPACE_LINES = 0.5
DEFAULT_LINE_SPACING_MULTIPLE = DEFAULT_BODY_LINE_SPACING_MULTIPLE
DEFAULT_TYPOGRAPHY_TOKENS = {
    "hero_title_font_pt": 24.0,
    "section_title_font_pt": 20.0,
    "page_title_font_pt": 24.0,
    "subtitle_font_pt": 18.0,
    "minor_title_font_pt": 16.0,
    "body_font_pt": 14.0,
    "label_font_pt": 12.0,
    "caption_font_pt": 12.0,
    "title_line_spacing_multiple": DEFAULT_TITLE_LINE_SPACING_MULTIPLE,
    "body_line_spacing_multiple": DEFAULT_BODY_LINE_SPACING_MULTIPLE,
    "title_paragraph_space_lines": DEFAULT_TITLE_PARAGRAPH_SPACE_LINES,
}


@dataclass(frozen=True)
class NodeStyle:
    """定义原生 shape 节点的视觉风格。"""

    fill_rgb: tuple[int, int, int]
    line_rgb: tuple[int, int, int]
    text_rgb: tuple[int, int, int] | None = None
    font_size: float = DEFAULT_TYPOGRAPHY_TOKENS["label_font_pt"]
    bold: bool = True
    shape_type: MSO_SHAPE = MSO_SHAPE.RECTANGLE


def default_palette() -> dict[str, tuple[int, int, int]]:
    """返回 deck 默认调色盘。"""
    return {
        "bg": (248, 250, 252),
        "title": (28, 45, 74),
        "subtitle": (93, 109, 136),
        "line": (71, 85, 105),
        "muted": (148, 163, 184),
        "blue": (37, 99, 235),
        "emerald": (16, 185, 129),
        "amber": (245, 158, 11),
        "violet": (124, 58, 237),
        "rose": (225, 29, 72),
        "teal": (13, 148, 136),
        "slate": (100, 116, 139),
    }


def default_typography_tokens() -> dict[str, float]:
    """返回 deck 默认字号 token。"""
    return dict(DEFAULT_TYPOGRAPHY_TOKENS)


def tint(rgb: tuple[int, int, int], alpha: float) -> tuple[int, int, int]:
    """把给定颜色按 alpha 与白色混合。"""
    red, green, blue = rgb
    return (
        int(round(alpha * red + (1 - alpha) * 255)),
        int(round(alpha * green + (1 - alpha) * 255)),
        int(round(alpha * blue + (1 - alpha) * 255)),
    )


def pick_contrast_text_rgb(fill_rgb: tuple[int, int, int]) -> tuple[int, int, int]:
    """在深色与白色候选中选择对比度更高的文字色。"""
    dark = (25, 38, 56)
    light = (255, 255, 255)

    def _srgb_to_linear(channel_8bit: int) -> float:
        value = channel_8bit / 255.0
        if value <= 0.04045:
            return value / 12.92
        return ((value + 0.055) / 1.055) ** 2.4

    def _relative_luminance(rgb: tuple[int, int, int]) -> float:
        red, green, blue = rgb
        return (
            0.2126 * _srgb_to_linear(red)
            + 0.7152 * _srgb_to_linear(green)
            + 0.0722 * _srgb_to_linear(blue)
        )

    def _contrast_ratio(rgb1: tuple[int, int, int], rgb2: tuple[int, int, int]) -> float:
        lum1 = _relative_luminance(rgb1)
        lum2 = _relative_luminance(rgb2)
        lighter, darker = (lum1, lum2) if lum1 >= lum2 else (lum2, lum1)
        return (lighter + 0.05) / (darker + 0.05)

    return dark if _contrast_ratio(fill_rgb, dark) >= _contrast_ratio(fill_rgb, light) else light


def panel_body_top_offset() -> float:
    """返回 panel 从顶部到正文内容区的偏移。"""
    return PANEL_LABEL_HEIGHT + PANEL_LABEL_GAP


def panel_label_width(title: str, panel_width: float) -> float:
    """按标题长度估算 panel label 宽度。"""
    estimate = 0.092 * len(title) + 0.64
    return min(panel_width, max(1.2, estimate))


def panel_content_box(left: float, top: float, width: float, height: float) -> tuple[float, float, float, float]:
    """返回 panel 内适合放正文和图表的内容框。"""
    content_left = left + PANEL_CONTENT_INSET_X
    content_top = top + panel_body_top_offset() + PANEL_CONTENT_INSET_Y
    content_width = max(0.2, width - PANEL_CONTENT_INSET_X * 2)
    content_height = max(0.2, height - panel_body_top_offset() - PANEL_CONTENT_INSET_Y * 2)
    return content_left, content_top, content_width, content_height


def new_presentation() -> Presentation:
    """创建统一 16:9 宽屏 Presentation。"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    return prs


def add_slide_header(slide, figure_tag: str, title: str, subtitle: str) -> None:
    """为页面添加统一标题头。"""
    palette = default_palette()
    tokens = default_typography_tokens()
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*palette["bg"])

    title_box = slide.shapes.add_textbox(Inches(0.72), Inches(0.16), Inches(14.3), Inches(0.72))
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title
    title_para.font.bold = True
    title_para.font.size = Pt(tokens["page_title_font_pt"])
    title_para.font.name = DEFAULT_FONT_NAME
    title_para.line_spacing = tokens["title_line_spacing_multiple"]
    title_para.space_before = Pt(tokens["page_title_font_pt"] * tokens["title_paragraph_space_lines"])
    title_para.space_after = Pt(tokens["page_title_font_pt"] * tokens["title_paragraph_space_lines"])
    title_para.font.color.rgb = RGBColor(*palette["title"])

    sub_box = slide.shapes.add_textbox(Inches(0.74), Inches(0.82), Inches(14.1), Inches(0.40))
    sub_para = sub_box.text_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(tokens["subtitle_font_pt"])
    sub_para.font.name = DEFAULT_FONT_NAME
    sub_para.line_spacing = tokens["title_line_spacing_multiple"]
    sub_para.space_before = Pt(tokens["subtitle_font_pt"] * tokens["title_paragraph_space_lines"])
    sub_para.space_after = Pt(tokens["subtitle_font_pt"] * tokens["title_paragraph_space_lines"])
    sub_para.font.color.rgb = RGBColor(*palette["subtitle"])

    tag_box = slide.shapes.add_textbox(Inches(14.85), Inches(8.43), Inches(0.46), Inches(0.18))
    tag_para = tag_box.text_frame.paragraphs[0]
    tag_para.text = figure_tag
    tag_para.alignment = PP_ALIGN.RIGHT
    tag_para.font.bold = True
    tag_para.font.size = Pt(tokens["caption_font_pt"])
    tag_para.font.name = DEFAULT_FONT_NAME
    tag_para.line_spacing = DEFAULT_LINE_SPACING_MULTIPLE
    tag_para.font.color.rgb = RGBColor(*palette["muted"])


def add_caption(slide, text: str) -> None:
    """在页底添加简短说明。"""
    palette = default_palette()
    tokens = default_typography_tokens()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(8.45), Inches(15.1), Inches(0.22))
    para = box.text_frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(tokens["caption_font_pt"])
    para.font.name = DEFAULT_FONT_NAME
    para.line_spacing = DEFAULT_LINE_SPACING_MULTIPLE
    para.font.color.rgb = RGBColor(*palette["subtitle"])


def add_text_block(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float | None = None,
    bold: bool = False,
    color_rgb: tuple[int, int, int] | None = None,
) -> None:
    """添加纯文本块。"""
    palette = default_palette()
    tokens = default_typography_tokens()
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    para = box.text_frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(font_size if font_size is not None else tokens["body_font_pt"])
    para.font.bold = bold
    para.font.name = DEFAULT_FONT_NAME
    para.line_spacing = DEFAULT_LINE_SPACING_MULTIPLE
    para.font.color.rgb = RGBColor(*(color_rgb or palette["subtitle"]))


def add_panel(
    slide,
    title: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
    body_fill_rgb: tuple[int, int, int] | None = None,
) -> None:
    """添加亮色 label 与浅色正文底块组合的 panel。"""
    tokens = default_typography_tokens()
    body_top = top + panel_body_top_offset()
    body_height = max(0.24, height - panel_body_top_offset())

    body = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(body_top),
        Inches(width),
        Inches(body_height),
    )
    body.text = ""
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(*(body_fill_rgb or tint(accent_rgb, 0.10)))
    body.line.color.rgb = RGBColor(*tint(accent_rgb, 0.60))
    body.line.width = Pt(1.0)

    label_width = panel_label_width(title, width)
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(label_width),
        Inches(PANEL_LABEL_HEIGHT),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*accent_rgb)
    header.line.color.rgb = RGBColor(*accent_rgb)
    header.text_frame.clear()
    header.text_frame.word_wrap = False
    header.text_frame.margin_left = Inches(0.08)
    header.text_frame.margin_right = Inches(0.08)
    header.text_frame.margin_top = Inches(0.00)
    header.text_frame.margin_bottom = Inches(0.00)
    header.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    para = header.text_frame.paragraphs[0]
    para.text = title
    para.alignment = PP_ALIGN.LEFT
    para.font.bold = True
    para.font.size = Pt(tokens["minor_title_font_pt"])
    para.font.name = DEFAULT_FONT_NAME
    para.line_spacing = 1.0
    para.font.color.rgb = RGBColor(*pick_contrast_text_rgb(accent_rgb))


def add_picture_card(
    slide,
    title: str,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
    caption: str | None = None,
) -> None:
    """添加承载高 DPI 图片的 panel 卡片。"""
    add_panel(slide, title, left, top, width, height, accent_rgb, body_fill_rgb=(255, 255, 255))
    content_left, content_top, content_width, content_height = panel_content_box(left, top, width, height)
    image_height = content_height - (0.28 if caption else 0.0)
    max_width = content_width

    with Image.open(image_path) as image:
        src_width, src_height = image.size

    src_ratio = src_width / src_height
    target_ratio = max_width / image_height
    if src_ratio >= target_ratio:
        fit_width = max_width
        fit_height = max_width / src_ratio
    else:
        fit_height = image_height
        fit_width = image_height * src_ratio

    offset_left = content_left + (max_width - fit_width) / 2
    offset_top = content_top + (image_height - fit_height) / 2
    slide.shapes.add_picture(
        str(image_path),
        Inches(offset_left),
        Inches(offset_top),
        width=Inches(fit_width),
        height=Inches(fit_height),
    )

    if caption:
        add_text_block(
            slide,
            caption,
            content_left,
            content_top + image_height + 0.08,
            content_width,
            0.18,
            font_size=DEFAULT_TYPOGRAPHY_TOKENS["caption_font_pt"],
            color_rgb=(99, 115, 141),
        )


def add_native_chart_card(
    slide,
    title: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
    categories: list[str],
    series_list: list[tuple[str, list[float] | tuple[float, ...]]],
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.BAR_CLUSTERED,
    number_format: str = "0",
    show_legend: bool = False,
    legend_position: XL_LEGEND_POSITION = XL_LEGEND_POSITION.BOTTOM,
    series_colors: list[tuple[int, int, int]] | None = None,
):
    """添加 PowerPoint 原生 chart 卡片，并返回 chart 对象。"""
    add_panel(slide, title, left, top, width, height, accent_rgb, body_fill_rgb=(255, 255, 255))
    content_left, content_top, content_width, content_height = panel_content_box(left, top, width, height)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for series_name, values in series_list:
        chart_data.add_series(series_name, values)

    graphic_frame = slide.shapes.add_chart(
        chart_type,
        Inches(content_left),
        Inches(content_top),
        Inches(content_width),
        Inches(content_height),
        chart_data,
    )
    chart = graphic_frame.chart
    chart.chart_style = 10
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = legend_position
        chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.vary_by_categories = False
    plot.has_data_labels = True
    plot.data_labels.number_format = number_format
    plot.data_labels.font.size = Pt(DEFAULT_TYPOGRAPHY_TOKENS["label_font_pt"])
    plot.data_labels.font.name = DEFAULT_FONT_NAME
    plot.data_labels.font.color.rgb = RGBColor(57, 70, 96)

    if hasattr(chart, "category_axis"):
        chart.category_axis.tick_labels.font.size = Pt(DEFAULT_TYPOGRAPHY_TOKENS["label_font_pt"])
        chart.category_axis.tick_labels.font.name = DEFAULT_FONT_NAME
        chart.category_axis.tick_labels.font.color.rgb = RGBColor(74, 85, 104)
        chart.category_axis.has_major_gridlines = False
    if hasattr(chart, "value_axis"):
        chart.value_axis.tick_labels.font.size = Pt(DEFAULT_TYPOGRAPHY_TOKENS["label_font_pt"])
        chart.value_axis.tick_labels.font.name = DEFAULT_FONT_NAME
        chart.value_axis.tick_labels.font.color.rgb = RGBColor(74, 85, 104)
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(226, 232, 240)
    if show_legend:
        chart.legend.font.size = Pt(DEFAULT_TYPOGRAPHY_TOKENS["label_font_pt"])
        chart.legend.font.name = DEFAULT_FONT_NAME
        chart.legend.font.color.rgb = RGBColor(57, 70, 96)

    if series_colors:
        for series, rgb in zip(chart.series, series_colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*rgb)
            series.format.line.color.rgb = RGBColor(*rgb)

    return chart


def add_node(
    slide,
    key: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    style: NodeStyle,
):
    """添加一个可被 connector 绑定的业务节点。"""
    shape = slide.shapes.add_shape(
        style.shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.name = key
    shape.text = text
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*style.fill_rgb)
    shape.line.color.rgb = RGBColor(*style.line_rgb)
    shape.line.width = Pt(1.25)
    text_rgb = style.text_rgb or pick_contrast_text_rgb(style.fill_rgb)
    for para in shape.text_frame.paragraphs:
        para.font.size = Pt(style.font_size)
        para.font.bold = style.bold
        para.line_spacing = DEFAULT_LINE_SPACING_MULTIPLE
        para.font.color.rgb = RGBColor(*text_rgb)
    return shape


def add_glued_connector(
    slide,
    from_shape,
    to_shape,
    from_site: str,
    to_site: str,
    line_rgb: tuple[int, int, int],
    connector_type: MSO_CONNECTOR = MSO_CONNECTOR.ELBOW,
    line_width: float = 1.35,
):
    """创建真正粘连到两个 shape 的 connector。"""
    if from_site not in SITE_INDEX or to_site not in SITE_INDEX:
        raise ValueError(f"未知连接点: from={from_site}, to={to_site}")

    connector = slide.shapes.add_connector(
        connector_type,
        from_shape.left,
        from_shape.top,
        to_shape.left,
        to_shape.top,
    )
    connector.begin_connect(from_shape, SITE_INDEX[from_site])
    connector.end_connect(to_shape, SITE_INDEX[to_site])
    connector.line.color.rgb = RGBColor(*line_rgb)
    connector.line.width = Pt(line_width)
    connector.shadow.inherit = False
    return connector


def save_presentation(prs: Presentation, output_path: Path) -> Path:
    """保存 PPT 并确保输出目录存在。"""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
