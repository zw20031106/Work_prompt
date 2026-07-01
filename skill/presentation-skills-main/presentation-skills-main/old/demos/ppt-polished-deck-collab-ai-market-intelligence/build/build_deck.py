#!/usr/bin/env python3
"""生成 AI 行业市场分析 polished deck demo。

定位与作用
----------
本脚本是 `old/demos/ppt-polished-deck-collab-ai-market-intelligence/` 的主构建入口，
用于把 deck 规划文档、示例数据、icon 资产和图表资产统一编译成一份可编辑 PPT。
它既服务业务表达，也服务 skill 验证，因此会同时覆盖 image hero、native chart、
Python figure、connector diagram、diagram visual、table-like matrix 和 icon accent。

大致流程
----------
1. 读取 `brief.md` 与 `deck_narrative.md`，派生 `slide_specs.yaml`；
2. 生成主题化 icon、Python figures 和 hero image；
3. 按页构建 6 张 slide；
4. 写出 `pptx` 与 `build_manifest.json`，供后续 connector 校验和 preview 导出复用。
"""

from __future__ import annotations

import json
import math
import subprocess
import sys
from pathlib import Path

import pandas as pd
import yaml
from PIL import Image, ImageDraw, ImageFilter
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

WORKSPACE_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = Path(__file__).resolve().parents[3]
SKILL_ROOT = REPO_ROOT / "ppt-polished-deck-collab"

sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from ppt_asset_helpers import (  # noqa: E402
    NodeStyle,
    add_caption,
    add_glued_connector,
    add_native_chart_card,
    add_node,
    add_panel,
    add_picture_card,
    add_slide_header,
    new_presentation,
    pick_contrast_text_rgb,
    save_presentation,
    tint,
)
from python_figure_helpers import (  # noqa: E402
    prepare_figure_dir,
    save_heatmap,
    save_ranked_bar,
    save_timeline_barh,
)

FONT_NAME = "Aptos"
BRIEF_PATH = WORKSPACE_DIR / "brief.md"
NARRATIVE_PATH = WORKSPACE_DIR / "deck_narrative.md"
PPTX_PATH = WORKSPACE_DIR / "build" / "pptx" / "ai_market_intelligence_demo.pptx"
ICON_THEME_DIR = WORKSPACE_DIR / "assets" / "icons" / "theme_consulting_light"
GENERATED_DIR = WORKSPACE_DIR / "build" / "rendered" / "generated"
GENERATED_SPEC_PATH = WORKSPACE_DIR / "build" / "generated" / "slide_specs.yaml"
MANIFEST_PATH = WORKSPACE_DIR / "validation" / "manifests" / "build_manifest.json"
SLIDE_WIDTH = 16.0
CONTENT_LEFT = 0.78
CONTENT_RIGHT = 15.22
CONTENT_WIDTH = CONTENT_RIGHT - CONTENT_LEFT
GRID_GAP = 0.28
FONT_TITLE = 20
FONT_SUBTITLE = 12.5
FONT_BODY = 14
FONT_LABEL = 12
FONT_BADGE = 11
FONT_CAPTION = 12
FONT_METRIC = 24
FONT_SECTION = 13
FONT_CALLOUT = 12
FONT_TAG = 10

PALETTE = {
    "bg": (248, 250, 252),
    "surface": (255, 255, 255),
    "ink": (15, 23, 42),
    "title": (23, 37, 66),
    "subtitle": (89, 104, 128),
    "line": (203, 213, 225),
    "blue": (37, 99, 235),
    "teal": (13, 148, 136),
    "emerald": (5, 150, 105),
    "amber": (245, 158, 11),
    "violet": (124, 58, 237),
    "rose": (225, 29, 72),
    "slate": (100, 116, 139),
    "good": (5, 150, 105),
    "warn": (180, 83, 9),
    "bad": (220, 38, 38),
}


def add_slide_header(slide, figure_tag: str, title: str, subtitle: str) -> None:
    """为当前 demo 添加左右等边距的标题头。"""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*PALETTE["bg"])

    title_box = slide.shapes.add_textbox(Inches(CONTENT_LEFT), Inches(0.16), Inches(CONTENT_WIDTH), Inches(0.72))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    style_paragraph(title_para, font_size=FONT_TITLE, color_rgb=PALETTE["title"], bold=True)

    subtitle_box = slide.shapes.add_textbox(
        Inches(CONTENT_LEFT),
        Inches(0.82),
        Inches(CONTENT_WIDTH),
        Inches(0.24),
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    style_paragraph(subtitle_para, font_size=FONT_SUBTITLE, color_rgb=PALETTE["subtitle"])

    tag_width = 0.50
    tag_left = CONTENT_RIGHT - tag_width
    tag_box = slide.shapes.add_textbox(Inches(tag_left), Inches(8.43), Inches(tag_width), Inches(0.18))
    tag_frame = tag_box.text_frame
    tag_frame.clear()
    tag_para = tag_frame.paragraphs[0]
    tag_para.text = figure_tag
    style_paragraph(tag_para, font_size=FONT_TAG, color_rgb=PALETTE["subtitle"], bold=True, align=PP_ALIGN.RIGHT)


def run_command(command: list[str]) -> None:
    """执行外部命令，并在失败时直接暴露 stderr。"""
    completed = subprocess.run(command, capture_output=True, text=True)
    if completed.returncode != 0:
        raise RuntimeError(
            "命令执行失败:\n"
            f"cmd={' '.join(command)}\n"
            f"stdout={completed.stdout}\n"
            f"stderr={completed.stderr}"
        )


def load_slide_specs() -> dict:
    """从总叙事文档派生并读取 slide specs。"""
    if not BRIEF_PATH.exists():
        raise FileNotFoundError(f"缺少 brief.md: {BRIEF_PATH}")
    if not NARRATIVE_PATH.exists():
        raise FileNotFoundError(f"缺少 deck_narrative.md: {NARRATIVE_PATH}")

    run_command(
        [
            sys.executable,
            str(SKILL_ROOT / "scripts" / "derive_slide_specs_from_narrative.py"),
            "--narrative",
            str(NARRATIVE_PATH),
            "--out-yaml",
            str(GENERATED_SPEC_PATH),
        ]
    )
    return yaml.safe_load(GENERATED_SPEC_PATH.read_text(encoding="utf-8"))


def read_csv(name: str) -> pd.DataFrame:
    """读取 processed data 下的 CSV。"""
    return pd.read_csv(WORKSPACE_DIR / "data" / "processed" / name)


def style_paragraph(
    paragraph,
    *,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """统一设置段落样式。"""
    paragraph.alignment = align
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = RGBColor(*color_rgb)
    paragraph.font.bold = bold
    paragraph.font.name = FONT_NAME


def add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """添加单段文本框。"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    style_paragraph(paragraph, font_size=font_size, color_rgb=color_rgb, bold=bold, align=align)


def add_multiline_textbox(
    slide,
    lines: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold_first: bool = False,
    gap_before: float = 5.0,
) -> None:
    """添加多段文本框。"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        if index > 0:
            paragraph.space_before = Pt(gap_before)
        style_paragraph(
            paragraph,
            font_size=font_size,
            color_rgb=color_rgb,
            bold=bold_first and index == 0,
        )


def add_badge(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
) -> None:
    """添加小型 badge。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*fill_rgb)
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    style_paragraph(
        paragraph,
        font_size=FONT_BADGE,
        color_rgb=pick_contrast_text_rgb(fill_rgb),
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_header_badge(
    slide,
    text: str,
    width: float,
    fill_rgb: tuple[int, int, int],
    *,
    top: float = 0.94,
    height: float = 0.22,
) -> None:
    """把页面级别的 badge 固定放到标题区右侧，避免和左侧 panel label 冲突。"""
    left = CONTENT_RIGHT - width
    add_badge(slide, text, left, top, width, height, fill_rgb)


def add_icon(slide, source_name: str, left: float, top: float, size: float) -> None:
    """按 source_name 插入已经渲染好的主题 icon。"""
    icon_path = ICON_THEME_DIR / f"{source_name}.png"
    if not icon_path.exists():
        raise FileNotFoundError(f"缺少主题化 icon: {icon_path}")
    slide.shapes.add_picture(
        str(icon_path),
        Inches(left),
        Inches(top),
        width=Inches(size),
        height=Inches(size),
    )


def add_metric_card(
    slide,
    *,
    icon_name: str,
    eyebrow: str,
    value: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加指标卡片。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*PALETTE["surface"])
    card.line.color.rgb = RGBColor(*tint(accent_rgb, 0.62))
    card.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.18, top + 0.16, 0.28)
    add_textbox(
        slide,
        eyebrow,
        left + 0.52,
        top + 0.14,
        width - 0.7,
        0.22,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["subtitle"],
        bold=True,
    )
    add_textbox(
        slide,
        value,
        left + 0.18,
        top + 0.46,
        width - 0.35,
        0.36,
        font_size=FONT_METRIC,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    add_multiline_textbox(
        slide,
        [body],
        left + 0.18,
        top + 0.90,
        width - 0.35,
        height - 1.02,
        font_size=FONT_BODY,
        color_rgb=PALETTE["subtitle"],
    )


def add_insight_chip(
    slide,
    *,
    icon_name: str,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加底部 insight chip。"""
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(*tint(accent_rgb, 0.10))
    chip.line.color.rgb = RGBColor(*tint(accent_rgb, 0.65))
    chip.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.16, top + 0.16, 0.24)
    add_textbox(
        slide,
        title,
        left + 0.48,
        top + 0.12,
        width - 0.6,
        0.22,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    add_multiline_textbox(
        slide,
        [body],
        left + 0.16,
        top + 0.42,
        width - 0.3,
        height - 0.5,
        font_size=FONT_BODY,
        color_rgb=PALETTE["subtitle"],
    )


def add_callout_box(
    slide,
    *,
    icon_name: str,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加适合窄侧栏的简洁 callout。"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*PALETTE["surface"])
    box.line.color.rgb = RGBColor(*tint(accent_rgb, 0.65))
    box.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.12, top + 0.15, 0.22)
    add_textbox(
        slide,
        title,
        left + 0.40,
        top + 0.14,
        width - 0.52,
        0.24,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    if body:
        add_multiline_textbox(
            slide,
            [body],
            left + 0.12,
            top + 0.40,
            width - 0.22,
            height - 0.46,
            font_size=FONT_CALLOUT,
            color_rgb=PALETTE["subtitle"],
        )


def score_fill(score: int) -> tuple[int, int, int]:
    """根据评分返回单元格底色。"""
    mapping = {
        1: (254, 226, 226),
        2: (254, 243, 199),
        3: (224, 231, 255),
        4: (204, 251, 241),
        5: (187, 247, 208),
    }
    return mapping[score]


def ensure_theme_icons() -> None:
    """把 curated icon 以当前 deck 主题渲染到 workspace。"""
    ICON_THEME_DIR.mkdir(parents=True, exist_ok=True)
    command = [
        sys.executable,
        str(SKILL_ROOT / "scripts" / "icon_registry.py"),
        "render",
        "--size",
        "128",
        "--color-mode",
        "auto",
        "--background-color",
        "#F8FAFC",
        "--accent-color",
        "#2563EB",
        "--out-dir",
        str(ICON_THEME_DIR),
    ]
    run_command(command)


def generate_hero_image() -> Path:
    """生成封面使用的抽象 AI 市场 hero image。"""
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    output_path = GENERATED_DIR / "hero_ai_market.png"

    width, height = 2200, 1400
    image = Image.new("RGBA", (width, height), (*PALETTE["bg"], 255))
    draw = ImageDraw.Draw(image, "RGBA")

    top_color = (241, 245, 249)
    bottom_color = (226, 232, 240)
    for y in range(height):
        ratio = y / max(1, height - 1)
        r = round(top_color[0] * (1 - ratio) + bottom_color[0] * ratio)
        g = round(top_color[1] * (1 - ratio) + bottom_color[1] * ratio)
        b = round(top_color[2] * (1 - ratio) + bottom_color[2] * ratio)
        draw.line((0, y, width, y), fill=(r, g, b, 255))

    accent_blobs = [
        ((250, 150, 980, 920), (*PALETTE["blue"], 42)),
        ((720, 60, 1700, 980), (*PALETTE["teal"], 38)),
        ((1380, 120, 2140, 880), (*PALETTE["amber"], 34)),
        ((980, 620, 2080, 1380), (*PALETTE["violet"], 34)),
    ]
    for box, fill in accent_blobs:
        draw.ellipse(box, fill=fill)

    node_centers = [
        (420, 1050),
        (760, 860),
        (1110, 720),
        (1480, 560),
        (1820, 420),
        (1730, 980),
        (1290, 1060),
    ]
    for start, end in zip(node_centers[:-1], node_centers[1:]):
        draw.line((*start, *end), fill=(51, 65, 85, 86), width=10)
    for index, center in enumerate(node_centers):
        radius = 34 if index not in {2, 4} else 44
        color = [PALETTE["blue"], PALETTE["teal"], PALETTE["amber"], PALETTE["violet"]][index % 4]
        draw.ellipse(
            (center[0] - radius, center[1] - radius, center[0] + radius, center[1] + radius),
            fill=(*color, 210),
            outline=(255, 255, 255, 220),
            width=8,
        )

    for x in range(1050, 2120, 130):
        draw.rounded_rectangle(
            (x, 1030, x + 84, 1118),
            radius=18,
            fill=(255, 255, 255, 116),
            outline=(148, 163, 184, 120),
            width=3,
        )
    for x in range(1120, 2040, 130):
        draw.rounded_rectangle(
            (x, 1160, x + 84, 1248),
            radius=18,
            fill=(255, 255, 255, 98),
            outline=(148, 163, 184, 110),
            width=3,
        )

    blurred = image.filter(ImageFilter.GaussianBlur(radius=0.6))
    blurred.convert("RGB").save(output_path)
    return output_path


def generate_python_figures() -> dict[str, Path]:
    """生成 ranked bar、heatmap 与 timeline 三类 Python figure。"""
    figure_dir = prepare_figure_dir(WORKSPACE_DIR)

    ranked_path = figure_dir / "enterprise_budget_rank.png"
    save_ranked_bar(
        output_path=ranked_path,
        data=read_csv("enterprise_budget_rank.csv"),
        label_col="Workflow",
        value_col="Budget_Bn",
        accent_rgb=PALETTE["amber"],
        title="Enterprise AI budget by workflow",
        x_label="$B",
        figsize=(6.0, 2.6),
    )

    heatmap_frame = read_csv("sector_adoption_heatmap.csv").set_index("Sector")
    heatmap_path = figure_dir / "sector_adoption_heatmap.png"
    save_heatmap(
        output_path=heatmap_path,
        frame=heatmap_frame,
        accent_rgb=PALETTE["violet"],
        title="Adoption intensity by sector",
        vmin=3.0,
        vmax=5.0,
        figsize=(6.0, 2.6),
    )

    timeline_path = figure_dir / "capability_rollout_timeline.png"
    save_timeline_barh(
        output_path=timeline_path,
        data=read_csv("capability_rollout_timeline.csv"),
        label_col="Wave",
        start_col="Start_Quarter",
        duration_col="Duration_Quarters",
        accent_rgb=PALETTE["teal"],
        title="Suggested market-entry path",
        x_label="Quarter from launch",
        figsize=(4.4, 2.2),
    )

    return {
        "ranked_bar": ranked_path,
        "heatmap": heatmap_path,
        "timeline": timeline_path,
    }


def add_matrix_table(slide) -> None:
    """构建策略 archetype 对比矩阵。"""
    data = read_csv("strategy_archetype_scores.csv")
    left = CONTENT_LEFT
    top = 1.22
    width = 9.80
    height = 6.35
    accent_rgb = PALETTE["blue"]

    add_panel(slide, "Archetype scorecard", left, top, width, height, accent_rgb, body_fill_rgb=(255, 255, 255))

    columns = ["Archetype", "Growth", "Margin", "Moat", "Capital Efficiency", "Time to Scale"]
    table_left = left + 0.14
    table_top = top + 0.55
    row_height = 0.84
    col_widths = [2.38, 1.02, 1.02, 0.94, 1.40, 1.48]

    cursor_x = table_left
    for index, column in enumerate(columns):
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(cursor_x),
            Inches(table_top),
            Inches(col_widths[index]),
            Inches(0.52),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*tint(accent_rgb, 0.12))
        header.line.color.rgb = RGBColor(*tint(accent_rgb, 0.55))
        header.line.width = Pt(0.8)
        paragraph = header.text_frame.paragraphs[0]
        paragraph.text = column
        style_paragraph(paragraph, font_size=FONT_LABEL, color_rgb=PALETTE["title"], bold=True, align=PP_ALIGN.CENTER)
        cursor_x += col_widths[index]

    for row_index, row in data.iterrows():
        row_top = table_top + 0.62 + row_index * row_height
        is_recommended = row["Recommendation"] == "Recommended"
        row_fill = tint(PALETTE["teal"], 0.12) if is_recommended else (255, 255, 255)

        name_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(table_left),
            Inches(row_top),
            Inches(col_widths[0]),
            Inches(0.68),
        )
        name_shape.fill.solid()
        name_shape.fill.fore_color.rgb = RGBColor(*row_fill)
        name_shape.line.color.rgb = RGBColor(*(PALETTE["teal"] if is_recommended else PALETTE["line"]))
        name_shape.line.width = Pt(1.0 if is_recommended else 0.8)

        add_icon(slide, row["Icon"], table_left + 0.10, row_top + 0.17, 0.22)
        add_textbox(
            slide,
            row["Archetype"],
            table_left + 0.38,
            row_top + 0.12,
            1.45,
            0.20,
            font_size=FONT_LABEL,
            color_rgb=PALETTE["title"],
            bold=True,
        )
        if is_recommended:
            add_badge(slide, "Pick", table_left + 1.62, row_top + 0.20, 0.44, 0.17, PALETTE["teal"])

        cursor_x = table_left + col_widths[0]
        for metric in columns[1:]:
            score = int(row[metric])
            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cursor_x + 0.08),
                Inches(row_top + 0.08),
                Inches(col_widths[columns.index(metric)] - 0.16),
                Inches(0.52),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*score_fill(score))
            cell.line.color.rgb = RGBColor(*tint(PALETTE["slate"], 0.45))
            cell.line.width = Pt(0.6)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.text = str(score)
            style_paragraph(paragraph, font_size=FONT_BODY, color_rgb=PALETTE["title"], bold=True, align=PP_ALIGN.CENTER)
            cursor_x += col_widths[columns.index(metric)]

    add_multiline_textbox(
        slide,
        [
            "Score meaning: 5 = strongest strategic position, 1 = structurally weak.",
            "Capital Efficiency treats a lighter capital model as more attractive.",
        ],
        left + 0.16,
        top + 5.96,
        width - 0.3,
        0.42,
        font_size=FONT_CAPTION,
        color_rgb=PALETTE["subtitle"],
    )


def build_slide_01(slide, spec: dict, hero_path: Path) -> None:
    """构建封面判断页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Executive View", 1.22, PALETTE["blue"])

    add_multiline_textbox(
        slide,
        [
            "The next AI winners will monetize workflows, not just model access.",
            "Supply-side capex remains necessary, but durable profit pools accrue where inference, domain data, governance and distribution reinforce each other.",
        ],
        CONTENT_LEFT,
        1.38,
        7.00,
        1.35,
        font_size=18,
        color_rgb=PALETTE["title"],
        bold_first=True,
        gap_before=8.0,
    )

    add_picture_card(
        slide,
        title="Where value is moving",
        image_path=hero_path,
        left=8.18,
        top=1.20,
        width=7.04,
        height=5.92,
        accent_rgb=PALETTE["blue"],
        caption="Generated locally inside the workspace to cover the image-hero route.",
    )

    summary = read_csv("summary_cards.csv")
    card_positions = [
        (CONTENT_LEFT, 3.22),
        (4.18, 3.22),
        (CONTENT_LEFT, 5.00),
        (4.18, 5.00),
    ]
    accents = [PALETTE["blue"], PALETTE["teal"], PALETTE["amber"], PALETTE["violet"]]
    for index, row in summary.iterrows():
        left, top = card_positions[index]
        add_metric_card(
            slide,
            icon_name=row["Icon"],
            eyebrow=row["Title"],
            value=row["Value"],
            body=row["Body"],
            left=left,
            top=top,
            width=3.24,
            height=1.62,
            accent_rgb=accents[index],
        )

    add_caption(slide, "Illustrative market view for skill demo only. Data is fabricated but internally consistent.")


def build_slide_02(slide, spec: dict) -> int:
    """构建 native chart 证据页，并返回原生图表数量。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Market Economics", 1.52, PALETTE["teal"])

    revenue = read_csv("market_layer_revenue.csv")
    chart_left = add_native_chart_card(
        slide=slide,
        title="Revenue by stack layer, $B",
        left=CONTENT_LEFT,
        top=1.24,
        width=7.12,
        height=3.48,
        accent_rgb=PALETTE["blue"],
        categories=[str(year) for year in revenue["Year"]],
        series_list=[
            ("Compute + Models", revenue["Compute_Models"].tolist()),
            ("Tools + Orchestration", revenue["Tools_Orchestration"].tolist()),
            ("Applications + Services", revenue["Applications_Services"].tolist()),
        ],
        chart_type=XL_CHART_TYPE.LINE,
        show_legend=True,
        legend_position=XL_LEGEND_POSITION.BOTTOM,
        series_colors=[PALETTE["blue"], PALETTE["teal"], PALETTE["amber"]],
    )
    chart_left.plots[0].has_data_labels = False
    chart_left.value_axis.maximum_scale = 470
    chart_left.value_axis.minimum_scale = 0
    chart_left.value_axis.major_unit = 100

    mix = read_csv("profit_pool_mix.csv")
    chart_right = add_native_chart_card(
        slide=slide,
        title="EBIT pool migration, % of industry EBIT",
        left=8.18,
        top=1.24,
        width=7.04,
        height=3.48,
        accent_rgb=PALETTE["teal"],
        categories=mix["Year"].tolist(),
        series_list=[
            ("Infra + Models", mix["Infra_Models"].tolist()),
            ("Tools + Control", mix["Tools_Control"].tolist()),
            ("Apps + Services", mix["Apps_Services"].tolist()),
        ],
        chart_type=XL_CHART_TYPE.BAR_STACKED,
        number_format="0%",
        show_legend=True,
        legend_position=XL_LEGEND_POSITION.BOTTOM,
        series_colors=[PALETTE["blue"], PALETTE["violet"], PALETTE["amber"]],
    )
    chart_right.value_axis.maximum_scale = 1.0
    chart_right.value_axis.minimum_scale = 0.0
    chart_right.value_axis.major_unit = 0.2

    chips = [
        ("chart-line", "Revenue CAGR 41%", "Revenue expands around 6x as pilots turn into funded programs.", PALETTE["blue"]),
        ("briefcase", "Apps/services EBIT share 56%", "Profit shifts toward recurring workflow software and services.", PALETTE["teal"]),
        ("brain", "Cost curves unlock new demand", "Cheaper inference lets more workflows clear ROI.", PALETTE["amber"]),
    ]
    for index, (icon_name, title, body, accent) in enumerate(chips):
        add_insight_chip(
            slide,
            icon_name=icon_name,
            title=title,
            body=body,
            left=CONTENT_LEFT + index * 4.62,
            top=5.02,
            width=4.30,
            height=1.52,
            accent_rgb=accent,
        )

    add_caption(slide, "Editable native charts are used intentionally here to preserve post-handoff data editing.")
    return 2


def build_slide_03(slide, spec: dict, figure_paths: dict[str, Path]) -> int:
    """构建 Python figure 研究表达页，并返回 Python figure 数量。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Demand Shape", 1.15, PALETTE["amber"])

    add_picture_card(
        slide,
        title="2026 enterprise AI budget by workflow, $B",
        image_path=figure_paths["ranked_bar"],
        left=CONTENT_LEFT,
        top=1.24,
        width=7.12,
        height=3.42,
        accent_rgb=PALETTE["amber"],
        caption="Budget authority is concentrated in workflows with labor substitution or risk ownership.",
    )
    add_picture_card(
        slide,
        title="Adoption intensity by sector and use case",
        image_path=figure_paths["heatmap"],
        left=8.18,
        top=1.24,
        width=7.04,
        height=3.42,
        accent_rgb=PALETTE["violet"],
        caption="Adoption is broad, but the monetization path differs sharply by sector and workflow owner.",
    )

    add_panel(
        slide,
        "Why spending concentrates",
        CONTENT_LEFT,
        4.92,
        CONTENT_WIDTH,
        2.02,
        PALETTE["teal"],
        body_fill_rgb=(255, 255, 255),
    )
    add_multiline_textbox(
        slide,
        [
            "1. Budget-rich functions buy first: support, software delivery, risk and sales productivity have clear owners and measurable ROI.",
            "2. Cross-functional use cases adopt broadly, but broad adoption does not automatically create a monetization wedge.",
            "3. This is why the winning product is usually a workflow package with integration, guardrails and change management baked in.",
        ],
        0.98,
        5.40,
        13.90,
        1.20,
        font_size=FONT_BODY,
        color_rgb=PALETTE["title"],
    )

    add_caption(slide, "This page uses Python figure assets to preserve information density that would be awkward in native Office charts.")
    return 2


def build_slide_04(slide, spec: dict) -> int:
    """构建 connector-based 行业价值链页，并返回 connector 数量。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Control Architecture", 1.60, PALETTE["violet"])

    add_panel(slide, "AI industry control map", CONTENT_LEFT, 1.18, 10.46, 6.45, PALETTE["blue"], body_fill_rgb=(255, 255, 255))
    add_panel(slide, "Economic implications", 11.52, 1.18, 3.70, 6.45, PALETTE["rose"], body_fill_rgb=(255, 255, 255))

    node_style_main = {
        "compute": NodeStyle(fill_rgb=tint(PALETTE["blue"], 0.18), line_rgb=tint(PALETTE["blue"], 0.72), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        "model": NodeStyle(fill_rgb=tint(PALETTE["teal"], 0.18), line_rgb=tint(PALETTE["teal"], 0.72), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        "control": NodeStyle(fill_rgb=tint(PALETTE["violet"], 0.18), line_rgb=tint(PALETTE["violet"], 0.72), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        "apps": NodeStyle(fill_rgb=tint(PALETTE["amber"], 0.18), line_rgb=tint(PALETTE["amber"], 0.72), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        "channel": NodeStyle(fill_rgb=tint(PALETTE["rose"], 0.16), line_rgb=tint(PALETTE["rose"], 0.72), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
    }

    add_icon(slide, "database", 1.24, 1.88, 0.26)
    add_icon(slide, "brain", 3.12, 1.88, 0.26)
    add_icon(slide, "circles-relation", 5.03, 1.88, 0.26)
    add_icon(slide, "briefcase", 6.95, 1.88, 0.26)
    add_icon(slide, "presentation", 8.86, 1.88, 0.26)

    nodes = {
        "compute": add_node(slide, "Compute Fabric", "Compute\nFabric", 1.10, 2.20, 1.38, 0.84, node_style_main["compute"]),
        "models": add_node(slide, "Foundation Models", "Foundation\nModels", 2.98, 2.20, 1.44, 0.84, node_style_main["model"]),
        "control": add_node(slide, "Agent Control Layer", "Agent Control\nLayer", 4.95, 2.20, 1.60, 0.84, node_style_main["control"]),
        "apps": add_node(slide, "Vertical Workflow Apps", "Vertical Workflow\nApps", 7.00, 2.20, 1.64, 0.84, node_style_main["apps"]),
        "channel": add_node(slide, "Enterprise Distribution", "Enterprise\nDistribution", 8.98, 2.20, 1.52, 0.84, node_style_main["channel"]),
        "data": add_node(
            slide,
            "Proprietary Data",
            "Proprietary Data",
            4.42,
            4.28,
            1.70,
            0.68,
            NodeStyle(fill_rgb=tint(PALETTE["blue"], 0.10), line_rgb=tint(PALETTE["blue"], 0.55), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        ),
        "safety": add_node(
            slide,
            "Safety Governance",
            "Safety & Governance",
            6.62,
            4.28,
            1.72,
            0.68,
            NodeStyle(fill_rgb=tint(PALETTE["emerald"], 0.14), line_rgb=tint(PALETTE["emerald"], 0.62), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        ),
        "services": add_node(
            slide,
            "Implementation Services",
            "Implementation Services",
            2.86,
            5.42,
            1.90,
            0.68,
            NodeStyle(fill_rgb=tint(PALETTE["amber"], 0.12), line_rgb=tint(PALETTE["amber"], 0.58), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        ),
        "feedback": add_node(
            slide,
            "User Feedback Telemetry",
            "User Feedback\nTelemetry",
            8.22,
            5.42,
            1.96,
            0.74,
            NodeStyle(fill_rgb=tint(PALETTE["violet"], 0.12), line_rgb=tint(PALETTE["violet"], 0.58), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE),
        ),
    }

    for shape in nodes.values():
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.CENTER

    connectors = [
        add_glued_connector(slide, nodes["compute"], nodes["models"], "right", "left", PALETTE["slate"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["models"], nodes["control"], "right", "left", PALETTE["slate"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["control"], nodes["apps"], "right", "left", PALETTE["slate"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["apps"], nodes["channel"], "right", "left", PALETTE["slate"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["data"], nodes["control"], "top", "bottom", PALETTE["blue"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["data"], nodes["apps"], "right", "bottom", PALETTE["blue"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["safety"], nodes["control"], "left", "bottom", PALETTE["emerald"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["safety"], nodes["channel"], "top", "bottom", PALETTE["emerald"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["services"], nodes["apps"], "top", "bottom", PALETTE["amber"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["feedback"], nodes["models"], "top", "bottom", PALETTE["violet"], MSO_CONNECTOR.ELBOW),
        add_glued_connector(slide, nodes["channel"], nodes["feedback"], "bottom", "top", PALETTE["rose"], MSO_CONNECTOR.ELBOW),
    ]
    for connector in connectors:
        connector.line.width = Pt(1.55)

    add_multiline_textbox(
        slide,
        [
            "1. Models stay essential, but control economics sit one layer downstream where products orchestrate tools, policies and data.",
            "2. Proprietary data and governance make workflows sticky; distribution and user telemetry keep improving the product after sale.",
            "3. Integrators monetize deployment, but the compounding moat usually belongs to the product layer that owns the feedback loop.",
        ],
        11.70,
        1.74,
        3.10,
        2.25,
        font_size=FONT_BODY,
        color_rgb=PALETTE["title"],
    )

    implication_cards = [
        ("chart-line", "Capital stays upstream", "", PALETTE["blue"]),
        ("briefcase", "Profit shifts downstream", "", PALETTE["amber"]),
        ("shield-check", "Moat compounds in governance", "", PALETTE["emerald"]),
    ]
    for index, (icon_name, title, body, accent) in enumerate(implication_cards):
        add_callout_box(
            slide,
            icon_name=icon_name,
            title=title,
            body=body,
            left=11.70,
            top=4.20 + index * 0.92,
            width=3.06,
            height=0.66,
            accent_rgb=accent,
        )

    add_caption(slide, "This page intentionally uses real glued connectors so the diagram remains maintainable after handoff.")
    return len(connectors)


def build_slide_05(slide, spec: dict, figure_paths: dict[str, Path]) -> int:
    """构建 recommendation page，并返回新增 Python figure 数量。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Recommendation", 1.28, PALETTE["teal"])

    add_matrix_table(slide)

    add_panel(slide, "Investment thesis", 10.86, 1.22, 4.36, 2.12, PALETTE["teal"], body_fill_rgb=(255, 255, 255))
    add_multiline_textbox(
        slide,
        [
            "The best play is a domain-focused agent stack with a narrow wedge, clear budget owner and expansion path into adjacent workflows.",
            "That model captures software-like margin while still benefiting from falling model cost and rising enterprise adoption.",
        ],
        11.08,
        1.70,
        3.94,
        1.10,
        font_size=FONT_BODY,
        color_rgb=PALETTE["title"],
    )

    add_picture_card(
        slide,
        title="Suggested three-wave entry plan",
        image_path=figure_paths["timeline"],
        left=10.86,
        top=3.56,
        width=4.36,
        height=3.88,
        accent_rgb=PALETTE["teal"],
        caption="Start with a narrow copilot wedge, then expand into agentic workflow execution and workflow-network adjacencies.",
    )

    add_caption(slide, "The matrix is intentionally native-shape based, while the rollout path uses a Python figure to preserve density and clarity.")
    return 1


def add_workflow_step(
    slide,
    *,
    icon_name: str,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加 slide 6 左侧流程步骤卡片。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*PALETTE["surface"])
    card.line.color.rgb = RGBColor(*tint(accent_rgb, 0.66))
    card.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.12, top + 0.14, 0.24)
    add_textbox(
        slide,
        title,
        left + 0.42,
        top + 0.10,
        width - 0.5,
        0.20,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    add_multiline_textbox(
        slide,
        [body],
        left + 0.12,
        top + 0.44,
        width - 0.24,
        height - 0.52,
        font_size=FONT_BODY,
        color_rgb=PALETTE["subtitle"],
    )


def add_workflow_row(
    slide,
    *,
    icon_name: str,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加宽行工作流条目。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*PALETTE["surface"])
    card.line.color.rgb = RGBColor(*tint(accent_rgb, 0.66))
    card.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.12, top + 0.13, 0.22)
    add_textbox(
        slide,
        title,
        left + 0.42,
        top + 0.10,
        1.20,
        0.18,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    add_textbox(
        slide,
        body,
        left + 1.62,
        top + 0.10,
        width - 1.78,
        0.22,
        font_size=FONT_BODY,
        color_rgb=PALETTE["subtitle"],
    )


def add_coverage_card(
    slide,
    *,
    icon_name: str,
    title: str,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """添加 slide 6 右侧能力卡片。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*tint(accent_rgb, 0.10))
    card.line.color.rgb = RGBColor(*tint(accent_rgb, 0.64))
    card.line.width = Pt(1.0)

    add_icon(slide, icon_name, left + 0.12, top + 0.12, 0.24)
    add_textbox(
        slide,
        title,
        left + 0.42,
        top + 0.09,
        width - 0.5,
        0.20,
        font_size=FONT_LABEL,
        color_rgb=PALETTE["title"],
        bold=True,
    )
    add_multiline_textbox(
        slide,
        [body],
        left + 0.12,
        top + 0.40,
        width - 0.22,
        height - 0.46,
        font_size=FONT_BODY,
        color_rgb=PALETTE["subtitle"],
    )


def build_slide_06(slide, spec: dict, metrics: dict[str, int]) -> None:
    """构建 capability showcase 页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"])
    add_header_badge(slide, "Capability Proof", 1.42, PALETTE["blue"])

    add_panel(slide, "Reproducible deck workflow", CONTENT_LEFT, 1.20, 8.58, 5.90, PALETTE["blue"], body_fill_rgb=(255, 255, 255))
    add_panel(slide, "Asset coverage in this demo", 9.64, 1.20, 5.58, 5.90, PALETTE["teal"], body_fill_rgb=(255, 255, 255))

    workflow_steps = [
        ("search", "Brief", "Lock audience and message.", PALETTE["blue"]),
        ("presentation", "Plan", "Set archetype and validation.", PALETTE["teal"]),
        ("database", "Data", "Keep inputs stable and traceable.", PALETTE["amber"]),
        ("chart-line", "Build", "Generate charts, icons and slides.", PALETTE["violet"]),
        ("shield-check", "Validate", "Check connectors and previews.", PALETTE["emerald"]),
        ("rocket", "Deliver", "Ship a board-ready workspace.", PALETTE["rose"]),
    ]
    workflow_positions = [2.14, 2.88, 3.62, 4.36, 5.10, 5.84]
    for index, (icon_name, title, body, accent) in enumerate(workflow_steps):
        add_workflow_row(
            slide,
            icon_name=icon_name,
            title=title,
            body=body,
            left=0.98,
            top=workflow_positions[index],
            width=7.62,
            height=0.58,
            accent_rgb=accent,
        )

    add_multiline_textbox(
        slide,
        [
            "The value is not just automated slide drawing.",
            "Planning, asset routing, editable build and validation stay preserved in one workspace.",
        ],
        0.96,
        6.34,
        8.04,
        0.46,
        font_size=FONT_BODY,
        color_rgb=PALETTE["title"],
        bold_first=True,
        gap_before=8.0,
    )

    coverage_cards = [
        ("presentation", "Editable layout", "Text and shapes remain editable after handoff.", PALETTE["blue"]),
        ("chart-line", "Native charts", "Revenue and EBIT pages stay editable.", PALETTE["teal"]),
        ("circles-relation", "Connector diagram", "The industry map keeps real glued connectors.", PALETTE["violet"]),
        ("stack-2", "Python figures", "Research figures keep density and clarity.", PALETTE["amber"]),
        ("shield-check", "Validation bundle", "Manifest, previews and connector checks stay together.", PALETTE["emerald"]),
        ("route", "Reusable workflow", "The workspace is built for iteration.", PALETTE["rose"]),
    ]
    for index, card in enumerate(coverage_cards):
        row = index // 2
        col = index % 2
        add_coverage_card(
            slide,
            icon_name=card[0],
            title=card[1],
            body=card[2],
            left=9.92 + col * 2.68,
            top=2.00 + row * 1.46,
            width=2.54,
            height=1.22,
            accent_rgb=card[3],
        )

    stat_cards = [
        ("Slides", str(metrics["slides"]), PALETTE["blue"]),
        ("Native charts", str(metrics["native_charts"]), PALETTE["teal"]),
        ("Python figs", str(metrics["python_figures"]), PALETTE["amber"]),
        ("Conns", str(metrics["connectors"]), PALETTE["violet"]),
        ("Icons", "1", PALETTE["rose"]),
    ]
    for index, (label, value, accent) in enumerate(stat_cards):
        left = CONTENT_LEFT + index * 2.89
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(7.42),
            Inches(2.68),
            Inches(0.76),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*tint(accent, 0.12))
        card.line.color.rgb = RGBColor(*tint(accent, 0.65))
        card.line.width = Pt(1.0)
        add_textbox(slide, label, left + 0.12, 7.56, 1.18, 0.16, font_size=FONT_LABEL, color_rgb=PALETTE["subtitle"], bold=True)
        add_textbox(slide, value, left + 1.40, 7.46, 0.96, 0.24, font_size=FONT_METRIC, color_rgb=PALETTE["title"], bold=True, align=PP_ALIGN.RIGHT)

    add_caption(slide, "This final page turns the demo itself into proof that the skill can deliver both insight and production workflow.")


def build_deck() -> tuple[Path, dict]:
    """构建整套 deck，并返回输出路径与 manifest。"""
    specs = load_slide_specs()
    ensure_theme_icons()
    hero_path = generate_hero_image()
    figure_paths = generate_python_figures()

    presentation = new_presentation()
    slide_map = {slide_spec["slide_id"]: slide_spec for slide_spec in specs["slides"]}

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_slide_01(slide, slide_map["S01"], hero_path)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    native_chart_count = build_slide_02(slide, slide_map["S02"])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    python_figure_count = build_slide_03(slide, slide_map["S03"], figure_paths)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    connector_count = build_slide_04(slide, slide_map["S04"])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    python_figure_count += build_slide_05(slide, slide_map["S05"], figure_paths)

    metrics = {
        "slides": 6,
        "native_charts": native_chart_count,
        "python_figures": python_figure_count,
        "connectors": connector_count,
    }
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_slide_06(slide, slide_map["S06"], metrics)

    output_path = save_presentation(presentation, PPTX_PATH)
    manifest = {
        "deck_title": specs["deck"]["title"],
        "brief": str(BRIEF_PATH),
        "narrative": str(NARRATIVE_PATH),
        "derived_slide_specs": str(GENERATED_SPEC_PATH),
        "output_pptx": str(output_path),
        "slide_count": metrics["slides"],
        "native_chart_count": metrics["native_charts"],
        "python_figure_count": metrics["python_figures"],
        "connector_count": metrics["connectors"],
        "generated_assets": {
            "hero_image": str(hero_path),
            "ranked_bar": str(figure_paths["ranked_bar"]),
            "heatmap": str(figure_paths["heatmap"]),
            "timeline": str(figure_paths["timeline"]),
            "icon_theme_dir": str(ICON_THEME_DIR),
        },
    }
    MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path, manifest


def main() -> int:
    """CLI 入口。"""
    output_path, manifest = build_deck()
    print(f"[OK] PPT 已生成: {output_path}")
    print(f"[OK] Manifest 已写入: {MANIFEST_PATH}")
    print(json.dumps(manifest, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
