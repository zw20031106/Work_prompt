#!/usr/bin/env python3
"""生成“标准战不是技术评测” polished deck 的主构建脚本。

定位与作用
----------
本脚本把 `brief.md`、`deck_narrative.md`、结构化数据、主题化 icon 和
页面级版式逻辑统一编译成一份可编辑 `pptx`。它服务的不是单页出图，而是
一个可复跑、可验证、可继续维护的 deck workspace。

大致流程
----------
1. 从 `deck_narrative.md` 派生 `slide_specs.yaml` 并读取 deck 级主题参数。
2. 渲染当前主题下的 icon 资产，并准备数据输入。
3. 逐页构建 12 张 slide，覆盖 hero、decision logic、research note、matrix 和 diagram。
4. 写出 `pptx` 与 `build_manifest.json`，供后续 preview 导出与 connector 校验复用。
"""

from __future__ import annotations

import json
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile
import xml.etree.ElementTree as ET

import pandas as pd
import yaml
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

WORKSPACE_DIR = Path(__file__).resolve().parents[1]
REPO_ROOT = Path(__file__).resolve().parents[3]
SKILL_ROOT = REPO_ROOT / "ppt-polished-deck-collab"

sys.path.insert(0, str(SKILL_ROOT / "scripts"))

from ppt_asset_helpers import (  # noqa: E402
    NodeStyle,
    add_native_chart_card,
    add_picture_card,
    add_glued_connector,
    add_node,
    new_presentation,
    pick_contrast_text_rgb,
    save_presentation,
    tint,
)
from python_figure_helpers import prepare_figure_dir, save_heatmap  # noqa: E402

BRIEF_PATH = WORKSPACE_DIR / "brief.md"
NARRATIVE_PATH = WORKSPACE_DIR / "deck_narrative.md"
GENERATED_SPEC_PATH = WORKSPACE_DIR / "build" / "generated" / "slide_specs.yaml"
ICON_THEME_DIR = WORKSPACE_DIR / "assets" / "icons" / "theme_market_vote"
PPTX_PATH = WORKSPACE_DIR / "build" / "pptx" / "standard_wars_executive_deck.pptx"
MANIFEST_PATH = WORKSPACE_DIR / "validation" / "manifests" / "build_manifest.json"

SLIDE_WIDTH = 16.0
SLIDE_HEIGHT = 9.0
CONTENT_LEFT = 0.78
CONTENT_RIGHT = 15.22
CONTENT_WIDTH = CONTENT_RIGHT - CONTENT_LEFT
BODY_TOP = 1.42
BODY_BOTTOM = 8.28
GRID_GAP = 0.22

FONT_TITLE = "Arial"
FONT_BODY = "Arial"
FONT_LATIN = "Arial"
FONT_EAST_ASIA = "黑体"

PALETTE: dict[str, tuple[int, int, int]] = {
    "bg": (246, 244, 238),
    "surface": (252, 251, 247),
    "surface_alt": (242, 239, 232),
    "ink": (21, 28, 37),
    "title": (28, 36, 49),
    "subtitle": (89, 94, 101),
    "line": (199, 193, 182),
    "muted": (140, 135, 126),
    "blue": (34, 78, 145),
    "warm": (190, 97, 47),
    "sage": (68, 109, 92),
    "gold": (181, 140, 56),
    "rose": (176, 68, 78),
    "good": (64, 126, 98),
    "neutral": (128, 123, 116),
}


@dataclass
class DeckContext:
    """封装 build 过程中复用的 deck 级上下文。"""

    deck_spec: dict[str, Any]
    slides_by_id: dict[str, dict[str, Any]]
    factor_frame: pd.DataFrame
    comparison_frame: pd.DataFrame
    milestones_frame: pd.DataFrame
    home_video_scorecard: pd.DataFrame
    networking_scorecard: pd.DataFrame
    networking_heatmap_path: Path


# ---------------------------------------------------------------------------
# 基础工具
# ---------------------------------------------------------------------------


def run_command(command: list[str]) -> None:
    """执行外部命令，并在失败时正确暴露标准输出与错误输出。"""
    completed = subprocess.run(command, capture_output=True, text=True)
    if completed.returncode != 0:
        raise RuntimeError(
            "命令执行失败:\n"
            f"cmd={' '.join(command)}\n"
            f"stdout={completed.stdout}\n"
            f"stderr={completed.stderr}"
        )


def configure_theme(theme_tokens: dict[str, Any]) -> None:
    """根据 narrative frontmatter 调整主题参数。"""
    global FONT_TITLE, FONT_BODY, FONT_LATIN, FONT_EAST_ASIA, CONTENT_LEFT, CONTENT_RIGHT, CONTENT_WIDTH, PALETTE

    FONT_TITLE = theme_tokens.get("title_font_name", FONT_TITLE)
    FONT_BODY = theme_tokens.get("body_font_name", FONT_BODY)
    FONT_LATIN = theme_tokens.get("latin_font_name", FONT_LATIN)
    FONT_EAST_ASIA = theme_tokens.get("east_asia_font_name", FONT_EAST_ASIA)
    CONTENT_LEFT = float(theme_tokens.get("left_margin_in", CONTENT_LEFT))
    CONTENT_RIGHT = float(theme_tokens.get("right_margin_in", CONTENT_RIGHT))
    CONTENT_WIDTH = CONTENT_RIGHT - CONTENT_LEFT

    background_rgb = tuple(theme_tokens.get("background_rgb", PALETTE["bg"]))
    ink_rgb = tuple(theme_tokens.get("ink_rgb", PALETTE["ink"]))
    accent_rgb = tuple(theme_tokens.get("accent_rgb", PALETTE["blue"]))
    accent_warm_rgb = tuple(theme_tokens.get("accent_warm_rgb", PALETTE["warm"]))
    PALETTE = {
        **PALETTE,
        "bg": background_rgb,
        "ink": ink_rgb,
        "title": ink_rgb,
        "blue": accent_rgb,
        "warm": accent_warm_rgb,
        "surface": tint(background_rgb, 0.72),
        "surface_alt": tint(background_rgb, 0.88),
    }


def ensure_theme_icons() -> None:
    """把当前 deck 需要的 icon 渲染到 workspace。"""
    ICON_THEME_DIR.mkdir(parents=True, exist_ok=True)
    run_command(
        [
            sys.executable,
            str(SKILL_ROOT / "scripts" / "icon_registry.py"),
            "render",
            "--size",
            "128",
            "--color-mode",
            "auto",
            "--background-color",
            rgb_hex(PALETTE["bg"]),
            "--accent-color",
            rgb_hex(PALETTE["blue"]),
            "--out-dir",
            str(ICON_THEME_DIR),
        ]
    )


def load_slide_specs() -> dict[str, Any]:
    """从总叙事文档派生并读取 slide specs。"""
    if not BRIEF_PATH.exists():
        raise FileNotFoundError(f"缺少 brief.md: {BRIEF_PATH}")
    if not NARRATIVE_PATH.exists():
        raise FileNotFoundError(f"缺少 deck_narrative.md: {NARRATIVE_PATH}")

    run_command(
        [
            sys.executable,
            str(SKILL_ROOT / "scripts" / "derive_slide_specs_from_narrative.py"),
            "--narrative",
            str(NARRATIVE_PATH),
            "--out-yaml",
            str(GENERATED_SPEC_PATH),
        ]
    )
    return yaml.safe_load(GENERATED_SPEC_PATH.read_text(encoding="utf-8"))


def load_context() -> DeckContext:
    """加载 slide specs 与外部数据，并初始化主题。"""
    deck_spec = load_slide_specs()
    configure_theme(deck_spec["deck"].get("theme_tokens", {}))
    ensure_theme_icons()

    slides_by_id = {slide["slide_id"]: slide for slide in deck_spec["slides"]}
    factor_frame = pd.read_csv(WORKSPACE_DIR / "data" / "processed" / "five_factor_framework.csv")
    comparison_frame = pd.read_csv(WORKSPACE_DIR / "data" / "processed" / "comparison_matrix.csv")
    milestones_frame = pd.read_csv(WORKSPACE_DIR / "data" / "processed" / "case_milestones.csv")
    home_video_scorecard = pd.read_csv(WORKSPACE_DIR / "data" / "processed" / "home_video_scorecard.csv")
    networking_scorecard = pd.read_csv(WORKSPACE_DIR / "data" / "processed" / "networking_scorecard.csv")
    figure_dir = prepare_figure_dir(WORKSPACE_DIR)
    networking_heatmap_path = figure_dir / "networking_factor_heatmap.png"
    heatmap_frame = pd.DataFrame(
        [[5, 2, 2, 2], [3, 5, 5, 5]],
        index=["OSI", "TCP/IP"],
        columns=["Integrity", "Speed", "Cost", "Ops"],
    )
    save_heatmap(
        output_path=networking_heatmap_path,
        frame=heatmap_frame,
        accent_rgb=PALETTE["blue"],
        title="Protocol adoption score by factor",
        figsize=(4.8, 1.9),
        vmin=1,
        vmax=5,
    )
    return DeckContext(
        deck_spec=deck_spec,
        slides_by_id=slides_by_id,
        factor_frame=factor_frame,
        comparison_frame=comparison_frame,
        milestones_frame=milestones_frame,
        home_video_scorecard=home_video_scorecard,
        networking_scorecard=networking_scorecard,
        networking_heatmap_path=networking_heatmap_path,
    )


def rgb_hex(rgb: tuple[int, int, int]) -> str:
    """把 RGB 三元组转成十六进制颜色。"""
    return "#{:02X}{:02X}{:02X}".format(*rgb)


def apply_font_policy_to_pptx(pptx_path: Path) -> None:
    """对最终 PPTX 做一次字体槽位后处理。

    目标是让所有文本对象显式写入：
    - `latin = Arial`
    - `ea = 黑体`
    - `cs = Arial`

    这样普通文本与 Office chart 都能稳定遵守“中文黑体、英文 Arial”的策略。
    """

    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    text_prop_tags = {
        f"{{{a_ns}}}rPr",
        f"{{{a_ns}}}defRPr",
        f"{{{a_ns}}}endParaRPr",
    }
    latin_tag = f"{{{a_ns}}}latin"
    ea_tag = f"{{{a_ns}}}ea"
    cs_tag = f"{{{a_ns}}}cs"
    sym_tag = f"{{{a_ns}}}sym"
    font_tag = f"{{{a_ns}}}font"
    hlink_click_tag = f"{{{a_ns}}}hlinkClick"
    hlink_hover_tag = f"{{{a_ns}}}hlinkMouseOver"
    rtl_tag = f"{{{a_ns}}}rtl"
    ext_lst_tag = f"{{{a_ns}}}extLst"
    theme_targets = {
        f"{{{a_ns}}}majorFont",
        f"{{{a_ns}}}minorFont",
    }

    def _upsert_typeface(parent: ET.Element, tag: str, typeface: str, successor_tags: tuple[str, ...]) -> None:
        child = parent.find(tag)
        if child is None:
            child = ET.Element(tag)
            insert_at = len(parent)
            for index, existing in enumerate(list(parent)):
                if existing.tag in successor_tags:
                    insert_at = index
                    break
            parent.insert(insert_at, child)
        child.set("typeface", typeface)

    with ZipFile(pptx_path, "r") as source:
        entries: dict[str, bytes] = {}
        for info in source.infolist():
            data = source.read(info.filename)
            if info.filename.startswith("ppt/") and info.filename.endswith(".xml"):
                try:
                    root = ET.fromstring(data)
                except ET.ParseError:
                    entries[info.filename] = data
                    continue

                changed = False
                for element in root.iter():
                    if element.tag in text_prop_tags:
                        _upsert_typeface(
                            element,
                            latin_tag,
                            FONT_LATIN,
                            (ea_tag, cs_tag, sym_tag, hlink_click_tag, hlink_hover_tag, rtl_tag, ext_lst_tag),
                        )
                        _upsert_typeface(
                            element,
                            ea_tag,
                            FONT_EAST_ASIA,
                            (cs_tag, sym_tag, hlink_click_tag, hlink_hover_tag, rtl_tag, ext_lst_tag),
                        )
                        _upsert_typeface(
                            element,
                            cs_tag,
                            FONT_LATIN,
                            (sym_tag, hlink_click_tag, hlink_hover_tag, rtl_tag, ext_lst_tag),
                        )
                        changed = True
                    if element.tag in theme_targets:
                        _upsert_typeface(element, latin_tag, FONT_LATIN, (ea_tag, cs_tag, font_tag))
                        _upsert_typeface(element, ea_tag, FONT_EAST_ASIA, (cs_tag, font_tag))
                        _upsert_typeface(element, cs_tag, FONT_LATIN, (font_tag,))
                        changed = True
                entries[info.filename] = ET.tostring(root, encoding="utf-8", xml_declaration=True) if changed else data
            else:
                entries[info.filename] = data

    temp_path = pptx_path.with_suffix(".fontfix.tmp")
    with ZipFile(temp_path, "w", compression=ZIP_DEFLATED) as target:
        for filename, data in entries.items():
            target.writestr(filename, data)
    temp_path.replace(pptx_path)


def set_paragraph_style(
    paragraph,
    *,
    font_name: str,
    font_size: float,
    color_rgb: tuple[int, int, int],
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """统一设置段落样式。"""
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = RGBColor(*color_rgb)
    paragraph.font.bold = bold
    paragraph.font.italic = italic


def fill_text_frame(
    text_frame,
    paragraphs: list[dict[str, Any]],
    *,
    margin_left: float = 0.14,
    margin_right: float = 0.14,
    margin_top: float = 0.10,
    margin_bottom: float = 0.08,
    vertical_anchor: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
) -> None:
    """向文本框或 shape 填充多段文本。"""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(margin_left)
    text_frame.margin_right = Inches(margin_right)
    text_frame.margin_top = Inches(margin_top)
    text_frame.margin_bottom = Inches(margin_bottom)
    text_frame.vertical_anchor = vertical_anchor

    for index, item in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = item["text"]
        paragraph.space_before = Pt(item.get("space_before", 0))
        paragraph.space_after = Pt(item.get("space_after", 0))
        set_paragraph_style(
            paragraph,
            font_name=item.get("font_name", FONT_BODY),
            font_size=item.get("font_size", 14),
            color_rgb=item.get("color_rgb", PALETTE["subtitle"]),
            bold=item.get("bold", False),
            italic=item.get("italic", False),
            align=item.get("align", PP_ALIGN.LEFT),
        )


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    paragraphs: list[dict[str, Any]],
    *,
    margin_left: float = 0.02,
    margin_right: float = 0.02,
    margin_top: float = 0.00,
    margin_bottom: float = 0.00,
    vertical_anchor: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    """添加多段文本框。"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    fill_text_frame(
        box.text_frame,
        paragraphs,
        margin_left=margin_left,
        margin_right=margin_right,
        margin_top=margin_top,
        margin_bottom=margin_bottom,
        vertical_anchor=vertical_anchor,
    )
    return box


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill_rgb: tuple[int, int, int] | None = None,
    line_rgb: tuple[int, int, int] | None = None,
    line_width: float = 1.0,
    rounded: bool = True,
):
    """添加卡片底板 shape。"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*(fill_rgb or PALETTE["surface"]))
    shape.line.color.rgb = RGBColor(*(line_rgb or PALETTE["line"]))
    shape.line.width = Pt(line_width)
    return shape


def add_accent_bar(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color_rgb: tuple[int, int, int],
    *,
    rounded: bool = False,
):
    """添加装饰性 accent bar。"""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color_rgb)
    shape.line.color.rgb = RGBColor(*color_rgb)
    return shape


def add_card_text(
    shape,
    paragraphs: list[dict[str, Any]],
    *,
    margin_left: float = 0.14,
    margin_right: float = 0.14,
    margin_top: float = 0.10,
    margin_bottom: float = 0.08,
    vertical_anchor: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
) -> None:
    """给卡片 shape 填充文本。"""
    fill_text_frame(
        shape.text_frame,
        paragraphs,
        margin_left=margin_left,
        margin_right=margin_right,
        margin_top=margin_top,
        margin_bottom=margin_bottom,
        vertical_anchor=vertical_anchor,
    )


def add_pill(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_rgb: tuple[int, int, int],
    *,
    font_size: float = 11.0,
    bold: bool = True,
) -> None:
    """添加紧凑型 pill。"""
    shape = add_card(
        slide,
        left,
        top,
        width,
        height,
        fill_rgb=fill_rgb,
        line_rgb=fill_rgb,
        line_width=0.8,
        rounded=True,
    )
    add_card_text(
        shape,
        [
            {
                "text": text,
                "font_name": FONT_BODY,
                "font_size": font_size,
                "bold": bold,
                "align": PP_ALIGN.CENTER,
                "color_rgb": pick_contrast_text_rgb(fill_rgb),
            }
        ],
        margin_left=0.04,
        margin_right=0.04,
        margin_top=0.00,
        margin_bottom=0.00,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )


def add_icon(slide, icon_name: str, left: float, top: float, size: float) -> None:
    """插入主题化 icon。"""
    icon_path = ICON_THEME_DIR / f"{icon_name}.png"
    if not icon_path.exists():
        raise FileNotFoundError(f"缺少 icon: {icon_path}")
    slide.shapes.add_picture(
        str(icon_path),
        Inches(left),
        Inches(top),
        width=Inches(size),
        height=Inches(size),
    )


def add_slide_header(slide, slide_id: str, title: str, subtitle: str, eyebrow: str, accent_rgb: tuple[int, int, int]) -> None:
    """统一添加标题头。"""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*PALETTE["bg"])

    eyebrow_width = min(2.30, max(1.18, 0.11 * len(eyebrow) + 0.34))
    add_pill(slide, eyebrow, CONTENT_LEFT, 0.20, eyebrow_width, 0.24, accent_rgb, font_size=10.0)
    add_textbox(
        slide,
        CONTENT_LEFT,
        0.54,
        CONTENT_WIDTH - 1.1,
        0.56,
        [
            {
                "text": title,
                "font_name": FONT_TITLE,
                "font_size": 24,
                "bold": True,
                "color_rgb": PALETTE["title"],
            }
        ],
    )
    add_textbox(
        slide,
        CONTENT_LEFT,
        0.98,
        CONTENT_WIDTH - 1.0,
        0.26,
        [
            {
                "text": subtitle,
                "font_name": FONT_BODY,
                "font_size": 12.5,
                "color_rgb": PALETTE["subtitle"],
            }
        ],
    )
    add_accent_bar(slide, CONTENT_LEFT, 1.22, 1.85, 0.04, accent_rgb)
    add_textbox(
        slide,
        CONTENT_RIGHT - 0.56,
        8.50,
        0.48,
        0.16,
        [
            {
                "text": slide_id,
                "font_name": FONT_BODY,
                "font_size": 10,
                "bold": True,
                "color_rgb": PALETTE["muted"],
                "align": PP_ALIGN.RIGHT,
            }
        ],
    )


def add_footer_note(slide, text: str) -> None:
    """添加页面底部说明。"""
    add_textbox(
        slide,
        CONTENT_LEFT,
        8.48,
        13.7,
        0.18,
        [
            {
                "text": text,
                "font_name": FONT_BODY,
                "font_size": 10.5,
                "color_rgb": PALETTE["muted"],
            }
        ],
    )


def add_background_orbs(slide, circles: list[dict[str, Any]]) -> None:
    """在背景层添加轻量彩色圆形，增强 page atmosphere。"""
    for circle in circles:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(circle["left"]),
            Inches(circle["top"]),
            Inches(circle["width"]),
            Inches(circle["height"]),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*circle["color_rgb"])
        shape.fill.transparency = circle.get("transparency", 0.62)
        shape.line.color.rgb = RGBColor(*circle["color_rgb"])
        shape.line.transparency = 1.0


def add_event_strip(
    slide,
    events: list[dict[str, Any]],
    left: float,
    top: float,
    width: float,
    accent_rgb: tuple[int, int, int],
    *,
    title_font_size: float = 12.0,
    detail_font_size: float = 11.5,
) -> None:
    """添加三节点事件带。"""
    segment_width = width / max(1, len(events))
    add_accent_bar(slide, left, top + 0.32, width, 0.03, tint(accent_rgb, 0.65))
    for index, event in enumerate(events):
        cursor_left = left + index * segment_width
        add_pill(
            slide,
            str(event["year"]),
            cursor_left + 0.02,
            top,
            0.66,
            0.24,
            accent_rgb if event.get("winner_signal") == "winner" else tint(accent_rgb, 0.55),
            font_size=10.5,
        )
        node = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cursor_left + 0.22),
            Inches(top + 0.28),
            Inches(0.16),
            Inches(0.16),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(*accent_rgb)
        node.line.color.rgb = RGBColor(*accent_rgb)
        add_textbox(
            slide,
            cursor_left + 0.02,
            top + 0.50,
            segment_width - 0.08,
            0.20,
            [
                {
                    "text": event["title"],
                    "font_name": FONT_BODY,
                    "font_size": title_font_size,
                    "bold": True,
                    "color_rgb": PALETTE["title"],
                }
            ],
            margin_left=0.00,
            margin_right=0.00,
        )
        add_textbox(
            slide,
            cursor_left + 0.02,
            top + 0.72,
            segment_width - 0.08,
            0.42,
            [
                {
                    "text": event["detail"],
                    "font_name": FONT_BODY,
                    "font_size": detail_font_size,
                    "color_rgb": PALETTE["subtitle"],
                }
            ],
            margin_left=0.00,
            margin_right=0.00,
        )


def add_case_card(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
    icon_name: str,
    case_label: str,
    title: str,
    body: str,
) -> None:
    """添加封面与结尾页使用的 case anchor card。"""
    card = add_card(slide, left, top, width, height, fill_rgb=(255, 255, 255), line_rgb=tint(accent_rgb, 0.68))
    add_accent_bar(slide, left, top, 0.08, height, accent_rgb, rounded=True)
    add_icon(slide, icon_name, left + 0.18, top + 0.18, 0.28)
    add_textbox(
        slide,
        left + 0.52,
        top + 0.15,
        width - 0.70,
        0.16,
        [{"text": case_label, "font_name": FONT_BODY, "font_size": 10.5, "bold": True, "color_rgb": accent_rgb}],
    )
    add_textbox(
        slide,
        left + 0.18,
        top + 0.48,
        width - 0.32,
        0.22,
        [{"text": title, "font_name": FONT_BODY, "font_size": 14, "bold": True, "color_rgb": PALETTE["title"]}],
    )
    add_textbox(
        slide,
        left + 0.18,
        top + 0.78,
        width - 0.32,
        height - 0.90,
        [{"text": body, "font_name": FONT_BODY, "font_size": 11.5, "color_rgb": PALETTE["subtitle"]}],
    )


def add_simple_table(
    slide,
    *,
    left: float,
    top: float,
    col_widths: list[float],
    row_height: float,
    headers: list[str],
    rows: list[list[dict[str, Any]]],
    accent_rgb: tuple[int, int, int],
    body_fill_rgb: tuple[int, int, int] = (255, 255, 255),
) -> None:
    """绘制用于 comparison page 的原生 PowerPoint 表格。"""
    total_width = sum(col_widths)
    total_rows = len(rows) + 1
    graphic_frame = slide.shapes.add_table(
        total_rows,
        len(headers),
        Inches(left),
        Inches(top),
        Inches(total_width),
        Inches(row_height * total_rows),
    )
    table = graphic_frame.table

    for index, width in enumerate(col_widths):
        table.columns[index].width = Inches(width)
    for index in range(total_rows):
        table.rows[index].height = Inches(row_height)

    def _fill_cell(cell, spec: dict[str, Any], *, default_fill: tuple[int, int, int]) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*spec.get("fill_rgb", default_fill))
        cell.margin_left = Inches(spec.get("margin_left", 0.08))
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.03)
        fill_text_frame(
            cell.text_frame,
            [
                {
                    "text": spec["text"],
                    "font_name": spec.get("font_name", FONT_BODY),
                    "font_size": spec.get("font_size", 11.5),
                    "bold": spec.get("bold", False),
                    "align": spec.get("align", PP_ALIGN.LEFT),
                    "color_rgb": spec.get("color_rgb", PALETTE["subtitle"]),
                }
            ],
            margin_left=0.00,
            margin_right=0.00,
            margin_top=0.00,
            margin_bottom=0.00,
            vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
        )

    for col_index, header in enumerate(headers):
        _fill_cell(
            table.cell(0, col_index),
            {
                "text": header,
                "font_name": FONT_BODY,
                "font_size": 12,
                "bold": True,
                "align": PP_ALIGN.CENTER,
                "color_rgb": PALETTE["title"],
                "margin_left": 0.04,
            },
            default_fill=tint(accent_rgb, 0.18),
        )

    for row_index, row in enumerate(rows, start=1):
        for col_index, cell_spec in enumerate(row):
            _fill_cell(
                table.cell(row_index, col_index),
                cell_spec,
                default_fill=body_fill_rgb,
            )


# ---------------------------------------------------------------------------
# 逐页构建
# ---------------------------------------------------------------------------


def build_slide_01(slide, spec: dict[str, Any]) -> None:
    """构建封面页。"""
    add_background_orbs(
        slide,
        [
            {"left": 9.4, "top": -0.6, "width": 4.2, "height": 4.2, "color_rgb": tint(PALETTE["blue"], 0.35), "transparency": 0.55},
            {"left": 12.0, "top": 2.0, "width": 3.5, "height": 3.5, "color_rgb": tint(PALETTE["warm"], 0.35), "transparency": 0.64},
            {"left": 8.5, "top": 4.8, "width": 4.8, "height": 4.8, "color_rgb": tint(PALETTE["sage"], 0.25), "transparency": 0.72},
        ],
    )
    add_pill(slide, "Deck Thesis", CONTENT_LEFT, 0.44, 1.18, 0.24, PALETTE["blue"], font_size=10.5)
    add_textbox(
        slide,
        CONTENT_LEFT,
        1.02,
        8.0,
        1.62,
        [
            {"text": "标准不会奖励最优技术，", "font_name": FONT_TITLE, "font_size": 28, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "标准奖励最容易被生态一起采用的方案。", "font_name": FONT_TITLE, "font_size": 28, "bold": True, "color_rgb": PALETTE["title"], "space_before": 5},
        ],
    )
    add_textbox(
        slide,
        CONTENT_LEFT,
        2.78,
        7.25,
        0.66,
        [
            {
                "text": "从 Betamax、OSI 到 Blu-ray，市场真正结算的是协同成本、切换风险、补充品联盟、部署速度与预期管理。",
                "font_name": FONT_BODY,
                "font_size": 15,
                "color_rgb": PALETTE["subtitle"],
            }
        ],
    )

    add_pill(slide, "误判纠正", CONTENT_LEFT, 3.80, 0.92, 0.22, tint(PALETTE["warm"], 0.94), font_size=10)
    add_pill(slide, "三场案例", CONTENT_LEFT + 1.06, 3.80, 0.92, 0.22, tint(PALETTE["blue"], 0.90), font_size=10)
    add_pill(slide, "一个框架", CONTENT_LEFT + 2.12, 3.80, 0.92, 0.22, tint(PALETTE["sage"], 0.90), font_size=10)

    add_textbox(
        slide,
        CONTENT_LEFT,
        4.28,
        6.9,
        1.34,
        [
            {"text": "市场不会先问“哪项指标更漂亮”，市场先问“整个系统愿不愿意一起下注”。", "font_name": FONT_BODY, "font_size": 18, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "这也是为什么标准战的赢家，经常不是规格表上最耀眼的一方。", "font_name": FONT_BODY, "font_size": 14, "color_rgb": PALETTE["subtitle"], "space_before": 8},
        ],
    )

    add_case_card(
        slide,
        left=9.15,
        top=1.04,
        width=5.35,
        height=1.92,
        accent_rgb=PALETTE["warm"],
        icon_name="presentation",
        case_label="CLAIM 1 | Home Video",
        title="任务匹配和联盟速度会压过峰值规格",
        body="Betamax vs VHS 证明：客厅先为“能不能录完整场比赛”和“厂商愿不愿意站队”投票。",
    )
    add_case_card(
        slide,
        left=9.45,
        top=3.15,
        width=5.05,
        height=1.92,
        accent_rgb=PALETTE["blue"],
        icon_name="network",
        case_label="CLAIM 2 | Networking",
        title="先跑起来并得到反馈的体系会定义现实",
        body="OSI vs TCP/IP 证明：deployment loop 会比理论完整性更早决定 adoption。",
    )
    add_case_card(
        slide,
        left=9.15,
        top=5.28,
        width=5.35,
        height=1.92,
        accent_rgb=PALETTE["sage"],
        icon_name="device-desktop-analytics",
        case_label="CLAIM 3 | HD Optical",
        title="补充品和 bundle 会在不透明市场替用户下注",
        body="Blu-ray vs HD DVD 证明：内容联盟、默认装机和 dominant standard 预期会替用户做判断。",
    )
    add_footer_note(slide, "这套 deck 把三场历史标准战压成一个今天仍然可用的判断框架。")
    add_textbox(
        slide,
        CONTENT_RIGHT - 0.56,
        8.50,
        0.48,
        0.16,
        [{"text": spec["slide_id"], "font_name": FONT_BODY, "font_size": 10, "bold": True, "align": PP_ALIGN.RIGHT, "color_rgb": PALETTE["muted"]}],
    )


def build_slide_02(slide, spec: dict[str, Any]) -> None:
    """构建误判纠正页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Reset The Lens", PALETTE["warm"])

    left_card = add_card(slide, CONTENT_LEFT, BODY_TOP, 6.5, 4.55, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["warm"], 0.65))
    add_accent_bar(slide, CONTENT_LEFT, BODY_TOP, 6.5, 0.10, tint(PALETTE["warm"], 0.92))
    add_icon(slide, "zoom-code", CONTENT_LEFT + 0.18, BODY_TOP + 0.20, 0.26)
    add_card_text(
        left_card,
        [
            {"text": "我们习惯把市场想成技术评委", "font_name": FONT_BODY, "font_size": 16, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "更高清、更优雅、更完整的方案，理应拿到更高分。", "font_name": FONT_BODY, "font_size": 13.5, "color_rgb": PALETTE["subtitle"], "space_before": 8},
            {"text": "于是技术人会先看规格表、协议层级、编码参数、设计纯度。", "font_name": FONT_BODY, "font_size": 13.5, "color_rgb": PALETTE["subtitle"], "space_before": 6},
        ],
        margin_left=0.22,
        margin_right=0.18,
        margin_top=0.18,
    )
    for idx, bullet in enumerate(
        [
            "更高峰值性能 = 更大获胜概率",
            "更完整架构 = 更合理的长期标准",
            "更漂亮的设计 = 更强的市场吸引力",
        ]
    ):
        add_pill(slide, bullet, CONTENT_LEFT + 0.20, BODY_TOP + 1.56 + idx * 0.48, 5.86, 0.28, tint(PALETTE["warm"], 0.88), font_size=11)

    right_left = CONTENT_LEFT + 6.78
    right_card = add_card(slide, right_left, BODY_TOP, 7.66, 4.55, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["blue"], 0.64))
    add_accent_bar(slide, right_left, BODY_TOP, 7.66, 0.10, tint(PALETTE["blue"], 0.92))
    add_icon(slide, "users", right_left + 0.18, BODY_TOP + 0.20, 0.26)
    add_card_text(
        right_card,
        [
            {"text": "市场真正定价的是采用风险", "font_name": FONT_BODY, "font_size": 16, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "标准战里，买方、渠道、补充品和制造方需要同时判断“现在下注会不会错”。", "font_name": FONT_BODY, "font_size": 13.5, "color_rgb": PALETTE["subtitle"], "space_before": 8},
        ],
        margin_left=0.22,
        margin_right=0.18,
        margin_top=0.18,
    )
    claim_cards = [
        ("presentation", "Claim 1 | 任务匹配和联盟速度会压过峰值规格", "Home Video 会证明：用户任务与阵营扩张会先于画质神话决定标准走向。"),
        ("rocket", "Claim 2 | 先跑起来并得到反馈的体系会定义现实", "Networking 会证明：更早落地并获得反馈的协议，会重写“应该如何”的辩论。"),
        ("circles-relation", "Claim 3 | 补充品和 bundle 会在不透明市场替用户下注", "HD Optical 会证明：当用户看不懂技术差距时，内容、渠道与默认装机会替他们表态。"),
    ]
    for idx, (icon_name, title, body) in enumerate(claim_cards):
        card = add_card(
            slide,
            right_left + 0.20,
            BODY_TOP + 1.24 + idx * 1.04,
            7.12,
            0.88,
            fill_rgb=tint(PALETTE["blue"], 0.08),
            line_rgb=tint(PALETTE["blue"], 0.55),
        )
        add_icon(slide, icon_name, right_left + 0.34, BODY_TOP + 1.48 + idx * 1.04, 0.22)
        add_textbox(
            slide,
            right_left + 0.66,
            BODY_TOP + 1.38 + idx * 1.04,
            6.32,
            0.18,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.2, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            right_left + 0.66,
            BODY_TOP + 1.62 + idx * 1.04,
            6.18,
            0.18,
            [{"text": body, "font_name": FONT_BODY, "font_size": 10.9, "color_rgb": PALETTE["subtitle"]}],
        )

    bottom = add_card(slide, CONTENT_LEFT, 6.18, CONTENT_WIDTH, 1.38, fill_rgb=tint(PALETTE["sage"], 0.10), line_rgb=tint(PALETTE["sage"], 0.58))
    add_icon(slide, "arrows-transfer-up-down", CONTENT_LEFT + 0.20, 6.38, 0.28)
    add_card_text(
        bottom,
        [
            {"text": "接下来 6 页的工作，就是把这三条 claim 一条一条坐实。", "font_name": FONT_BODY, "font_size": 15, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "而第 3 页会给出这些 claim 背后的共同解释器：采用风险与网络效应如何放大早期差异。", "font_name": FONT_BODY, "font_size": 13, "color_rgb": PALETTE["subtitle"], "space_before": 7},
        ],
        margin_left=0.56,
        margin_right=0.14,
        margin_top=0.15,
    )
    add_footer_note(slide, "从这一页开始，整个 deck 的判断轴已经从“性能对比”切换到“采用风险对比”。")


def build_slide_03(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建五因素 connector 框架页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Five-Factor Frame", PALETTE["blue"])

    panel = add_card(slide, CONTENT_LEFT, BODY_TOP, 9.35, 6.26, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["blue"], 0.62))
    add_card_text(
        panel,
        [
            {"text": "市场投票路径", "font_name": FONT_BODY, "font_size": 12, "bold": True, "color_rgb": PALETTE["blue"]},
        ],
        margin_left=0.18,
        margin_top=0.08,
        margin_bottom=0.02,
    )
    add_textbox(
        slide,
        CONTENT_LEFT + 0.20,
        BODY_TOP + 0.42,
        8.9,
        0.30,
        [{"text": "局部技术优势只有在被兼容路径、联盟速度和部署反馈接住时，才会真正转化成标准地位。", "font_name": FONT_BODY, "font_size": 12.5, "color_rgb": PALETTE["subtitle"]}],
    )

    node_style_main = NodeStyle(fill_rgb=tint(PALETTE["blue"], 0.12), line_rgb=tint(PALETTE["blue"], 0.76), text_rgb=PALETTE["title"], font_size=13)
    node_style_alt = NodeStyle(fill_rgb=tint(PALETTE["warm"], 0.12), line_rgb=tint(PALETTE["warm"], 0.76), text_rgb=PALETTE["title"], font_size=13)
    node_style_end = NodeStyle(fill_rgb=tint(PALETTE["sage"], 0.12), line_rgb=tint(PALETTE["sage"], 0.76), text_rgb=PALETTE["title"], font_size=13)

    node_tech = add_node(slide, "Node Tech", "局部技术优势", CONTENT_LEFT + 0.40, BODY_TOP + 1.58, 1.46, 0.78, node_style_main)
    node_c1 = add_node(slide, "Node Compatibility", "兼容与迁移", CONTENT_LEFT + 2.18, BODY_TOP + 1.24, 1.68, 0.92, node_style_alt)
    node_c2 = add_node(slide, "Node Base", "装机基盘", CONTENT_LEFT + 4.10, BODY_TOP + 1.96, 1.60, 0.92, node_style_main)
    node_c3 = add_node(slide, "Node Complements", "补充品联盟", CONTENT_LEFT + 6.00, BODY_TOP + 1.24, 1.82, 0.92, node_style_alt)
    node_c4 = add_node(slide, "Node Speed", "部署速度", CONTENT_LEFT + 7.95, BODY_TOP + 1.96, 1.56, 0.92, node_style_main)
    node_c5 = add_node(slide, "Node Expectation", "预期管理", CONTENT_LEFT + 5.95, BODY_TOP + 3.26, 1.82, 0.92, node_style_alt)
    node_end = add_node(slide, "Node Result", "标准结果", CONTENT_LEFT + 3.86, BODY_TOP + 4.26, 1.82, 1.00, node_style_end)

    for from_shape, to_shape, from_site, to_site in [
        (node_tech, node_c1, "right", "left"),
        (node_c1, node_c2, "right", "left"),
        (node_c2, node_c3, "right", "left"),
        (node_c3, node_c4, "right", "left"),
        (node_c4, node_c5, "bottom", "right"),
        (node_c5, node_end, "left", "right"),
        (node_c2, node_end, "bottom", "top"),
    ]:
        add_glued_connector(slide, from_shape, to_shape, from_site, to_site, tint(PALETTE["blue"], 0.70), connector_type=MSO_CONNECTOR.ELBOW, line_width=1.6)

    add_pill(slide, "规格表只解释左边", CONTENT_LEFT + 0.44, BODY_TOP + 5.44, 1.52, 0.22, tint(PALETTE["warm"], 0.92), font_size=10)
    add_pill(slide, "采用结果由整条路径共同决定", CONTENT_LEFT + 2.18, BODY_TOP + 5.44, 2.32, 0.22, tint(PALETTE["blue"], 0.90), font_size=10)

    sidebar_left = CONTENT_LEFT + 9.62
    add_textbox(
        slide,
        sidebar_left,
        BODY_TOP,
        5.25,
        0.28,
        [{"text": "五个变量各自压住什么风险", "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["title"]}],
    )
    for index, row in ctx.factor_frame.iterrows():
        card = add_card(
            slide,
            sidebar_left,
            BODY_TOP + 0.42 + index * 1.10,
            5.18,
            0.94,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_pill(
            slide,
            row["short_label"],
            sidebar_left + 0.12,
            BODY_TOP + 0.54 + index * 1.10,
            0.88,
            0.18,
            PALETTE["blue"] if index % 2 == 0 else PALETTE["warm"],
            font_size=9.5,
        )
        add_textbox(
            slide,
            sidebar_left + 1.12,
            BODY_TOP + 0.50 + index * 1.10,
            3.90,
            0.18,
            [{"text": row["market_question"], "font_name": FONT_BODY, "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            sidebar_left + 1.12,
            BODY_TOP + 0.74 + index * 1.10,
            3.90,
            0.24,
            [{"text": row["why_it_matters"], "font_name": FONT_BODY, "font_size": 10.8, "color_rgb": PALETTE["subtitle"]}],
        )

    add_footer_note(slide, "结构图使用真 connector，后续校验会确认连线确实绑定在业务节点上。")


def build_slide_04(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建 Betamax vs VHS 案例页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 01 | Home Video", PALETTE["warm"])

    verdict = add_card(slide, CONTENT_LEFT, BODY_TOP, 14.44, 1.10, fill_rgb=tint(PALETTE["warm"], 0.10), line_rgb=tint(PALETTE["warm"], 0.62))
    add_icon(slide, "presentation", CONTENT_LEFT + 0.18, BODY_TOP + 0.22, 0.30)
    add_card_text(
        verdict,
        [
            {"text": "客厅不会为画质投票，它先为“整场球赛能不能录完”和“厂商愿不愿意一起押注”投票。", "font_name": FONT_BODY, "font_size": 16, "bold": True, "color_rgb": PALETTE["title"]},
        ],
        margin_left=0.60,
        margin_right=0.16,
        margin_top=0.18,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_pill(slide, "正在验证 Claim 1", CONTENT_LEFT + 9.08, BODY_TOP + 0.22, 1.56, 0.20, tint(PALETTE["sage"], 0.92), font_size=9.5)
    add_pill(slide, "1 小时 vs 2 小时", CONTENT_LEFT + 10.90, BODY_TOP + 0.22, 1.18, 0.20, PALETTE["warm"], font_size=9.5)
    add_pill(slide, "谨慎授权 vs 更广阵营", CONTENT_LEFT + 12.18, BODY_TOP + 0.22, 1.98, 0.20, PALETTE["blue"], font_size=9.5)

    events = ctx.milestones_frame.query("case == 'Home Video'").to_dict("records")
    add_event_strip(slide, events, CONTENT_LEFT, BODY_TOP + 1.42, 8.72, PALETTE["warm"])

    lens_specs = [
        ("users", "家庭任务匹配", "VHS 一上来就更接近主流家庭录制整场节目、赛事和电影的任务需求。"),
        ("stack-2", "制造成本与介质经济性", "更少组件和更低制造复杂度，让 VHS 更容易形成规模化供给。"),
        ("circles-relation", "授权与阵营扩张", "更广泛的厂商支持更快堆出 installed base，并放大后续补充品供给。"),
    ]
    base_left = CONTENT_LEFT + 9.08
    for index, (icon_name, title, body) in enumerate(lens_specs):
        card = add_card(
            slide,
            base_left,
            BODY_TOP + 1.42 + index * 1.54,
            5.36,
            1.34,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_icon(slide, icon_name, base_left + 0.16, BODY_TOP + 1.62 + index * 1.54, 0.24)
        add_textbox(
            slide,
            base_left + 0.48,
            BODY_TOP + 1.58 + index * 1.54,
            4.60,
            0.18,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            base_left + 0.16,
            BODY_TOP + 1.88 + index * 1.54,
            4.92,
            0.48,
            [{"text": body, "font_name": FONT_BODY, "font_size": 11.4, "color_rgb": PALETTE["subtitle"]}],
        )

    add_textbox(
        slide,
        CONTENT_LEFT,
        BODY_TOP + 4.98,
        8.66,
        1.30,
        [
            {"text": "为什么“更好技术”的神话会误导判断", "font_name": FONT_BODY, "font_size": 13, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "Betamax 后来经常和专业格式 Betacam 的品质记忆混在一起，但家庭录像市场真正放大的，是更低风险的任务完成路径和更大的阵营扩张速度。", "font_name": FONT_BODY, "font_size": 12.5, "color_rgb": PALETTE["subtitle"], "space_before": 6},
        ],
    )
    add_footer_note(slide, "Sources: Sony corporate history; Smithsonian review; AskHistorians discussion summarized in the design brief.")


def build_slide_05(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建 VHS lesson 页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 01 | Lesson", PALETTE["warm"])

    loop_panel = add_card(slide, CONTENT_LEFT, BODY_TOP, 7.30, 5.98, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["warm"], 0.62))
    add_textbox(
        slide,
        CONTENT_LEFT + 0.18,
        BODY_TOP + 0.16,
        6.80,
        0.18,
        [{"text": "VHS 如何闭合“更低风险下注”的循环", "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["title"]}],
    )
    loop_boxes = [
        (CONTENT_LEFT + 0.40, BODY_TOP + 0.76, 2.70, 1.00, "先满足家庭任务", "更长录制时长先解决了主流使用场景。"),
        (CONTENT_LEFT + 3.98, BODY_TOP + 0.76, 2.64, 1.00, "再形成厂商联盟", "更低制造复杂度和更广授权让阵营扩大更快。"),
        (CONTENT_LEFT + 3.98, BODY_TOP + 2.42, 2.64, 1.00, "再堆大 installed base", "设备数、内容数与渠道信号开始相互强化。"),
        (CONTENT_LEFT + 0.40, BODY_TOP + 2.42, 2.70, 1.00, "最后让市场预期锁定", "“大家都在用 VHS” 让后续下注更轻松。"),
    ]
    shapes = []
    for idx, (left, top, width, height, title, body) in enumerate(loop_boxes):
        fill = tint(PALETTE["warm"], 0.10) if idx % 2 == 0 else tint(PALETTE["blue"], 0.10)
        line = tint(PALETTE["warm"], 0.72) if idx % 2 == 0 else tint(PALETTE["blue"], 0.72)
        card = add_card(slide, left, top, width, height, fill_rgb=fill, line_rgb=line)
        add_card_text(
            card,
            [
                {"text": title, "font_name": FONT_BODY, "font_size": 13, "bold": True, "color_rgb": PALETTE["title"]},
                {"text": body, "font_name": FONT_BODY, "font_size": 11.5, "color_rgb": PALETTE["subtitle"], "space_before": 5},
            ],
        )
        shapes.append(card)

    arrow_specs = [
        (CONTENT_LEFT + 3.18, BODY_TOP + 1.15, 0.52, 0.12, MSO_AUTO_SHAPE_TYPE.CHEVRON),
        (CONTENT_LEFT + 5.08, BODY_TOP + 1.92, 0.16, 0.38, MSO_AUTO_SHAPE_TYPE.DOWN_ARROW),
        (CONTENT_LEFT + 3.18, BODY_TOP + 2.82, 0.52, 0.12, MSO_AUTO_SHAPE_TYPE.CHEVRON),
        (CONTENT_LEFT + 1.52, BODY_TOP + 1.92, 0.16, 0.38, MSO_AUTO_SHAPE_TYPE.UP_ARROW),
    ]
    for left, top, width, height, shape_type in arrow_specs:
        arrow = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(*tint(PALETTE["warm"], 0.78))
        arrow.line.color.rgb = RGBColor(*tint(PALETTE["warm"], 0.78))

    takeaway_left = CONTENT_LEFT + 7.56
    categories = ["峰值规格", "任务匹配", "制造经济性", "联盟速度"]
    video_chart = add_native_chart_card(
        slide,
        title="综合评分 | 分析性合成",
        left=takeaway_left,
        top=BODY_TOP,
        width=6.88,
        height=3.06,
        accent_rgb=PALETTE["blue"],
        categories=categories,
        series_list=[
            ("Betamax", ctx.home_video_scorecard.loc[0, categories].tolist()),
            ("VHS", ctx.home_video_scorecard.loc[1, categories].tolist()),
        ],
        show_legend=True,
        series_colors=[tint(PALETTE["warm"], 0.94), PALETTE["blue"]],
    )
    video_chart.value_axis.minimum_scale = 0
    video_chart.value_axis.maximum_scale = 5
    video_chart.value_axis.major_unit = 1
    add_textbox(
        slide,
        takeaway_left + 0.14,
        BODY_TOP + 2.86,
        6.26,
        0.16,
        [{"text": "评分用于把案例里的机制差异显式化，不代表原始市场统计。", "font_name": FONT_BODY, "font_size": 9.8, "color_rgb": PALETTE["muted"]}],
    )

    takeaways = [
        ("target-arrow", "任务匹配优先于峰值性能", "标准战的第一步不是赢技术论坛，而是赢主流用户要完成的任务。"),
        ("circles-relation", "更开放的授权会更快堆出 installed base", "标准化不是孤勇者竞赛，而是联盟速度竞赛。"),
    ]
    for index, (icon_name, title, body) in enumerate(takeaways):
        card = add_card(
            slide,
            takeaway_left,
            BODY_TOP + 3.34 + index * 1.28,
            6.88,
            1.04,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_pill(
            slide,
            f"Takeaway {index + 1}",
            takeaway_left + 0.14,
            BODY_TOP + 3.48 + index * 1.28,
            0.88,
            0.18,
            PALETTE["warm"] if index == 0 else PALETTE["blue"],
            font_size=9.8,
        )
        add_icon(slide, icon_name, takeaway_left + 0.16, BODY_TOP + 3.74 + index * 1.28, 0.22)
        add_textbox(
            slide,
            takeaway_left + 0.46,
            BODY_TOP + 3.70 + index * 1.28,
            6.00,
            0.16,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.4, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            takeaway_left + 0.16,
            BODY_TOP + 3.96 + index * 1.28,
            6.34,
            0.22,
            [{"text": body, "font_name": FONT_BODY, "font_size": 10.9, "color_rgb": PALETTE["subtitle"]}],
        )

    add_footer_note(slide, "Betamax/VHS 的 lesson 是：先闭合采用风险，规格优势才有机会兑现。图表为分析性合成评分。")


def build_slide_06(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建 OSI vs TCP/IP 案例页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 02 | Networking", PALETTE["blue"])

    hero = add_card(slide, CONTENT_LEFT, BODY_TOP, 14.44, 0.96, fill_rgb=tint(PALETTE["blue"], 0.10), line_rgb=tint(PALETTE["blue"], 0.62))
    add_icon(slide, "network", CONTENT_LEFT + 0.18, BODY_TOP + 0.18, 0.28)
    add_card_text(
        hero,
        [{"text": "机房会奖励先跑起来、可免费实现、能更快得到反馈的协议栈。", "font_name": FONT_BODY, "font_size": 16, "bold": True, "color_rgb": PALETTE["title"]}],
        margin_left=0.58,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_pill(slide, "正在验证 Claim 2", CONTENT_LEFT + 12.00, BODY_TOP + 0.20, 1.56, 0.20, tint(PALETTE["warm"], 0.92), font_size=9.5)

    events = []
    for event in ctx.milestones_frame.query("case == 'Networking'").to_dict("records"):
        normalized_event = dict(event)
        if normalized_event["title"].startswith("ARPANET"):
            normalized_event["title"] = "TCP/IP 落地"
        elif normalized_event["title"].startswith("OSI"):
            normalized_event["title"] = "OSI 推进仍迟缓"
        elif normalized_event["title"].startswith("Is OSI Too Late"):
            normalized_event["title"] = "《Too Late?》"
            normalized_event["detail"] = "行业开始接受“太晚了”的判断，时间窗口正式关闭。"
        events.append(normalized_event)
    left_panel = add_card(slide, CONTENT_LEFT, BODY_TOP + 1.26, 5.00, 3.04, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["blue"], 0.60))
    add_textbox(
        slide,
        CONTENT_LEFT + 0.18,
        BODY_TOP + 1.44,
        4.50,
        0.18,
        [{"text": "时间窗口如何锁定现实路径", "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["title"]}],
    )
    add_event_strip(
        slide,
        events,
        CONTENT_LEFT + 0.18,
        BODY_TOP + 1.80,
        4.56,
        PALETTE["blue"],
        title_font_size=10.0,
        detail_font_size=10.7,
    )
    add_picture_card(
        slide,
        title="Heatmap",
        image_path=ctx.networking_heatmap_path,
        left=CONTENT_LEFT,
        top=BODY_TOP + 4.48,
        width=5.00,
        height=2.18,
        accent_rgb=PALETTE["warm"],
        caption="热力图用于显式化 OSI / TCP/IP 在四个机制维度上的强弱，不代表原始市场统计。",
    )

    compare_left = CONTENT_LEFT + 5.30
    headers = ["维度", "OSI", "TCP/IP"]
    rows = [
        [
            {"text": "开发与协同速度", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "更慢，治理和标准开发摩擦更大", "font_size": 11.2, "color_rgb": PALETTE["subtitle"]},
            {"text": "更快，现实网络先跑起来", "font_size": 11.2, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "实现获取成本", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "标准文档和实现门槛更重", "font_size": 11.2, "color_rgb": PALETTE["subtitle"]},
            {"text": "政府资助 + Berkeley 分发，免费实现更普及", "font_size": 11.2, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "运维与诊断", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "更理想化，但现实诊断链路更重", "font_size": 11.2, "color_rgb": PALETTE["subtitle"]},
            {"text": "更易 eyeball、调试和部署", "font_size": 11.2, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "政策与现实", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "政策层被鼓励，但现实 adoption 迟缓", "font_size": 11.2, "color_rgb": PALETTE["subtitle"]},
            {"text": "现实部署先行，政策最终被现实追认", "font_size": 11.2, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
    ]
    add_simple_table(slide, left=compare_left, top=BODY_TOP + 1.26, col_widths=[1.90, 3.05, 3.95], row_height=0.84, headers=headers, rows=rows, accent_rgb=PALETTE["blue"])

    callout = add_card(slide, compare_left, BODY_TOP + 5.02, 8.90, 1.24, fill_rgb=tint(PALETTE["warm"], 0.10), line_rgb=tint(PALETTE["warm"], 0.68))
    add_icon(slide, "messages", compare_left + 0.18, BODY_TOP + 5.20, 0.22)
    add_card_text(
        callout,
        [
            {"text": "关键情绪拐点", "font_name": FONT_BODY, "font_size": 12, "bold": True, "color_rgb": PALETTE["warm"]},
            {"text": "《Is OSI Too Late?》之所以击中行业，是因为现实已经让时间窗口关闭。", "font_name": FONT_BODY, "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"], "space_before": 4},
            {"text": "标准开发越慢，先跑起来的一方越容易反过来定义事实标准。", "font_name": FONT_BODY, "font_size": 11.2, "color_rgb": PALETTE["subtitle"], "space_before": 4},
        ],
        margin_left=0.58,
        margin_right=0.12,
        margin_top=0.12,
        margin_bottom=0.06,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_footer_note(slide, "Sources: IEEE Spectrum on OSI; policy and Berkeley/TCP-IP implementation history summarized in the design brief. 本页含 Python figure 证据图。")


def build_slide_07(slide, spec: dict[str, Any]) -> None:
    """构建 OSI vs TCP/IP lesson 页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 02 | Lesson", PALETTE["blue"])

    headers = ["维度", "OSI 逻辑", "TCP/IP 逻辑"]
    rows = [
        [
            {"text": "价值主张", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "更完整、分层更理想", "font_size": 11.5, "color_rgb": PALETTE["subtitle"]},
            {"text": "更快进入现实网络并获得反馈", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "采用门槛", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "协调与理解成本更高", "font_size": 11.5, "color_rgb": PALETTE["subtitle"]},
            {"text": "实现、分发与调试门槛更低", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "现实回路", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "更晚进入真实运维场景", "font_size": 11.5, "color_rgb": PALETTE["subtitle"]},
            {"text": "更早进入真实运维场景", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
        [
            {"text": "最终结果", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "体系优雅未转化成标准地位", "font_size": 11.5, "color_rgb": PALETTE["subtitle"]},
            {"text": "现实可用性反过来定义了事实标准", "font_size": 11.5, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["blue"], 0.08)},
        ],
    ]
    add_simple_table(slide, left=CONTENT_LEFT, top=BODY_TOP, col_widths=[1.60, 4.16, 4.54], row_height=0.86, headers=headers, rows=rows, accent_rgb=PALETTE["blue"])

    right_left = CONTENT_LEFT + 10.66
    network_chart = add_native_chart_card(
        slide,
        title="分析性合成评分",
        left=right_left,
        top=BODY_TOP,
        width=4.56,
        height=2.18,
        accent_rgb=PALETTE["warm"],
        categories=["完整性", "速度", "成本", "运维"],
        series_list=[
            ("OSI", [5, 2, 2, 2]),
            ("TCP/IP", [3, 5, 5, 5]),
        ],
        show_legend=False,
        series_colors=[tint(PALETTE["warm"], 0.94), PALETTE["blue"]],
    )
    network_chart.value_axis.minimum_scale = 0
    network_chart.value_axis.maximum_scale = 5
    network_chart.value_axis.major_unit = 1
    add_textbox(
        slide,
        right_left + 0.16,
        BODY_TOP + 2.06,
        4.10,
        0.16,
        [{"text": "蓝色 = TCP/IP，橙色 = OSI。评分用于显式化机制差异，不代表原始市场统计。", "font_name": FONT_BODY, "font_size": 9.6, "color_rgb": PALETTE["muted"]}],
    )

    lessons = [
        ("rocket", "free beats elegant", "免费实现与分发路径会极大改变 adoption 曲线。"),
        ("sparkles", "feedback beats purity", "更早拿到真实反馈，胜过更晚拿到完整设计。"),
        ("shield-check", "ops matters", "能否被运维和工程团队低成本地用起来，是标准地位的关键前提。"),
    ]
    for idx, (icon_name, title, body) in enumerate(lessons):
        card = add_card(
            slide,
            right_left,
            BODY_TOP + 2.34 + idx * 1.22,
            4.56,
            0.98,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_pill(slide, f"Lesson {idx + 1}", right_left + 0.14, BODY_TOP + 2.48 + idx * 1.22, 0.74, 0.18, PALETTE["blue"], font_size=9.5)
        add_icon(slide, icon_name, right_left + 0.16, BODY_TOP + 2.70 + idx * 1.22, 0.22)
        add_textbox(
            slide,
            right_left + 0.46,
            BODY_TOP + 2.68 + idx * 1.22,
            3.78,
            0.18,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.3, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            right_left + 0.14,
            BODY_TOP + 2.92 + idx * 1.22,
            4.18,
            0.28,
            [{"text": body, "font_name": FONT_BODY, "font_size": 11.2, "color_rgb": PALETTE["subtitle"]}],
        )

    add_footer_note(slide, "第二个案例说明：架构完整性不会自动转化成标准地位，现实 deployment loop 才会。")


def build_slide_08(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建 Blu-ray vs HD DVD 案例页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 03 | HD Optical", PALETTE["sage"])

    board = add_card(slide, CONTENT_LEFT, BODY_TOP, 9.18, 6.06, fill_rgb=(255, 255, 255), line_rgb=tint(PALETTE["sage"], 0.60))
    add_textbox(
        slide,
        CONTENT_LEFT + 0.18,
        BODY_TOP + 0.16,
        8.68,
        0.18,
        [{"text": "谁在替市场做决定", "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["title"]}],
    )
    lane_specs = [
        ("book", "内容方 / Studio", "当技术差距不显著时，独家内容站队会极大影响格式命运。"),
        ("device-desktop-analytics", "默认装机 / Console", "PS3 把 Blu-ray 带进用户家中，installed base 从设备端被加速。"),
        ("briefcase", "渠道 / Retail", "零售和租赁体系会根据 dominant standard 预期重新配置资源。"),
        ("users", "消费者 / Households", "消费者不想买错标准，于是会等待生态 signal 帮自己下注。"),
    ]
    for idx, (icon_name, title, body) in enumerate(lane_specs):
        top = BODY_TOP + 0.56 + idx * 1.26
        card = add_card(
            slide,
            CONTENT_LEFT + 0.18,
            top,
            8.74,
            0.94,
            fill_rgb=tint(PALETTE["sage"], 0.08) if idx % 2 == 0 else (255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_icon(slide, icon_name, CONTENT_LEFT + 0.34, top + 0.18, 0.24)
        add_textbox(
            slide,
            CONTENT_LEFT + 0.66,
            top + 0.14,
            2.34,
            0.18,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.3, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            CONTENT_LEFT + 3.12,
            top + 0.14,
            5.40,
            0.28,
            [{"text": body, "font_name": FONT_BODY, "font_size": 11.2, "color_rgb": PALETTE["subtitle"]}],
        )

    add_pill(slide, "dominant standard 预期", CONTENT_LEFT + 3.16, BODY_TOP + 5.32, 1.52, 0.22, PALETTE["sage"], font_size=10)
    add_pill(slide, "内容联盟", CONTENT_LEFT + 0.80, BODY_TOP + 5.32, 0.90, 0.22, PALETTE["warm"], font_size=10)
    add_pill(slide, "PS3 装机", CONTENT_LEFT + 1.84, BODY_TOP + 5.32, 0.90, 0.22, PALETTE["blue"], font_size=10)
    add_pill(slide, "Toshiba 退出", CONTENT_LEFT + 4.92, BODY_TOP + 5.32, 1.08, 0.22, PALETTE["warm"], font_size=10)

    right_left = CONTENT_LEFT + 9.44
    highlight = add_card(slide, right_left, BODY_TOP, 5.78, 2.18, fill_rgb=tint(PALETTE["warm"], 0.10), line_rgb=tint(PALETTE["warm"], 0.68))
    add_card_text(
        highlight,
        [
            {"text": "PS3 是这场战争里最关键的 bundle engine", "font_name": FONT_BODY, "font_size": 14.5, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "它把本来需要单独说服消费者购买的新格式，直接塞进了一个已经有强需求的设备路径里。", "font_name": FONT_BODY, "font_size": 11.8, "color_rgb": PALETTE["subtitle"], "space_before": 8},
        ],
        margin_left=0.16,
        margin_right=0.14,
        margin_top=0.16,
    )
    add_pill(slide, "正在验证 Claim 3", right_left + 3.92, BODY_TOP + 0.16, 1.56, 0.20, tint(PALETTE["blue"], 0.92), font_size=9.5)

    signal_cards = [
        ("内容联盟先表态", "当消费者无法判断技术优越性时，studio 和租赁渠道的站队本身就是决策信号。"),
        ("预期一旦形成就会自强化", "越多人相信 Blu-ray 会赢，越多人不愿意押注 HD DVD。"),
        ("退出声明是战争收束点", "Toshiba 在 2008 年宣布停止 HD DVD 业务，市场预期从此彻底锁定。"),
    ]
    for idx, (title, body) in enumerate(signal_cards):
        card = add_card(
            slide,
            right_left,
            BODY_TOP + 2.40 + idx * 1.20,
            5.78,
            0.98,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_textbox(
            slide,
            right_left + 0.14,
            BODY_TOP + 2.54 + idx * 1.20,
            5.20,
            0.18,
            [{"text": title, "font_name": FONT_BODY, "font_size": 12.2, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            right_left + 0.14,
            BODY_TOP + 2.78 + idx * 1.20,
            5.30,
            0.22,
            [{"text": body, "font_name": FONT_BODY, "font_size": 11.1, "color_rgb": PALETTE["subtitle"]}],
        )

    events = ctx.milestones_frame.query("case == 'HD Optical'").to_dict("records")
    compact_titles = [
        "格式战正式开启",
        "PS3 带来默认装机",
        "Toshiba 宣布退出",
    ]
    for idx, event in enumerate(events):
        tile_left = right_left + idx * 1.96
        tile = add_card(
            slide,
            tile_left,
            BODY_TOP + 6.02,
            1.84,
            1.02,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["sage"], 0.56),
            line_width=0.8,
        )
        add_pill(slide, str(event["year"]), tile_left + 0.10, BODY_TOP + 6.14, 0.64, 0.18, PALETTE["sage"], font_size=9.8)
        add_textbox(
            slide,
            tile_left + 0.10,
            BODY_TOP + 6.40,
            1.56,
            0.34,
            [{"text": compact_titles[idx], "font_name": FONT_BODY, "font_size": 10.8, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            tile_left + 0.10,
            BODY_TOP + 6.72,
            1.56,
            0.16,
            [{"text": event["detail"], "font_name": FONT_BODY, "font_size": 9.5, "color_rgb": PALETTE["subtitle"]}],
        )
    add_footer_note(slide, "Sources: Christ & Slowak on Blu-ray vs HD DVD; historical timeline summarized in the design brief.")


def build_slide_09(slide, spec: dict[str, Any]) -> None:
    """构建 Blu-ray lesson 页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Case 03 | Lesson", PALETTE["sage"])

    stage_specs = [
        ("当用户看不懂技术差距", "技术优越性在消费者层面难以被直接感知时，单点参数的说服力会迅速下降。"),
        ("补充品 gatekeeper 先替市场表态", "内容方、渠道和默认入口会把“哪一边更像未来标准”的信号传给所有后续参与者。"),
        ("捆绑与默认装机会加速预期收敛", "一旦某个方案获得 bundle engine 和 dominant standard 预期，自我强化就会开始。"),
    ]
    box_width = 4.42
    for idx, (title, body) in enumerate(stage_specs):
        left = CONTENT_LEFT + idx * (box_width + 0.56)
        card = add_card(
            slide,
            left,
            BODY_TOP + 0.76,
            box_width,
            2.12,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["sage"], 0.58) if idx == 1 else tint(PALETTE["line"], 0.95),
            line_width=0.9 if idx == 1 else 0.8,
        )
        add_pill(slide, f"Stage {idx + 1}", left + 0.14, BODY_TOP + 0.90, 0.72, 0.18, PALETTE["sage"] if idx == 1 else PALETTE["blue"], font_size=9.8)
        add_textbox(
            slide,
            left + 0.14,
            BODY_TOP + 1.20,
            box_width - 0.28,
            0.36,
            [{"text": title, "font_name": FONT_BODY, "font_size": 13.4, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            left + 0.14,
            BODY_TOP + 1.66,
            box_width - 0.28,
            0.70,
            [{"text": body, "font_name": FONT_BODY, "font_size": 11.4, "color_rgb": PALETTE["subtitle"]}],
        )
        if idx < 2:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + box_width + 0.16),
                Inches(BODY_TOP + 1.56),
                Inches(0.24),
                Inches(0.28),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(*tint(PALETTE["sage"], 0.76))
            arrow.line.color.rgb = RGBColor(*tint(PALETTE["sage"], 0.76))

    memo = add_card(slide, CONTENT_LEFT, BODY_TOP + 3.20, 14.44, 1.84, fill_rgb=tint(PALETTE["warm"], 0.08), line_rgb=tint(PALETTE["warm"], 0.66))
    add_card_text(
        memo,
        [
            {"text": "现代含义：当产品差异对大多数买方并不透明时，市场会让生态替自己做判断。", "font_name": FONT_BODY, "font_size": 15, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "这意味着平台战、AI 工具战、终端能力战里，内容/数据联盟、默认装机、渠道和政策杠杆，经常比单点性能更早改变 adoption 曲线。", "font_name": FONT_BODY, "font_size": 12.4, "color_rgb": PALETTE["subtitle"], "space_before": 8},
        ],
        margin_left=0.18,
        margin_right=0.16,
        margin_top=0.16,
    )

    for idx, message in enumerate(
        [
            "补充品 gatekeeper 决定能否形成生态信号。",
            "bundle engine 决定 adopted base 的爬坡速度。",
            "dominant standard 预期决定等待行为何时转成下注行为。",
        ]
    ):
        card = add_card(
            slide,
            CONTENT_LEFT + idx * 4.88,
            BODY_TOP + 5.32,
            4.48,
            0.76,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["line"], 0.95),
            line_width=0.8,
        )
        add_textbox(
            slide,
            CONTENT_LEFT + idx * 4.88 + 0.14,
            BODY_TOP + 5.56,
            4.08,
            0.24,
            [{"text": message, "font_name": FONT_BODY, "font_size": 10.9, "bold": True, "color_rgb": PALETTE["title"]}],
        )
    add_footer_note(slide, "第三个案例强调：当技术差距不透明时，联盟、捆绑和预期会替市场做决定。")


def build_slide_10(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建三案合并矩阵页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Cross-Case Matrix", PALETTE["blue"])

    headers = ["关键维度", "Betamax vs VHS", "OSI vs TCP/IP", "Blu-ray vs HD DVD", "真正被放大的规律"]
    rows: list[list[dict[str, Any]]] = []
    winner_color = tint(PALETTE["blue"], 0.08)
    neutral_color = tint(PALETTE["surface_alt"], 0.95)
    for _, row in ctx.comparison_frame.iterrows():
        row_cells = [
            {"text": row["factor"], "font_size": 11.4, "bold": True, "color_rgb": PALETTE["title"], "fill_rgb": tint(PALETTE["surface_alt"], 0.90)},
            {"text": row["home_video"], "font_size": 11.4, "bold": True if row["home_video"] == "VHS" else False, "align": PP_ALIGN.CENTER, "color_rgb": PALETTE["title"], "fill_rgb": winner_color if row["home_video"] == "VHS" else neutral_color},
            {"text": row["networking"], "font_size": 11.4, "bold": True if row["networking"] == "TCP/IP" else False, "align": PP_ALIGN.CENTER, "color_rgb": PALETTE["title"], "fill_rgb": winner_color if row["networking"] == "TCP/IP" else neutral_color},
            {"text": row["hd_optical"], "font_size": 11.4, "bold": True if row["hd_optical"] == "Blu-ray" else False, "align": PP_ALIGN.CENTER, "color_rgb": PALETTE["title"], "fill_rgb": winner_color if row["hd_optical"] == "Blu-ray" else neutral_color},
            {"text": row["takeaway"], "font_size": 11.2, "color_rgb": PALETTE["subtitle"]},
        ]
        rows.append(row_cells)

    add_simple_table(
        slide,
        left=CONTENT_LEFT,
        top=BODY_TOP,
        col_widths=[2.12, 2.48, 2.44, 2.54, 4.86],
        row_height=0.80,
        headers=headers,
        rows=rows,
        accent_rgb=PALETTE["blue"],
    )

    theorem = add_card(slide, CONTENT_LEFT, 7.08, 14.44, 1.20, fill_rgb=tint(PALETTE["warm"], 0.10), line_rgb=tint(PALETTE["warm"], 0.66))
    add_icon(slide, "binary-tree", CONTENT_LEFT + 0.20, 7.34, 0.24)
    add_card_text(
        theorem,
        [
            {"text": "三条 claim 最终都在说同一件事：输家更容易在局部性能上显得更强，赢家更早闭合生态协调与 adoption risk。", "font_name": FONT_BODY, "font_size": 14.2, "bold": True, "color_rgb": PALETTE["title"]},
        ],
        margin_left=0.56,
        margin_right=0.16,
        margin_top=0.20,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    add_footer_note(slide, "这一页把三场战争压成显式 matrix，避免“每个案例都很特殊”掩盖重复规律。")


def build_slide_11(slide, spec: dict[str, Any], ctx: DeckContext) -> None:
    """构建管理层五问页。"""
    add_slide_header(slide, spec["slide_id"], spec["title"], spec["key_message"], "Management Checklist", PALETTE["warm"])

    positions = [
        (CONTENT_LEFT + 0.10, BODY_TOP + 0.40),
        (CONTENT_LEFT + 4.96, BODY_TOP + 0.10),
        (CONTENT_LEFT + 9.82, BODY_TOP + 0.40),
        (CONTENT_LEFT + 2.54, BODY_TOP + 2.46),
        (CONTENT_LEFT + 7.40, BODY_TOP + 2.46),
    ]
    question_icons = ["route", "device-desktop-analytics", "circles-relation", "rocket", "target-arrow"]
    for index, row in ctx.factor_frame.iterrows():
        left, top = positions[index]
        card = add_card(
            slide,
            left,
            top,
            4.32,
            1.84,
            fill_rgb=(255, 255, 255),
            line_rgb=tint(PALETTE["warm"], 0.68) if index in {0, 2, 4} else tint(PALETTE["blue"], 0.68),
            line_width=0.9,
        )
        add_pill(slide, f"Q{index + 1}", left + 0.14, top + 0.14, 0.48, 0.18, PALETTE["warm"] if index in {0, 2, 4} else PALETTE["blue"], font_size=9.8)
        add_icon(slide, question_icons[index], left + 0.16, top + 0.52, 0.22)
        add_textbox(
            slide,
            left + 0.46,
            top + 0.42,
            3.64,
            0.36,
            [{"text": row["management_question"], "font_name": FONT_BODY, "font_size": 13, "bold": True, "color_rgb": PALETTE["title"]}],
        )
        add_textbox(
            slide,
            left + 0.46,
            top + 0.92,
            3.64,
            0.52,
            [{"text": row["why_it_matters"], "font_name": FONT_BODY, "font_size": 11.3, "color_rgb": PALETTE["subtitle"]}],
        )

    arrow_pairs = [
        (CONTENT_LEFT + 4.44, BODY_TOP + 0.88, 0.32, 0.18),
        (CONTENT_LEFT + 9.32, BODY_TOP + 0.88, 0.32, 0.18),
        (CONTENT_LEFT + 5.88, BODY_TOP + 2.08, 0.18, 0.28),
        (CONTENT_LEFT + 10.74, BODY_TOP + 2.08, 0.18, 0.28),
    ]
    for left, top, width, height in arrow_pairs:
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON if width > height else MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(*tint(PALETTE["warm"], 0.78))
        arrow.line.color.rgb = RGBColor(*tint(PALETTE["warm"], 0.78))

    memo = add_card(slide, CONTENT_LEFT, 6.00, 14.44, 1.48, fill_rgb=tint(PALETTE["blue"], 0.08), line_rgb=tint(PALETTE["blue"], 0.62))
    add_card_text(
        memo,
        [
            {"text": "Apply now", "font_name": FONT_BODY, "font_size": 11.5, "bold": True, "color_rgb": PALETTE["blue"]},
            {"text": "投资时看谁更容易吃到 installed base 与 bundle engine。产品决策时看谁更顺着用户任务和现有流程演进。平台决策时看谁能更快拉齐 complementors 与渠道。", "font_name": FONT_BODY, "font_size": 12.6, "color_rgb": PALETTE["title"], "space_before": 7},
        ],
        margin_left=0.18,
        margin_right=0.16,
        margin_top=0.14,
    )
    add_footer_note(slide, "这一页是整套 deck 的方法论出口：它把历史故事转成了今天可直接执行的管理问题。")


def build_slide_12(slide, spec: dict[str, Any]) -> None:
    """构建结尾页。"""
    add_background_orbs(
        slide,
        [
            {"left": -0.8, "top": 4.8, "width": 4.0, "height": 4.0, "color_rgb": tint(PALETTE["warm"], 0.30), "transparency": 0.68},
            {"left": 11.2, "top": -0.3, "width": 5.0, "height": 5.0, "color_rgb": tint(PALETTE["blue"], 0.26), "transparency": 0.72},
        ],
    )
    add_pill(slide, "Closing Statement", CONTENT_LEFT, 0.44, 1.32, 0.24, PALETTE["warm"], font_size=10.5)
    add_textbox(
        slide,
        CONTENT_LEFT,
        1.18,
        9.10,
        1.24,
        [
            {"text": "更好的技术，", "font_name": FONT_TITLE, "font_size": 29, "bold": True, "color_rgb": PALETTE["title"]},
            {"text": "只有在它同时更容易被一起采用时，", "font_name": FONT_TITLE, "font_size": 29, "bold": True, "color_rgb": PALETTE["title"], "space_before": 5},
            {"text": "才更可能成为标准。", "font_name": FONT_TITLE, "font_size": 29, "bold": True, "color_rgb": PALETTE["title"], "space_before": 5},
        ],
    )
    add_textbox(
        slide,
        CONTENT_LEFT,
        3.12,
        8.80,
        0.46,
        [{"text": "市场不奖励孤立的技术优雅，市场奖励被整个生态低风险采纳的能力。", "font_name": FONT_BODY, "font_size": 15.5, "bold": True, "color_rgb": PALETTE["subtitle"]}],
    )
    add_case_card(
        slide,
        left=9.40,
        top=1.34,
        width=5.02,
        height=1.54,
        accent_rgb=PALETTE["warm"],
        icon_name="presentation",
        case_label="CLAIM 1",
        title="任务匹配与联盟速度会压过峰值规格",
        body="Betamax vs VHS 把第一条 claim 坐实：市场先为主流任务与阵营扩张投票。",
    )
    add_case_card(
        slide,
        left=9.60,
        top=3.10,
        width=4.82,
        height=1.54,
        accent_rgb=PALETTE["blue"],
        icon_name="briefcase",
        case_label="CLAIM 2",
        title="先跑起来并得到反馈的体系会定义现实",
        body="OSI vs TCP/IP 把第二条 claim 坐实：deployment loop 比理论完整性更早决定 adoption。",
    )
    add_case_card(
        slide,
        left=9.40,
        top=4.86,
        width=5.02,
        height=1.54,
        accent_rgb=PALETTE["sage"],
        icon_name="network",
        case_label="CLAIM 3",
        title="补充品与 bundle 会在不透明市场替用户下注",
        body="Blu-ray vs HD DVD 把第三条 claim 坐实：联盟、默认装机与预期会替用户判断。",
    )
    add_textbox(
        slide,
        CONTENT_LEFT,
        5.28,
        7.80,
        1.22,
        [
            {"text": "三条 claim 给出的共同答案是：", "font_name": FONT_BODY, "font_size": 12.5, "bold": True, "color_rgb": PALETTE["subtitle"]},
            {"text": "如果你只能记住一句话，就记住 deck 开头那句总论点。它不是修辞，它是三条 claim 共同指向的市场判断公式。", "font_name": FONT_BODY, "font_size": 13.2, "color_rgb": PALETTE["title"], "space_before": 8},
        ],
    )
    add_footer_note(slide, "End state: 规格表只是开始，生态采纳能力才是标准战的真正结算项。")
    add_textbox(
        slide,
        CONTENT_RIGHT - 0.56,
        8.50,
        0.48,
        0.16,
        [{"text": spec["slide_id"], "font_name": FONT_BODY, "font_size": 10, "bold": True, "align": PP_ALIGN.RIGHT, "color_rgb": PALETTE["muted"]}],
    )


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------


def build_deck() -> Path:
    """构建完整 deck，并输出 build manifest。"""
    ctx = load_context()
    prs = new_presentation()

    slide_sequence = [
        ("S01", build_slide_01),
        ("S02", build_slide_02),
        ("S03", lambda slide, spec: build_slide_03(slide, spec, ctx)),
        ("S04", lambda slide, spec: build_slide_04(slide, spec, ctx)),
        ("S05", lambda slide, spec: build_slide_05(slide, spec, ctx)),
        ("S06", lambda slide, spec: build_slide_06(slide, spec, ctx)),
        ("S07", build_slide_07),
        ("S08", lambda slide, spec: build_slide_08(slide, spec, ctx)),
        ("S09", build_slide_09),
        ("S10", lambda slide, spec: build_slide_10(slide, spec, ctx)),
        ("S11", lambda slide, spec: build_slide_11(slide, spec, ctx)),
        ("S12", build_slide_12),
    ]

    for slide_id, builder in slide_sequence:
        spec = ctx.slides_by_id[slide_id]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        builder(slide, spec)

    output_path = save_presentation(prs, PPTX_PATH)
    apply_font_policy_to_pptx(output_path)
    manifest = {
        "deck_title": ctx.deck_spec["deck"]["title"],
        "workspace": str(WORKSPACE_DIR),
        "pptx": str(output_path),
        "slide_count": len(prs.slides),
        "slides": [{"slide_id": slide_id, "title": ctx.slides_by_id[slide_id]["title"]} for slide_id, _ in slide_sequence],
        "inputs": {
            "brief": str(BRIEF_PATH),
            "narrative": str(NARRATIVE_PATH),
            "slide_specs": str(GENERATED_SPEC_PATH),
        },
        "assets": {
            "icons": str(ICON_THEME_DIR),
            "factor_csv": str(WORKSPACE_DIR / "data" / "processed" / "five_factor_framework.csv"),
            "comparison_csv": str(WORKSPACE_DIR / "data" / "processed" / "comparison_matrix.csv"),
            "milestones_csv": str(WORKSPACE_DIR / "data" / "processed" / "case_milestones.csv"),
        },
        "font_policy": {
            "latin": FONT_LATIN,
            "east_asia": FONT_EAST_ASIA,
        },
    }
    MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path


def main() -> int:
    """执行完整 build。"""
    output_path = build_deck()
    print(f"[OK] 已生成 deck: {output_path}")
    print(f"[INFO] manifest: {MANIFEST_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
