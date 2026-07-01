"""生成复杂股票分析框架的可编辑 PPT Demo。

定位与作用
----------
本脚本用于演示 `ppt-complex-diagram-collab` 技能的核心能力：
1. 将复杂业务架构用 PowerPoint 形状和连接线程序化表达；
2. 用“分层卡片 + 递进箭头”的版式表达层级递进，不使用 connector 连线；
3. 产出可交由连接器校验脚本做程序化验收的 pptx 文件。

大致流程
----------
1. 初始化单页画布与标题；
2. 绘制四层架构卡片（每层含若干节点）；
3. 绘制层间递进箭头与协作入口；
4. 保存到 `pptx/stock_architecture_complex_demo.pptx`。
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


@dataclass
class NodeStyle:
    """节点样式配置。"""

    fill_rgb: tuple[int, int, int]
    line_rgb: tuple[int, int, int]

@dataclass(frozen=True)
class LayerSpec:
    """每一层架构卡片的描述。"""

    key: str
    title: str
    accent_rgb: tuple[int, int, int]
    nodes: list[str]


def pick_contrast_text_rgb(fill_rgb: tuple[int, int, int]) -> tuple[int, int, int]:
    """根据对比度自动选择文字色（WCAG 思路）。"""
    dark = (25, 38, 56)
    light = (255, 255, 255)

    def _srgb_to_linear(channel_8bit: int) -> float:
        c = channel_8bit / 255.0
        if c <= 0.04045:
            return c / 12.92
        return ((c + 0.055) / 1.055) ** 2.4

    def _relative_luminance(rgb: tuple[int, int, int]) -> float:
        red, green, blue = rgb
        r = _srgb_to_linear(red)
        g = _srgb_to_linear(green)
        b = _srgb_to_linear(blue)
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    def _contrast_ratio(rgb1: tuple[int, int, int], rgb2: tuple[int, int, int]) -> float:
        l1 = _relative_luminance(rgb1)
        l2 = _relative_luminance(rgb2)
        lighter, darker = (l1, l2) if l1 >= l2 else (l2, l1)
        return (lighter + 0.05) / (darker + 0.05)

    return dark if _contrast_ratio(fill_rgb, dark) >= _contrast_ratio(fill_rgb, light) else light


def tint(rgb: tuple[int, int, int], alpha: float) -> tuple[int, int, int]:
    """将颜色按 alpha 与白色混合（alpha 越小越接近白色）。"""
    red, green, blue = rgb
    r = int(round(alpha * red + (1 - alpha) * 255))
    g = int(round(alpha * green + (1 - alpha) * 255))
    b = int(round(alpha * blue + (1 - alpha) * 255))
    return (r, g, b)


def add_card(
    slide,
    title: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent_rgb: tuple[int, int, int],
) -> None:
    """新增分层卡片（背景 + 顶部标题条）。"""
    outer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    outer.text = ""
    outer.fill.solid()
    outer.fill.fore_color.rgb = RGBColor(*tint(accent_rgb, 0.14))
    outer.line.color.rgb = RGBColor(*tint(accent_rgb, 0.55))
    outer.line.width = Pt(1.0)

    header_h = 0.32
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(header_h),
    )
    header.text = title
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(*accent_rgb)
    header.line.color.rgb = RGBColor(*accent_rgb)
    header.line.width = Pt(0.5)

    para = header.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(13)
    para.font.color.rgb = RGBColor(*pick_contrast_text_rgb(accent_rgb))


def add_node(slide, key: str, text: str, left: float, top: float, width: float, height: float, style: NodeStyle):
    """新增业务节点并返回 shape。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.name = key
    shape.text = text
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*style.fill_rgb)
    shape.line.color.rgb = RGBColor(*style.line_rgb)
    shape.line.width = Pt(1.25)
    para = shape.text_frame.paragraphs[0]
    para.font.size = Pt(10.5)
    para.font.bold = True
    para.font.color.rgb = RGBColor(*pick_contrast_text_rgb(style.fill_rgb))
    return shape


def add_progress_arrow(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    arrow_text: str,
    fill_rgb: tuple[int, int, int],
    line_rgb: tuple[int, int, int],
) -> None:
    """新增用于表达层级递进关系的箭头图形。"""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    arrow.text = arrow_text
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(*fill_rgb)
    arrow.line.color.rgb = RGBColor(*line_rgb)
    arrow.line.width = Pt(1.25)
    para = arrow.text_frame.paragraphs[0]
    para.font.bold = True
    para.font.size = Pt(10)
    para.font.color.rgb = RGBColor(*pick_contrast_text_rgb(fill_rgb))


def build_ppt(output_path: Path) -> Path:
    """构建 demo PPT 并返回输出路径。"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(15.2), Inches(0.5))
    title_box.text_frame.text = "Frame-based Equity Intelligence & Execution Framework"
    title_para = title_box.text_frame.paragraphs[0]
    title_para.font.size = Pt(22)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(28, 45, 74)

    subtitle = slide.shapes.add_textbox(Inches(0.3), Inches(0.62), Inches(15.2), Inches(0.32))
    subtitle.text_frame.text = "Data -> Representation -> Portfolio -> Execution -> Attribution Feedback"
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.size = Pt(12)
    subtitle_para.font.color.rgb = RGBColor(93, 109, 136)

    palette = {
        "data": (37, 99, 235),  # blue
        "research": (16, 185, 129),  # emerald
        "portfolio": (245, 158, 11),  # amber
        "execution": (124, 58, 237),  # violet
        "neutral_line": (100, 116, 139),  # slate
        "neutral_fill": (148, 163, 184),
    }

    layers: list[LayerSpec] = [
        LayerSpec(
            key="l1",
            title="L1 Data Acquisition & Quality",
            accent_rgb=palette["data"],
            nodes=["Market Data Lake", "Alt Data Hub", "Corporate Events", "Quality & Lineage Gate"],
        ),
        LayerSpec(
            key="l2",
            title="L2 Representation & Signal Factory",
            accent_rgb=palette["research"],
            nodes=["Feature Store", "Factor Engine", "Regime Detector", "Signal Blender", "Model Registry"],
        ),
        LayerSpec(
            key="l3",
            title="L3 Portfolio Construction & Risk",
            accent_rgb=palette["portfolio"],
            nodes=["Constraint Compiler", "Robust Optimizer", "Risk Scenario Lab", "Capacity & Cost Model", "Compliance Rules"],
        ),
        LayerSpec(
            key="l4",
            title="L4 Execution & Attribution",
            accent_rgb=palette["execution"],
            nodes=["Execution Router", "Broker Simulator", "TCA Monitor", "PnL Attribution", "Feedback Scheduler"],
        ),
    ]

    card_left = 0.45
    card_w = 15.1
    card_h = 1.6
    gap = 0.16
    start_top = 1.05

    def _node_label(text: str) -> str:
        if len(text) <= 16:
            return text
        parts = text.split(" ")
        if len(parts) <= 1:
            return text
        return f"{parts[0]}\n{' '.join(parts[1:])}"

    def add_layer(layer: LayerSpec, top: float) -> None:
        add_card(slide, layer.title, card_left, top, card_w, card_h, layer.accent_rgb)
        node_style = NodeStyle((255, 255, 255), tint(layer.accent_rgb, 0.75))

        inner_left = card_left + 0.25
        inner_top = top + 0.45
        inner_w = card_w - 0.5
        node_h = 0.85

        cols = 5 if len(layer.nodes) >= 5 else 4
        col_w = inner_w / cols
        node_w = col_w - 0.16

        for idx, node_text in enumerate(layer.nodes):
            col = idx % cols
            row = idx // cols
            x = inner_left + col * col_w
            y = inner_top + row * (node_h + 0.12)
            add_node(
                slide,
                f"{layer.key}_n{idx}",
                _node_label(node_text),
                x,
                y,
                node_w,
                node_h,
                node_style,
            )

    for i, layer in enumerate(layers):
        top = start_top + i * (card_h + gap)
        add_layer(layer, top)
        if i < len(layers) - 1:
            down = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(7.65),
                Inches(top + card_h + 0.02),
                Inches(0.7),
                Inches(gap - 0.02),
            )
            down.text = ""
            down.fill.solid()
            down.fill.fore_color.rgb = RGBColor(*palette["neutral_fill"])
            down.line.color.rgb = RGBColor(*palette["neutral_line"])
            down.line.width = Pt(1.0)

    role_style = NodeStyle((255, 255, 255), (148, 163, 184))
    add_node(slide, "role_researcher", "Researcher", 0.55, 8.2, 2.25, 0.62, role_style)
    add_node(slide, "role_pm", "PM", 3.0, 8.2, 1.4, 0.62, role_style)
    add_node(slide, "role_risk", "Risk", 4.65, 8.2, 1.6, 0.62, role_style)
    add_progress_arrow(slide, 6.6, 8.2, 3.0, 0.62, "Human-in-the-loop -> Governance", (226, 232, 240), (100, 116, 139))
    add_progress_arrow(slide, 9.85, 8.2, 5.7, 0.62, "Model Registry / Audit / Reproducibility", (226, 232, 240), (100, 116, 139))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    """CLI 入口。"""
    output_path = Path(__file__).parent / "pptx" / "stock_architecture_complex_demo.pptx"
    final_path = build_ppt(output_path)
    print(f"[OK] 已生成: {final_path}")


if __name__ == "__main__":
    main()
