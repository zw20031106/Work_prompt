"""生成 Apple FY2025 财报点评可编辑 PPT。

本脚本读取 `build/generated/slide_specs.yaml` 与 `data/processed/` 中的底稿数据，
用 `python-pptx` 直接生成可编辑的研究报告风格 PPT。整体流程为：
1. 载入核心财务、产品线、地区、毛利率和来源数据；
2. 计算同比、占比、利润率和现金流等页面指标；
3. 按 slide_specs 逐页生成标题区、图表、表格、结论卡片和页脚；
4. 将 PPTX 写入 `build/pptx/` 与 `final/`。
"""

from __future__ import annotations

import json
import shutil
import tempfile
import zipfile
from datetime import date
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

import pandas as pd
import yaml
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


WORKSPACE = Path(__file__).resolve().parents[1]
DATA_DIR = WORKSPACE / "data" / "processed"
BUILD_DIR = WORKSPACE / "build"
PPTX_DIR = BUILD_DIR / "pptx"
FINAL_DIR = WORKSPACE / "final"
SPEC_PATH = BUILD_DIR / "generated" / "slide_specs.yaml"

OUT_PPTX = PPTX_DIR / "apple_fy2025_financial_report_review.pptx"
FINAL_PPTX = FINAL_DIR / "apple_fy2025_financial_report_review.pptx"

PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)

CITIC_RED = RGBColor(155, 23, 35)
DEEP_RED = RGBColor(111, 18, 28)
GOLD = RGBColor(180, 145, 70)
DARK_GOLD = RGBColor(132, 104, 45)
INK = RGBColor(35, 31, 32)
GRAY = RGBColor(92, 92, 92)
LIGHT_GRAY = RGBColor(247, 244, 239)
MID_GRAY = RGBColor(217, 208, 194)
TABLE_ALT = RGBColor(252, 249, 244)
TEXT = INK
WHITE = RGBColor(255, 255, 255)
WARM_BROWN = RGBColor(126, 80, 62)
NAVY = DEEP_RED
RED = CITIC_RED
GREEN = DARK_GOLD
BLUE = WARM_BROWN
PAPER = RGBColor(254, 253, 250)
PANEL_BG = RGBColor(255, 254, 251)
PANEL_LINE = RGBColor(218, 210, 196)
GRID_LINE = RGBColor(225, 221, 213)
MUTED_TEXT = RGBColor(112, 108, 101)
SOFT_RED_BG = RGBColor(249, 241, 241)
SOFT_GOLD_BG = RGBColor(248, 245, 235)
REPORT_LEFT = 0.48
REPORT_RIGHT = 12.86
REPORT_WIDTH = REPORT_RIGHT - REPORT_LEFT
CONTENT_TOP = 1.08
FOOTER_Y = 7.08

FONT_CN = "SimSun"
FONT_CN_DISPLAY = "宋体"
FONT_EN = "Times New Roman"
BODY_FONT_PT = 12
TABLE_FONT_PT = 10.5
BODY_LINE_SPACING = 1.5
BODY_SPACE_PT = 6
FIRST_LINE_INDENT_IN = 0.24


def load_inputs() -> dict[str, pd.DataFrame | dict]:
    """读取所有构建输入。"""

    return {
        "spec": yaml.safe_load(SPEC_PATH.read_text(encoding="utf-8")),
        "financials": pd.read_csv(DATA_DIR / "apple_financials_fy2021_fy2025.csv"),
        "products": pd.read_csv(DATA_DIR / "apple_product_net_sales_fy2021_fy2025.csv"),
        "regions": pd.read_csv(DATA_DIR / "apple_region_net_sales_fy2021_fy2025.csv"),
        "gross_margin": pd.read_csv(DATA_DIR / "apple_gross_margin_by_type_fy2021_fy2025.csv"),
        "sources": pd.read_csv(DATA_DIR / "apple_10k_sources.csv"),
    }


def usd_bn(value: float) -> float:
    """美元金额转为十亿美元。"""

    return float(value) / 1_000_000_000


def pct(value: float) -> float:
    """比例转百分数。"""

    return float(value) * 100


def yoy(series: pd.Series, current_year: int, previous_year: int) -> float:
    """计算同比增速。"""

    current = float(series.loc[current_year])
    previous = float(series.loc[previous_year])
    return current / previous - 1


def make_prs() -> Presentation:
    """创建宽屏 PPT。"""

    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H
    return prs


def blank_slide(prs: Presentation):
    """添加空白页。"""

    return prs.slides.add_slide(prs.slide_layouts[6])


def set_text_frame_style(shape, margin: float = 0.06, vertical: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    """设置文本框基础边距与垂直对齐。"""

    tf = shape.text_frame
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = vertical
    tf.word_wrap = True


def apply_mixed_font(run, size: float, color: RGBColor, bold: bool = False) -> None:
    """设置中英混排字体：中文宋体，英文 Times New Roman。"""

    run.font.name = FONT_EN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", FONT_EN), ("a:ea", FONT_CN_DISPLAY), ("a:cs", FONT_EN)):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", typeface)


def apply_paragraph_style(
    paragraph,
    size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    line_spacing: float = BODY_LINE_SPACING,
    space_before_pt: float = BODY_SPACE_PT,
    space_after_pt: float = BODY_SPACE_PT,
    first_line_indent: bool = False,
) -> None:
    """设置段落字号、行距、段前段后和首行缩进。"""

    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_before = Pt(space_before_pt)
    paragraph.space_after = Pt(space_after_pt)
    if first_line_indent:
        paragraph.first_line_indent = Inches(FIRST_LINE_INDENT_IN)
    for run in paragraph.runs:
        apply_mixed_font(run, size=size, color=color, bold=bold)


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 11,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    fill: RGBColor | None = None,
    margin: float = 0.05,
    line_spacing: float | None = None,
    space_before_pt: float | None = None,
    space_after_pt: float | None = None,
    first_line_indent: bool | None = None,
):
    """添加单段文本框。"""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_frame_style(box, margin=margin)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    is_body = size >= 11 and not bold
    apply_paragraph_style(
        p,
        size=size,
        color=color,
        bold=bold,
        align=align,
        line_spacing=line_spacing if line_spacing is not None else (BODY_LINE_SPACING if is_body else 1.0),
        space_before_pt=space_before_pt if space_before_pt is not None else (BODY_SPACE_PT if is_body else 0),
        space_after_pt=space_after_pt if space_after_pt is not None else (BODY_SPACE_PT if is_body else 0),
        first_line_indent=first_line_indent if first_line_indent is not None else (is_body and align == PP_ALIGN.LEFT),
    )
    return box


def add_multiline(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Iterable[str],
    size: float = BODY_FONT_PT,
    color: RGBColor = TEXT,
    bullet: bool = False,
    leading: float = BODY_LINE_SPACING,
    space_before_pt: float = BODY_SPACE_PT,
    space_after_pt: float = BODY_SPACE_PT,
    first_line_indent: bool = True,
):
    """添加多行正文文本。"""

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text_frame_style(box, margin=0.06)
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"· {line}" if bullet else line
        apply_paragraph_style(
            p,
            size=size,
            color=color,
            line_spacing=leading,
            space_before_pt=space_before_pt,
            space_after_pt=space_after_pt,
            first_line_indent=first_line_indent,
        )
    return box


def add_rect(slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None):
    """添加矩形。"""

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_title(slide, title: str, page_no: int) -> None:
    """添加正文页标题区。"""

    add_rect(slide, 0, 0, 13.333, 0.22, NAVY)
    add_text(slide, 0.48, 0.34, 10.6, 0.44, title, size=18.5, color=NAVY, bold=True, margin=0)
    add_rect(slide, 0.48, 0.86, 1.25, 0.035, RED)
    add_text(slide, 11.7, 0.35, 1.15, 0.25, f"{page_no:02d}", size=8.5, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)


def add_footer(slide, page_no: int, source: str = "数据来源：SEC companyfacts、Apple 10-K；金额单位除特别说明外为十亿美元。") -> None:
    """添加页脚。"""

    add_rect(slide, 0.48, 7.08, 12.35, 0.01, MID_GRAY)
    add_text(slide, 0.48, 7.13, 8.8, 0.18, source, size=6.6, color=GRAY, margin=0)
    add_text(slide, 9.1, 7.13, 3.7, 0.18, "仅供学术交流使用，不代表任何组织机构、机构和个人观点和立场", size=6.3, color=GRAY, align=PP_ALIGN.RIGHT, margin=0)


def add_kpi_card(slide, x: float, y: float, w: float, h: float, label: str, value: str, note: str, accent: RGBColor = RED) -> None:
    """添加关键指标卡片。"""

    shape = add_rect(slide, x, y, w, h, LIGHT_GRAY, MID_GRAY)
    shape.line.width = Pt(0.8)
    add_rect(slide, x, y, 0.05, h, accent)
    add_text(slide, x + 0.15, y + 0.12, w - 0.25, 0.18, label, size=8.4, color=GRAY, bold=True, margin=0)
    add_text(slide, x + 0.15, y + 0.38, w - 0.25, 0.34, value, size=18, color=NAVY, bold=True, margin=0)
    add_text(slide, x + 0.15, y + 0.82, w - 0.25, 0.36, note, size=7.4, color=GRAY, margin=0)


def add_note_box(slide, x: float, y: float, w: float, h: float, title: str, lines: list[str], accent: RGBColor = NAVY) -> None:
    """添加观点说明框。"""

    add_rect(slide, x, y, w, h, RGBColor(255, 253, 248), MID_GRAY)
    add_text(slide, x + 0.16, y + 0.12, w - 0.32, 0.3, title, size=13, color=accent, bold=True, margin=0)
    add_multiline(slide, x + 0.16, y + 0.52, w - 0.32, h - 0.6, lines, size=BODY_FONT_PT, color=TEXT)


def add_chart(
    slide,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    categories: list[str],
    series: dict[str, list[float]],
    legend: bool = True,
    value_format: str = "0",
):
    """添加原生 Office chart。"""

    data = CategoryChartData()
    data.categories = categories
    for name, values in series.items():
        data.add_series(name, values)
    graphic = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data)
    chart = graphic.chart
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.number_format = value_format
    chart.value_axis.tick_labels.font.size = Pt(9.5)
    chart.value_axis.tick_labels.font.name = FONT_EN
    chart.category_axis.tick_labels.font.size = Pt(9.5)
    chart.category_axis.tick_labels.font.name = FONT_EN
    palette = [CITIC_RED, DARK_GOLD, WARM_BROWN, GRAY, RGBColor(184, 174, 151), DEEP_RED]
    for idx, series_item in enumerate(chart.series):
        color = palette[idx % len(palette)]
        series_item.format.fill.solid()
        series_item.format.fill.fore_color.rgb = color
        series_item.format.line.color.rgb = color
    return chart


def add_table(slide, x: float, y: float, w: float, h: float, headers: list[str], rows: list[list[str]], font_size: float = TABLE_FONT_PT):
    """添加原生表格。"""

    font_size = TABLE_FONT_PT
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_RED
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            apply_paragraph_style(
                p,
                size=font_size,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
                line_spacing=1.0,
                space_before_pt=0,
                space_after_pt=0,
            )
    for ridx, row in enumerate(rows, start=1):
        for cidx, value in enumerate(row):
            cell = table.cell(ridx, cidx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ridx % 2 else TABLE_ALT
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                apply_paragraph_style(
                    p,
                    size=font_size,
                    color=TEXT,
                    line_spacing=1.0,
                    space_before_pt=0,
                    space_after_pt=0,
                )
    return table_shape


def add_report_shell(slide, title: str, page_no: int | None, source: str) -> None:
    """添加研报页共享页眉、标题区与页脚。"""

    add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
    add_rect(slide, 0, 0, 13.333, 0.115, NAVY)
    add_rect(slide, 0, 0.115, 13.333, 0.014, GOLD)
    add_text(
        slide,
        REPORT_LEFT,
        0.29,
        11.25,
        0.38,
        title,
        size=18.2,
        color=NAVY,
        bold=True,
        margin=0,
        line_spacing=1.0,
        first_line_indent=False,
    )
    add_rect(slide, REPORT_LEFT, 0.82, 0.78, 0.022, RED)
    add_rect(slide, REPORT_LEFT + 0.80, 0.82, 0.45, 0.022, GOLD)
    if page_no is not None:
        add_text(
            slide,
            12.08,
            0.30,
            0.72,
            0.18,
            f"{page_no:02d}",
            size=7.2,
            color=MUTED_TEXT,
            align=PP_ALIGN.RIGHT,
            margin=0,
            first_line_indent=False,
        )
    add_report_footer(slide, source)


def add_report_footer(slide, source: str) -> None:
    """添加正式研报风格页脚。"""

    add_rect(slide, REPORT_LEFT, FOOTER_Y, REPORT_WIDTH, 0.008, MID_GRAY)
    add_text(
        slide,
        REPORT_LEFT,
        FOOTER_Y + 0.055,
        7.8,
        0.18,
        source,
        size=6.1,
        color=MUTED_TEXT,
        margin=0,
        first_line_indent=False,
    )
    add_text(
        slide,
        9.0,
        FOOTER_Y + 0.055,
        3.86,
        0.18,
        "仅供学术交流使用，不代表任何组织机构、机构和个人观点和立场",
        size=6.1,
        color=MUTED_TEXT,
        align=PP_ALIGN.RIGHT,
        margin=0,
        first_line_indent=False,
    )


def add_report_panel(slide, x: float, y: float, w: float, h: float, fill: RGBColor = PANEL_BG) -> None:
    """添加无投影的研报内容面板。"""

    panel = add_rect(slide, x, y, w, h, fill, PANEL_LINE)
    panel.line.width = Pt(0.55)


def add_report_label(slide, x: float, y: float, text: str, color: RGBColor = NAVY) -> None:
    """添加小型栏目标签。"""

    add_rect(slide, x, y + 0.03, 0.05, 0.20, color)
    add_text(
        slide,
        x + 0.10,
        y,
        3.8,
        0.26,
        text,
        size=9.0,
        color=color,
        bold=True,
        margin=0,
        line_spacing=1.0,
        first_line_indent=False,
    )


def add_takeaway_strip(slide, x: float, y: float, w: float, h: float, title: str, body: str) -> None:
    """添加页内核心判断条。"""

    add_report_panel(slide, x, y, w, h, SOFT_RED_BG)
    add_rect(slide, x, y, 0.045, h, RED)
    add_text(
        slide,
        x + 0.16,
        y + 0.08,
        w - 0.28,
        0.20,
        title,
        size=9.0,
        color=NAVY,
        bold=True,
        margin=0,
        line_spacing=1.0,
        first_line_indent=False,
    )
    add_text(
        slide,
        x + 0.16,
        y + 0.34,
        w - 0.28,
        h - 0.42,
        body,
        size=9.4,
        color=TEXT,
        margin=0,
        line_spacing=1.25,
        space_before_pt=0,
        space_after_pt=0,
        first_line_indent=False,
    )


def add_report_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    headers: list[str],
    rows: list[list[str]],
    font_size: float = 7.3,
    header_fill: RGBColor = NAVY,
    col_widths: list[float] | None = None,
    numeric_from_col: int = 1,
    numeric_columns: set[int] | None = None,
):
    """添加紧凑研报表格。"""

    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for ridx in range(len(rows) + 1):
        table.rows[ridx].height = Inches(h / (len(rows) + 1))
    for cidx, header in enumerate(headers):
        cell = table.cell(0, cidx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.015)
        cell.margin_bottom = Inches(0.015)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            apply_paragraph_style(
                paragraph,
                size=font_size,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
                line_spacing=1.0,
                space_before_pt=0,
                space_after_pt=0,
            )
    for ridx, row in enumerate(rows, start=1):
        for cidx, value in enumerate(row):
            cell = table.cell(ridx, cidx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ridx % 2 else TABLE_ALT
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.012)
            cell.margin_bottom = Inches(0.012)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                is_numeric = cidx in numeric_columns if numeric_columns is not None else cidx >= numeric_from_col
                apply_paragraph_style(
                    paragraph,
                    size=font_size,
                    color=TEXT,
                    align=PP_ALIGN.RIGHT if is_numeric else PP_ALIGN.LEFT,
                    line_spacing=1.0,
                    space_before_pt=0,
                    space_after_pt=0,
                    first_line_indent=False,
                )
    return table_shape


def style_native_chart(chart, palette: list[RGBColor]) -> None:
    """统一原生 Office chart 的字体、网格线与系列颜色。"""

    chart.has_title = False
    chart.chart_style = 10
    try:
        chart.value_axis.major_gridlines.format.line.color.rgb = GRID_LINE
        chart.value_axis.major_gridlines.format.line.width = Pt(0.35)
    except Exception:
        pass
    for axis in (getattr(chart, "value_axis", None), getattr(chart, "category_axis", None)):
        if axis is None:
            continue
        try:
            axis.tick_labels.font.size = Pt(7.8)
            axis.tick_labels.font.name = FONT_EN
            axis.format.line.color.rgb = GRID_LINE
            axis.format.line.width = Pt(0.35)
        except Exception:
            pass
    if chart.has_legend:
        try:
            chart.legend.font.size = Pt(7.4)
            chart.legend.font.name = FONT_EN
        except Exception:
            pass
    for idx, series_item in enumerate(chart.series):
        color = palette[idx % len(palette)]
        try:
            series_item.format.fill.solid()
            series_item.format.fill.fore_color.rgb = color
        except Exception:
            pass
        try:
            series_item.format.line.color.rgb = color
            series_item.format.line.width = Pt(1.6)
        except Exception:
            pass


def add_report_chart(
    slide,
    chart_type: XL_CHART_TYPE,
    x: float,
    y: float,
    w: float,
    h: float,
    figure_title: str,
    unit: str,
    categories: list[str],
    series: dict[str, list[float]],
    value_format: str,
    legend: bool = True,
    palette: list[RGBColor] | None = None,
):
    """添加带图题和单位的研报图表面板。"""

    add_report_panel(slide, x, y, w, h, WHITE)
    add_text(
        slide,
        x + 0.16,
        y + 0.10,
        w - 1.6,
        0.20,
        figure_title,
        size=8.2,
        color=NAVY,
        bold=True,
        margin=0,
        first_line_indent=False,
    )
    add_text(
        slide,
        x + w - 1.38,
        y + 0.10,
        1.18,
        0.18,
        unit,
        size=6.8,
        color=MUTED_TEXT,
        align=PP_ALIGN.RIGHT,
        margin=0,
        first_line_indent=False,
    )
    chart = add_chart(
        slide,
        chart_type,
        x + 0.22,
        y + 0.44,
        w - 0.44,
        h - 0.66,
        categories,
        series,
        legend=legend,
        value_format=value_format,
    )
    style_native_chart(chart, palette or [RED, DARK_GOLD, WARM_BROWN, GRAY, RGBColor(184, 174, 151), DEEP_RED])
    return chart


def add_right_brief(slide, x: float, y: float, w: float, h: float, title: str, points: list[str]) -> None:
    """添加右侧研报摘要框。"""

    add_report_panel(slide, x, y, w, h, PAPER)
    add_rect(slide, x, y, w, 0.06, NAVY)
    add_text(
        slide,
        x + 0.16,
        y + 0.18,
        w - 0.32,
        0.24,
        title,
        size=10.3,
        color=NAVY,
        bold=True,
        margin=0,
        first_line_indent=False,
    )
    add_multiline(
        slide,
        x + 0.16,
        y + 0.55,
        w - 0.32,
        h - 0.72,
        points,
        size=BODY_FONT_PT,
        color=TEXT,
        bullet=True,
        leading=BODY_LINE_SPACING,
        space_before_pt=BODY_SPACE_PT,
        space_after_pt=BODY_SPACE_PT,
        first_line_indent=False,
    )


def add_body_panel(slide, x: float, y: float, w: float, h: float, title: str, paragraphs: list[str]) -> None:
    """添加正式正文段落面板。"""

    add_report_panel(slide, x, y, w, h, PAPER)
    add_rect(slide, x, y, w, 0.06, NAVY)
    add_text(slide, x + 0.16, y + 0.18, w - 0.32, 0.24, title, size=10.3, color=NAVY, bold=True, margin=0, first_line_indent=False)
    add_multiline(
        slide,
        x + 0.16,
        y + 0.62,
        w - 0.32,
        h - 0.82,
        paragraphs,
        size=BODY_FONT_PT,
        color=TEXT,
        bullet=False,
        leading=BODY_LINE_SPACING,
        space_before_pt=BODY_SPACE_PT,
        space_after_pt=BODY_SPACE_PT,
        first_line_indent=True,
    )


def fmt_bn(value: float) -> str:
    """把美元金额格式化为十亿美元字符串。"""

    return f"{usd_bn(value):.1f}"


def fmt_pct(value: float) -> str:
    """把比例格式化为百分比字符串。"""

    return f"{pct(value):.1f}%"


def pct_points(value: float) -> str:
    """把比例差格式化为百分点字符串。"""

    return f"{pct(value):+.1f}pct"


def slide_01(prs: Presentation, data: dict) -> None:
    """生成封面页。"""

    slide = blank_slide(prs)
    add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
    add_rect(slide, 0, 0, 13.333, 0.16, NAVY)
    add_rect(slide, 0, 0.16, 13.333, 0.018, GOLD)
    add_text(slide, 0.66, 0.55, 3.6, 0.22, "财报点评 | FY2025 Form 10-K", size=9.0, color=NAVY, bold=True, margin=0, first_line_indent=False)
    add_text(slide, 0.64, 1.05, 8.4, 0.62, "Apple Inc. FY2025 财报点评", size=28, color=NAVY, bold=True, margin=0, first_line_indent=False)
    add_text(
        slide,
        0.68,
        1.86,
        8.6,
        0.42,
        "稳健恢复延续，服务化与现金流仍是核心护城河",
        size=16.5,
        color=DARK_GOLD,
        bold=True,
        margin=0,
        first_line_indent=False,
    )
    add_rect(slide, 0.68, 2.48, 1.25, 0.035, RED)
    add_rect(slide, 1.96, 2.48, 0.72, 0.035, GOLD)
    add_takeaway_strip(
        slide,
        0.68,
        3.04,
        7.05,
        1.52,
        "核心判断",
        "FY2025 收入和净利润同步恢复，Services 高增长与高毛利继续支撑经营韧性；但 iPhone 依赖、Greater China 承压和监管变量仍决定后续增长斜率。",
    )
    add_report_table(
        slide,
        8.18,
        1.02,
        4.42,
        2.35,
        ["指标", "FY2025", "同比/变化"],
        [
            ["Revenue", fmt_bn(data["financials"].set_index("fiscal_year").loc[2025, "revenue"]), "+6.4%"],
            ["Net income", fmt_bn(data["financials"].set_index("fiscal_year").loc[2025, "net_income"]), "+19.5%"],
            ["Gross margin", fmt_pct(data["financials"].set_index("fiscal_year").loc[2025, "gross_margin"]), "+0.7pct"],
            ["FCF", fmt_bn(data["financials"].set_index("fiscal_year").loc[2025, "fcf"]), "98.8bn"],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.55, 1.30, 1.57],
        numeric_columns={1, 2},
    )
    add_right_brief(
        slide,
        8.18,
        3.66,
        4.42,
        1.50,
        "材料边界",
        [
            "基于 SEC companyfacts 与 Apple FY2021-FY2025 Form 10-K。",
            "不含评级、目标价、实时行情或未披露经营数据。",
        ],
    )
    add_report_panel(slide, 0.68, 5.55, 11.92, 0.62, SOFT_GOLD_BG)
    add_text(
        slide,
        0.88,
        5.75,
        11.5,
        0.18,
        "仅供学术交流使用，不代表任何组织机构，机构和个人的观点和立场。",
        size=8.8,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        margin=0,
        first_line_indent=False,
    )
    add_text(slide, 0.68, 6.72, 5.9, 0.20, f"生成日期：{date(2026, 5, 10).isoformat()} | Fiscal year ended September 27, 2025", size=6.8, color=MUTED_TEXT, margin=0, first_line_indent=False)


def slide_02(prs: Presentation, data: dict) -> None:
    """生成核心观点页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "FY2025 经营恢复明确，结构改善主要来自服务化和利润率", 2, "数据来源：SEC companyfacts、Apple FY2025 Form 10-K；金额单位为十亿美元。")
    fin = data["financials"].set_index("fiscal_year")
    products = data["products"]
    services_2025 = products[(products.fiscal_year == 2025) & (products.category == "Services")]["value_usd_mn"].iloc[0] / 1000
    services_2024 = products[(products.fiscal_year == 2024) & (products.category == "Services")]["value_usd_mn"].iloc[0] / 1000
    gc = data["regions"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    add_takeaway_strip(
        slide,
        REPORT_LEFT,
        1.10,
        REPORT_WIDTH,
        0.78,
        "本页结论",
        "Apple FY2025 的财务主线是收入恢复、利润率上行与现金流保持高位；结构主线是 Services 抬升经营韧性，风险主线是硬件周期和区域竞争仍未消失。",
    )
    add_report_label(slide, 0.58, 2.12, "关键指标快照")
    add_report_table(
        slide,
        0.62,
        2.48,
        5.85,
        1.65,
        ["指标", "FY2025", "同比/变化", "阅读含义"],
        [
            ["收入", fmt_bn(fin.loc[2025, "revenue"]), f"{yoy(fin['revenue'], 2025, 2024):+.1%}", "五年新高"],
            ["净利润", fmt_bn(fin.loc[2025, "net_income"]), f"{yoy(fin['net_income'], 2025, 2024):+.1%}", "恢复至千亿美元以上"],
            ["Services", f"{services_2025:.1f}", f"{services_2025 / services_2024 - 1:+.1%}", "高毛利增量来源"],
            ["Greater China", f"{gc.loc[2025, 'Greater China'] / 1000:.1f}", f"{gc.loc[2025, 'Greater China'] / gc.loc[2024, 'Greater China'] - 1:+.1%}", "区域压力项"],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.12, 1.00, 1.05, 2.68],
        numeric_columns={1, 2},
    )
    add_report_label(slide, 6.78, 2.12, "三条主线")
    add_report_table(
        slide,
        6.82,
        2.48,
        5.78,
        1.65,
        ["主线", "结论", "后续跟踪"],
        [
            ["业绩", "收入、净利润同步修复", "FY2026 硬件周期"],
            ["结构", "Services 占比与毛利率继续上行", "监管和云/AI 成本"],
            ["风险", "Greater China 承压", "渠道、价格带和竞争格局"],
            ["资本", "现金流支撑回购分红", "净现金收敛与投资节奏"],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[0.88, 2.34, 2.56],
        numeric_columns=set(),
    )
    add_report_panel(slide, 0.62, 4.55, 11.98, 1.08, PAPER)
    add_text(slide, 0.82, 4.77, 11.55, 0.24, "点评口径", size=10.0, color=NAVY, bold=True, margin=0, first_line_indent=False)
    add_text(
        slide,
        0.82,
        5.10,
        11.55,
        0.22,
        "本材料只讨论 FY2025 公开财报与历史结构变化，不外推评级、目标价或实时交易建议；所有后续判断均需要结合 Apple 后续定期报告和监管文件更新。",
        size=8.8,
        color=TEXT,
        margin=0,
        line_spacing=1.20,
        first_line_indent=False,
    )


def slide_03(prs: Presentation, data: dict) -> None:
    """生成财务摘要页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "收入创五年新高，净利润恢复至千亿美元以上", 3, "数据来源：SEC companyfacts、Apple 10-K；金额单位为十亿美元。")
    fin = data["financials"].set_index("fiscal_year")
    years = [f"FY{y}" for y in fin.index]
    add_report_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.65,
        1.18,
        7.75,
        3.88,
        "图1：收入与净利润五年趋势",
        "单位：十亿美元",
        years,
        {
            "收入": [usd_bn(v) for v in fin["revenue"]],
            "净利润": [usd_bn(v) for v in fin["net_income"]],
        },
        value_format='0"bn"',
    )
    add_report_table(
        slide,
        0.65,
        5.20,
        7.75,
        1.32,
        ["指标"] + years,
        [
            ["收入"] + [fmt_bn(v) for v in fin["revenue"]],
            ["净利润"] + [fmt_bn(v) for v in fin["net_income"]],
            ["经营利润率"] + [fmt_pct(v) for v in fin["operating_margin"]],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.20, 1.31, 1.31, 1.31, 1.31, 1.31],
    )
    add_right_brief(
        slide,
        8.75,
        1.18,
        3.75,
        4.98,
        "FY2025 核心数据",
        [
            f"收入 {fmt_bn(fin.loc[2025, 'revenue'])}bn，同比增长 {yoy(fin['revenue'], 2025, 2024):.1%}。",
            f"经营利润率 {fmt_pct(fin.loc[2025, 'operating_margin'])}，较 FY2024 提升 {pct(fin.loc[2025, 'operating_margin'] - fin.loc[2024, 'operating_margin']):.1f}pct。",
            f"净利润 {fmt_bn(fin.loc[2025, 'net_income'])}bn，同比增长 {yoy(fin['net_income'], 2025, 2024):.1%}。",
            f"自由现金流 {fmt_bn(fin.loc[2025, 'fcf'])}bn，仍处于全球科技龙头高位。",
        ],
    )


def slide_04(prs: Presentation, data: dict) -> None:
    """生成收入结构页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "iPhone 仍是收入中枢，Services 成为最稳定的增量来源", 4, "数据来源：Apple FY2025 与 FY2022 Form 10-K，Net sales by category。")
    pivot = data["products"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    categories = [f"FY{y}" for y in pivot.index]
    series_order = ["iPhone", "Mac", "iPad", "Wearables, Home and Accessories", "Services"]
    add_report_chart(
        slide,
        XL_CHART_TYPE.COLUMN_STACKED,
        0.65,
        1.18,
        7.75,
        3.88,
        "图2：产品线收入结构",
        "单位：十亿美元",
        categories,
        {name: [v / 1000 for v in pivot[name]] for name in series_order},
        value_format='0"bn"',
    )
    services_share = pivot.loc[2025, "Services"] / pivot.loc[2025, "Total net sales"]
    iphone_share = pivot.loc[2025, "iPhone"] / pivot.loc[2025, "Total net sales"]
    add_report_table(
        slide,
        0.65,
        5.20,
        7.75,
        1.32,
        ["类别", "FY2025", "占比", "同比"],
        [
            ["iPhone", f"{pivot.loc[2025, 'iPhone'] / 1000:.1f}", f"{iphone_share:.1%}", f"{pivot.loc[2025, 'iPhone'] / pivot.loc[2024, 'iPhone'] - 1:+.1%}"],
            ["Services", f"{pivot.loc[2025, 'Services'] / 1000:.1f}", f"{services_share:.1%}", f"{pivot.loc[2025, 'Services'] / pivot.loc[2024, 'Services'] - 1:+.1%}"],
            ["Wearables", f"{pivot.loc[2025, 'Wearables, Home and Accessories'] / 1000:.1f}", f"{pivot.loc[2025, 'Wearables, Home and Accessories'] / pivot.loc[2025, 'Total net sales']:.1%}", f"{pivot.loc[2025, 'Wearables, Home and Accessories'] / pivot.loc[2024, 'Wearables, Home and Accessories'] - 1:+.1%}"],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[2.1, 1.9, 1.8, 1.95],
    )
    add_right_brief(
        slide,
        8.75,
        1.18,
        3.75,
        4.98,
        "结构点评",
        [
            f"iPhone FY2025 收入占比 {iphone_share:.1%}，仍是收入中枢。",
            f"Services FY2025 收入占比 {services_share:.1%}，同比增长 {pivot.loc[2025, 'Services'] / pivot.loc[2024, 'Services'] - 1:.1%}。",
            "Wearables/Home/Accessories 连续两年下滑，反映非核心硬件承压。",
            "服务化提升收入可见度，但硬件装机规模仍决定长期服务池上限。",
        ],
    )


def slide_05(prs: Presentation, data: dict) -> None:
    """生成利润率结构页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "Services 高毛利扩张继续抬升综合毛利率", 5, "数据来源：Apple 10-K Gross margin percentage；SEC companyfacts R&D expense。")
    gm = data["gross_margin"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    fin = data["financials"].set_index("fiscal_year")
    add_report_chart(
        slide,
        XL_CHART_TYPE.LINE_MARKERS,
        0.65,
        1.18,
        7.75,
        3.88,
        "图3：Products、Services 与综合毛利率",
        "单位：%",
        [f"FY{y}" for y in gm.index],
        {
            "Products 毛利率": list(gm["Products"]),
            "Services 毛利率": list(gm["Services"]),
            "综合毛利率": list(gm["Total gross margin percentage"]),
        },
        value_format='0.0"%"',
        palette=[RED, DARK_GOLD, WARM_BROWN],
    )
    add_report_table(
        slide,
        0.65,
        5.20,
        7.75,
        1.32,
        ["指标", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025"],
        [
            ["综合毛利率"] + [f"{v:.1f}" for v in gm["Total gross margin percentage"]],
            ["Services 毛利率"] + [f"{v:.1f}" for v in gm["Services"]],
            ["R&D ratio"] + [fmt_pct(v) for v in fin["rd_ratio"]],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.30, 1.29, 1.29, 1.29, 1.29, 1.29],
    )
    add_right_brief(
        slide,
        8.75,
        1.18,
        3.75,
        4.98,
        "机制拆解",
        [
            "Services 毛利率显著高于 Products，收入占比提升会自然推高综合毛利率。",
            f"FY2025 综合毛利率 {gm.loc[2025, 'Total gross margin percentage']:.1f}%，较 FY2021 提升 {gm.loc[2025, 'Total gross margin percentage'] - gm.loc[2021, 'Total gross margin percentage']:.1f}pct。",
            f"R&D ratio 从 FY2021 的 {pct(fin.loc[2021, 'rd_ratio']):.1f}% 提升至 FY2025 的 {pct(fin.loc[2025, 'rd_ratio']):.1f}%。",
            "后续利润率风险在于监管分成、云与 AI 成本、汇率和硬件促销压力。",
        ],
    )


def slide_06(prs: Presentation, data: dict) -> None:
    """生成地区表现页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "美欧日及亚太恢复支撑增长，大中华区延续承压", 6, "数据来源：Apple FY2025 与 FY2022 Form 10-K，Net sales by reportable segment。")
    regions = data["regions"].pivot(index="fiscal_year", columns="category", values="value_usd_mn")
    region_order = ["Americas", "Europe", "Greater China", "Japan", "Rest of Asia Pacific"]
    add_report_chart(
        slide,
        XL_CHART_TYPE.BAR_CLUSTERED,
        0.65,
        1.18,
        7.20,
        3.90,
        "图4：FY2025 地区收入排序",
        "单位：十亿美元",
        region_order,
        {"FY2025 地区收入": [regions.loc[2025, name] / 1000 for name in region_order]},
        value_format='0"bn"',
        legend=False,
        palette=[RED],
    )
    rows = []
    for name in region_order:
        value = regions.loc[2025, name] / 1000
        growth = regions.loc[2025, name] / regions.loc[2024, name] - 1
        rows.append([name, f"{value:.1f}", f"{growth:+.1%}"])
    add_report_table(slide, 8.25, 1.18, 4.18, 1.65, ["地区", "FY2025收入", "同比"], rows, font_size=TABLE_FONT_PT, col_widths=[1.78, 1.18, 1.22], numeric_columns={1, 2})
    add_right_brief(
        slide,
        8.25,
        3.18,
        4.18,
        2.98,
        "地区观察",
        [
            "Americas 与 Europe 仍是收入规模最大的两个地区，FY2025 增速分别为 6.8% 与 9.6%。",
            "Greater China FY2025 同比 -3.8%，连续两年回落，是最需要单独跟踪的区域变量。",
            "Japan 和 Rest of Asia Pacific 恢复较快，但基数和战略含义弱于美欧与大中华区。",
        ],
    )


def slide_07(prs: Presentation, data: dict) -> None:
    """生成现金流与资本回报页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "自由现金流仍然厚实，资本回报维持高强度", 7, "数据来源：SEC companyfacts、Apple 10-K；金额单位为十亿美元。")
    fin = data["financials"].set_index("fiscal_year")
    add_report_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.65,
        1.18,
        7.75,
        3.88,
        "图5：经营现金流、自由现金流与资本回报",
        "单位：十亿美元",
        [f"FY{y}" for y in fin.index],
        {
            "经营现金流": [usd_bn(v) for v in fin["operating_cash_flow"]],
            "自由现金流": [usd_bn(v) for v in fin["fcf"]],
            "回购+分红": [usd_bn(v) for v in fin["capital_return"]],
        },
        value_format='0"bn"',
    )
    add_report_table(
        slide,
        0.65,
        5.20,
        7.75,
        1.32,
        ["指标", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025"],
        [
            ["自由现金流"] + [fmt_bn(v) for v in fin["fcf"]],
            ["资本回报"] + [fmt_bn(v) for v in fin["capital_return"]],
            ["FCF margin"] + [fmt_pct(v) for v in fin["fcf_margin"]],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.30, 1.29, 1.29, 1.29, 1.29, 1.29],
    )
    add_right_brief(
        slide,
        8.75,
        1.18,
        3.75,
        4.98,
        "资本配置点评",
        [
            f"FY2025 自由现金流 {fmt_bn(fin.loc[2025, 'fcf'])}bn，低于 FY2024 但仍保持强韧。",
            f"FY2025 回购与分红合计 {fmt_bn(fin.loc[2025, 'capital_return'])}bn，连续五年保持高强度。",
            "回购对 EPS 与股东回报有支撑，但也会持续压缩净现金头寸。",
            "AI 投资、供应链资本开支和监管成本上行时，资本回报节奏需要观察。",
        ],
    )


def slide_08(prs: Presentation, data: dict) -> None:
    """生成资产负债表页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "净现金收敛但仍为正，资产负债表保持弹性", 8, "数据来源：SEC companyfacts；总债务=Commercial paper + current long-term debt + non-current long-term debt。")
    fin = data["financials"].set_index("fiscal_year")
    add_report_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        0.65,
        1.18,
        7.75,
        3.88,
        "图6：现金及可出售证券、总债务与净现金",
        "单位：十亿美元",
        [f"FY{y}" for y in fin.index],
        {
            "现金及可出售证券": [usd_bn(v) for v in fin["cash_marketable_securities"]],
            "总债务": [usd_bn(v) for v in fin["total_debt"]],
            "净现金": [usd_bn(v) for v in fin["net_cash"]],
        },
        value_format='0"bn"',
    )
    add_report_table(
        slide,
        0.65,
        5.20,
        7.75,
        1.32,
        ["指标", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025"],
        [
            ["现金及可出售证券"] + [fmt_bn(v) for v in fin["cash_marketable_securities"]],
            ["总债务"] + [fmt_bn(v) for v in fin["total_debt"]],
            ["净现金"] + [fmt_bn(v) for v in fin["net_cash"]],
        ],
        font_size=TABLE_FONT_PT,
        col_widths=[1.50, 1.25, 1.25, 1.25, 1.25, 1.25],
    )
    add_right_brief(
        slide,
        8.75,
        1.18,
        3.75,
        4.98,
        "安全边际",
        [
            f"FY2025 现金及可出售证券 {fmt_bn(fin.loc[2025, 'cash_marketable_securities'])}bn，总债务 {fmt_bn(fin.loc[2025, 'total_debt'])}bn。",
            f"净现金 {fmt_bn(fin.loc[2025, 'net_cash'])}bn，较 FY2021 明显收敛但仍为正。",
            "总债务口径包含 commercial paper、current long-term debt 与 non-current long-term debt。",
            "现金储备为研发、供应链、回购和潜在监管冲击提供缓冲。",
        ],
    )


def slide_09(prs: Presentation, data: dict) -> None:
    """生成关键观察变量页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "后续跟踪重点落在 AI 终端、服务生态和区域竞争", 9, "数据来源：基于 Apple FY2025 10-K 披露风险因素与经营数据的研究框架整理。")
    add_takeaway_strip(
        slide,
        REPORT_LEFT,
        1.10,
        REPORT_WIDTH,
        0.72,
        "研究框架",
        "FY2025 的财务质量已经较强，后续再加速需要产品周期、服务生态变现和区域竞争压力三者共同改善。",
    )
    rows = [
        ["AI 终端与换机周期", "Apple Intelligence 与端侧 AI 体验能否转化为换机需求；iPhone 价格带、渠道库存和新兴市场渗透率。", "决定硬件收入弹性与 Products 毛利率。"],
        ["Services monetization", "App Store、iCloud、支付、广告和订阅组合的增长持续性；监管对抽佣和默认入口的影响。", "决定 Services 增速、毛利率和综合利润率中枢。"],
        ["区域与供应链风险", "Greater China 竞争格局、需求和渠道变化；印度、东南亚生产与销售端扩张效率。", "决定区域收入分化、成本韧性和地缘风险暴露。"],
        ["资本配置与监管", "回购、分红、AI 投资、供应链资本开支和监管成本之间的优先级变化。", "决定净现金收敛速度和股东回报持续性。"],
    ]
    add_report_table(
        slide,
        0.65,
        2.20,
        12.00,
        3.45,
        ["观察变量", "跟踪指标", "财报含义"],
        rows,
        font_size=TABLE_FONT_PT,
        col_widths=[2.28, 6.20, 3.52],
        numeric_columns=set(),
    )
    add_text(slide, 0.68, 5.95, 11.7, 0.22, "建议把上述变量纳入后续季度跟踪表，而不是只在年报点评中一次性描述。", size=7.8, color=MUTED_TEXT, margin=0, first_line_indent=False)


def slide_10(prs: Presentation, data: dict) -> None:
    """生成风险声明页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "风险声明：强现金流不消除经营、监管和估值波动风险", 10, "数据来源：Apple FY2025 Form 10-K Risk Factors 与本材料整理。")
    rows = [
        ["产品与创新", "iPhone、Mac、iPad、Wearables 等产品周期不及预期；AI 终端体验、生态协同或消费者换机意愿低于预期。"],
        ["竞争与区域", "智能手机、PC、服务生态和本地互联网服务竞争加剧；Greater China 等重点地区需求、渠道与竞争格局变化。"],
        ["供应链与地缘", "关键零部件、制造产能、物流、出口管制、关税和地缘政治事件影响供给、成本或销售。"],
        ["监管与法律", "App Store 抽佣、默认入口、数据隐私、反垄断、税务和内容监管政策变化影响收入与利润率。"],
        ["财务与市场", "汇率、利率、税率、市场流动性和估值波动影响财务表现与市场价格；回购节奏可能受现金与监管约束。"],
        ["数据与方法", "本材料基于公开 10-K 和 SEC XBRL 数据整理，表格抽取、科目映射和期间口径可能存在解释差异。"],
    ]
    add_report_table(slide, 0.65, 1.25, 12.00, 4.75, ["风险类型", "风险说明"], rows, font_size=TABLE_FONT_PT, col_widths=[2.05, 9.95], numeric_columns=set())
    add_text(slide, 0.68, 6.25, 11.7, 0.25, "以上风险并非完整清单，读者应结合 Apple 后续定期报告、监管文件及宏观市场环境持续更新判断。", size=7.8, color=MUTED_TEXT, margin=0, first_line_indent=False)


def slide_11(prs: Presentation, data: dict) -> None:
    """生成免责声明与来源页。"""

    slide = blank_slide(prs)
    add_report_shell(slide, "免责声明与数据来源", 11, "数据来源：SEC 官方 API 与 Apple 10-K；本页为合规与来源说明。")
    disclaimer = [
        "仅供学术交流使用，不代表任何组织机构，机构和个人的观点和立场。",
        "本材料不构成任何证券、基金、衍生品或其他金融产品的投资建议、评级、招揽、承诺或收益保证。",
        "本材料仅基于公开资料整理，未对 Apple Inc.、其审计师、监管机构或任何第三方进行访谈确认。",
        "读者应独立核验数据来源，并结合自身研究目标、风险承受能力和合规要求使用本材料。",
    ]
    add_body_panel(slide, 0.65, 1.18, 5.65, 4.78, "免责声明", disclaimer)
    sources = data["sources"].sort_values("fiscal_year", ascending=False)
    source_rows = [
        [f"FY{int(row.fiscal_year)}", str(row.filing_date), str(row.accession)]
        for _, row in sources.iterrows()
    ]
    add_report_table(slide, 6.72, 1.18, 5.65, 1.78, ["报告", "Filed", "Accession"], source_rows, font_size=TABLE_FONT_PT, col_widths=[0.86, 1.18, 3.61], numeric_columns=set())
    add_report_label(slide, 6.74, 3.28, "主要数据来源")
    add_text(
        slide,
        6.78,
        3.66,
        5.55,
        1.15,
        "主要数据来源：SEC submissions API、SEC companyfacts API、Apple Inc. FY2021-FY2025 Form 10-K。完整链接与本地文件路径见 data/processed/sources.md 与 apple_10k_sources.csv。",
        size=8.2,
        color=TEXT,
        margin=0.02,
        line_spacing=1.25,
        first_line_indent=False,
    )
    add_report_panel(slide, 6.72, 5.08, 5.65, 0.86, SOFT_GOLD_BG)
    add_text(slide, 6.90, 5.32, 5.28, 0.20, "金额单位默认使用美元；图表通常换算为十亿美元，表格来源以原始披露口径为准。", size=7.5, color=MUTED_TEXT, margin=0, line_spacing=1.15, first_line_indent=False)


def fix_docprops_slide_count(pptx_path: Path, slide_count: int) -> None:
    """修正 `docProps/app.xml` 中的 Slides 统计。"""

    ns = {"ep": "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"}
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        with zipfile.ZipFile(pptx_path, "r") as source_zip:
            source_zip.extractall(temp_path)
        app_xml = temp_path / "docProps" / "app.xml"
        tree = ET.parse(app_xml)
        root = tree.getroot()
        slides = root.find("ep:Slides", ns)
        if slides is None:
            slides = ET.SubElement(root, f"{{{ns['ep']}}}Slides")
        slides.text = str(slide_count)
        ET.register_namespace("", ns["ep"])
        tree.write(app_xml, encoding="UTF-8", xml_declaration=True)

        fixed_path = pptx_path.with_suffix(".fixed.pptx")
        with zipfile.ZipFile(fixed_path, "w", compression=zipfile.ZIP_DEFLATED) as target_zip:
            for file_path in temp_path.rglob("*"):
                if file_path.is_file():
                    target_zip.write(file_path, file_path.relative_to(temp_path).as_posix())
        shutil.move(fixed_path, pptx_path)


def main() -> None:
    """执行 PPT 构建。"""

    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    data = load_inputs()
    prs = make_prs()
    builders = [
        slide_01,
        slide_02,
        slide_03,
        slide_04,
        slide_05,
        slide_06,
        slide_07,
        slide_08,
        slide_09,
        slide_10,
        slide_11,
    ]
    for builder in builders:
        builder(prs, data)
    prs.save(OUT_PPTX)
    fix_docprops_slide_count(OUT_PPTX, len(builders))
    prs.save(FINAL_PPTX)
    fix_docprops_slide_count(FINAL_PPTX, len(builders))
    summary = {
        "pptx": str(OUT_PPTX.relative_to(WORKSPACE)),
        "final_pptx": str(FINAL_PPTX.relative_to(WORKSPACE)),
        "slides": len(builders),
        "built_at": date(2026, 5, 10).isoformat(),
        "source_spec": str(SPEC_PATH.relative_to(WORKSPACE)),
    }
    (BUILD_DIR / "generated" / "build_summary.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
